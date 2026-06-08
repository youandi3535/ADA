"""outputs.style.visual_kit - PPT shape/gradient helpers for consulting-grade design."""

from __future__ import annotations


def hex_to_rgb(h: str) -> tuple[int, int, int]:
    h = h.lstrip("#")
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_rect(slide, x_cm, y_cm, w_cm, h_cm, hex_color, *, line=False):
    """Solid filled rectangle. Returns shape."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Cm

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(hex_color))
    if not line:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, x_cm, y_cm, w_cm, h_cm, hex_color, *, line_hex=None):
    """Rounded rectangle - good for KPI cards / pills."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Cm

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(hex_color))
    if line_hex:
        shape.line.color.rgb = RGBColor(*hex_to_rgb(line_hex))
        shape.line.width = Cm(0.02)
    else:
        shape.line.fill.background()
    return shape


def add_oval(slide, x_cm, y_cm, w_cm, h_cm, hex_color):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Cm

    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(hex_color))
    shape.line.fill.background()
    return shape


def add_chevron(slide, x_cm, y_cm, w_cm, h_cm, hex_color):
    """Right-pointing arrow chevron - for process flow."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Cm

    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(hex_color))
    shape.line.fill.background()
    return shape


def set_text(
    shape, text: str, *, font="Malgun Gothic", size_pt=14, bold=False, color_hex="#0F172A", align="left", vcenter=True
):
    """Set text inside a shape with KO font."""
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Pt

    tf = shape.text_frame
    tf.word_wrap = True
    if vcenter:
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass
    tf.text = text
    for p in tf.paragraphs:
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
        for r in p.runs:
            r.font.size = Pt(size_pt)
            r.font.bold = bold
            r.font.color.rgb = RGBColor(*hex_to_rgb(color_hex))
            r.font.name = font
            _ea_font(r, font)


def _ea_font(run, family: str) -> None:
    """Set east-asian (Korean) typeface explicitly."""
    try:
        from lxml import etree
        from pptx.oxml.ns import qn

        rPr = run._r.get_or_add_rPr()
        for ea in rPr.findall(qn("a:ea")):
            rPr.remove(ea)
        ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", family)
    except Exception:
        pass


def add_text_box(slide, x_cm, y_cm, w_cm, h_cm, text: str, **style):
    """Convenience: textbox + set_text."""
    from pptx.util import Cm

    tx = slide.shapes.add_textbox(Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm))
    set_text(tx, text, **style)
    return tx


def add_horizontal_rule(slide, x_cm, y_cm, w_cm, hex_color, *, h_cm=0.05):
    """Thin horizontal divider line."""
    return add_rect(slide, x_cm, y_cm, w_cm, h_cm, hex_color)


def add_vertical_accent(slide, x_cm, y_cm, h_cm, hex_color, *, w_cm=0.18):
    """Vertical color bar - for slide title left-side accent."""
    return add_rect(slide, x_cm, y_cm, w_cm, h_cm, hex_color)


def add_gradient_rect(slide, x_cm, y_cm, w_cm, h_cm, hex_start, hex_end, *, angle=0):
    """Linear gradient rectangle via OXML."""
    from lxml import etree
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    from pptx.util import Cm

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm))
    shape.line.fill.background()
    sp = shape.fill._xPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill"):
        for el in sp.findall(qn(tag)):
            sp.remove(el)
    grad = etree.SubElement(sp, qn("a:gradFill"))
    grad.set("flip", "none")
    grad.set("rotWithShape", "1")
    gsLst = etree.SubElement(grad, qn("a:gsLst"))
    for pos, hx in ((0, hex_start), (100000, hex_end)):
        gs = etree.SubElement(gsLst, qn("a:gs"))
        gs.set("pos", str(pos))
        clr = etree.SubElement(gs, qn("a:srgbClr"))
        clr.set("val", hx.lstrip("#"))
    lin = etree.SubElement(grad, qn("a:lin"))
    lin.set("ang", str(int(angle * 60000)))
    lin.set("scaled", "1")
    etree.SubElement(grad, qn("a:tileRect"))
    return shape


def add_glyph(slide, x_cm, y_cm, w_cm, h_cm, glyph, color_hex="#FFFFFF", size_pt=20, bg_color=None):
    """Unicode glyph icon with optional circular background."""
    if bg_color:
        add_oval(slide, x_cm, y_cm, w_cm, h_cm, bg_color)
    add_text_box(
        slide,
        x_cm,
        y_cm,
        w_cm,
        h_cm,
        glyph,
        size_pt=size_pt,
        color_hex=color_hex,
        align="center",
        vcenter=True,
        bold=True,
    )


def add_triangle(slide, x_cm, y_cm, w_cm, h_cm, hex_color, *, direction="up"):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Cm

    shape_map = {
        "up": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "right": MSO_SHAPE.RIGHT_TRIANGLE,
        "diamond": MSO_SHAPE.DIAMOND,
        "hex": MSO_SHAPE.HEXAGON,
    }
    sh = shape_map.get(direction, MSO_SHAPE.ISOSCELES_TRIANGLE)
    shape = slide.shapes.add_shape(sh, Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(hex_color))
    shape.line.fill.background()
    return shape


GLYPHS = {
    "data": "▦",
    "chart": "▲",
    "report": "▤",
    "model": "◆",
    "warn": "⚠",
    "ok": "✓",
    "no": "✗",
    "info": "ⓘ",
    "arrow_r": "▶",
    "arrow_d": "▼",
    "arrow_u": "▲",
    "star": "★",
    "hex": "⬢",
    "circle": "●",
    "diamond": "◆",
    "bulb": "✦",
    "target": "◎",
    "settings": "⚙",
    "shield": "▼",
}


# ==============================================================
# HJ 2026-06-08 — 시각 품질 향상 Phase 1
# 그라데이션·그림자·라운드 카드 확장 (순수 python-pptx, CPU 100%)
# ==============================================================


def add_gradient_rect_multi(slide, x_cm, y_cm, w_cm, h_cm, stops, *, angle=0):
    """다중 stop 그라데이션 사각형. stops = [(position_pct, hex_color), ...].

    예: stops=[(0, "#1e40af"), (50, "#3b82f6"), (100, "#93c5fd")]
    """
    from lxml import etree
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    from pptx.util import Cm

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm))
    shape.line.fill.background()
    fill = shape.fill
    fill.gradient_stops  # ensure gradient
    fill_el = shape.fill._xPr.find(qn("a:gradFill"))
    if fill_el is None:
        # python-pptx 의 gradient 설정 한계 — XML 직접 조작
        sp_pr = shape.fill._xPr
        # 기존 fill 노드 제거
        for tag in ("a:solidFill", "a:noFill", "a:blipFill", "a:pattFill"):
            existing = sp_pr.find(qn(tag))
            if existing is not None:
                sp_pr.remove(existing)
        # gradFill 새로 생성
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        gf = etree.SubElement(sp_pr, qn("a:gradFill"), nsmap=nsmap)
        gs_lst = etree.SubElement(gf, qn("a:gsLst"))
        for pos_pct, hex_c in stops:
            gs = etree.SubElement(gs_lst, qn("a:gs"), pos=str(int(pos_pct * 1000)))
            r, g, b = hex_to_rgb(hex_c)
            etree.SubElement(gs, qn("a:srgbClr"), val=f"{r:02X}{g:02X}{b:02X}")
        # 방향: lin (angle in 60000ths of a degree)
        lin = etree.SubElement(gf, qn("a:lin"), ang=str(int(angle * 60000)), scaled="0")
        _ = lin  # noqa: F841
    # 외곽선 없애기
    try:
        shape.line.fill.background()
    except Exception:
        pass
    return shape


def add_drop_shadow(shape, *, blur_pt=8, distance_pt=4, alpha_pct=40, direction_deg=90):
    """python-pptx Shape 에 drop shadow 효과 적용 (XML 직접).

    Args:
        shape: pptx 의 도형 객체
        blur_pt: 그림자 흐림 정도 (포인트)
        distance_pt: 그림자 거리 (포인트)
        alpha_pct: 투명도 0~100
        direction_deg: 빛 방향 (0=오른쪽, 90=아래쪽)
    """
    from lxml import etree
    from pptx.oxml.ns import qn

    sp_pr = shape._element.find(qn("p:spPr")) or shape._element.find(qn("a:spPr"))
    if sp_pr is None:
        # fallback — 일반 shape
        sp_pr = shape.fill._xPr
    if sp_pr is None:
        return shape

    # effect list 추가
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    # 기존 effectLst 제거
    existing = sp_pr.find(qn("a:effectLst"))
    if existing is not None:
        sp_pr.remove(existing)
    effect_lst = etree.SubElement(sp_pr, qn("a:effectLst"), nsmap=nsmap)
    outer_shdw = etree.SubElement(
        effect_lst,
        qn("a:outerShdw"),
        blurRad=str(int(blur_pt * 12700)),
        dist=str(int(distance_pt * 12700)),
        dir=str(int(direction_deg * 60000)),
        rotWithShape="0",
        algn="ctr",
    )
    srgb = etree.SubElement(outer_shdw, qn("a:srgbClr"), val="000000")
    etree.SubElement(srgb, qn("a:alpha"), val=str(int(alpha_pct * 1000)))
    return shape


def add_rounded_card_with_shadow(slide, x_cm, y_cm, w_cm, h_cm, hex_color="#FFFFFF", border_hex=None, shadow=True):
    """그림자 있는 라운드 카드 — KPI 카드·정보 박스용."""
    shape = add_rounded_rect(slide, x_cm, y_cm, w_cm, h_cm, hex_color, line_hex=border_hex)
    if shadow:
        add_drop_shadow(shape, blur_pt=12, distance_pt=6, alpha_pct=25)
    return shape


def add_image_with_overlay(slide, x_cm, y_cm, w_cm, h_cm, image_path, overlay_color="#000000", overlay_alpha_pct=40):
    """이미지 + 반투명 overlay (표지·섹션 디바이더용)."""
    from pptx.util import Cm

    img = slide.shapes.add_picture(image_path, Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm))
    # 반투명 검정 overlay
    overlay = add_rect(slide, x_cm, y_cm, w_cm, h_cm, overlay_color, line=False)
    _set_shape_alpha(overlay, overlay_alpha_pct)
    return img, overlay


def _set_shape_alpha(shape, alpha_pct: int):
    """Shape 의 fill 에 alpha 적용 (XML 조작)."""
    from lxml import etree
    from pptx.oxml.ns import qn

    sp_pr = shape.fill._xPr
    solid = sp_pr.find(qn("a:solidFill"))
    if solid is None:
        return shape
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is None:
        return shape
    # 기존 alpha 제거
    existing = srgb.find(qn("a:alpha"))
    if existing is not None:
        srgb.remove(existing)
    etree.SubElement(srgb, qn("a:alpha"), val=str(int(alpha_pct * 1000)))
    return shape


def add_color_band(slide, x_cm, y_cm, w_cm, h_cm, hex_color="#2563eb", gradient_to=None, angle=0):
    """좌측·상단 색띠 (섹션 마커용). gradient_to 가 있으면 그라데이션."""
    if gradient_to:
        return add_gradient_rect_multi(
            slide,
            x_cm,
            y_cm,
            w_cm,
            h_cm,
            stops=[(0, hex_color), (100, gradient_to)],
            angle=angle,
        )
    return add_rect(slide, x_cm, y_cm, w_cm, h_cm, hex_color, line=False)
