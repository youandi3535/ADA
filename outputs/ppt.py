"""outputs.ppt — OUT-01 PowerPoint 발표자료 (python-pptx). ADR-008 L2 reattach 통합."""

from __future__ import annotations

from typing import Any

from outputs.base import OutputGenerator, get_theme, reattach_pii


def _set_title_color(slide: Any, rgb: tuple[int, int, int]) -> None:
    try:
        from pptx.dml.color import RGBColor

        for para in slide.shapes.title.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = RGBColor(*rgb)
    except Exception:
        pass


def _set_background(slide: Any, rgb: tuple[int, int, int]) -> None:
    try:
        from pptx.dml.color import RGBColor

        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*rgb)
    except Exception:
        pass


class PresentationGenerator(OutputGenerator):
    output_code = "OUT-01"
    extension = "pptx"

    def generate(
        self,
        *,
        insights: str,
        best_model: dict[str, Any],
        eda_charts: list[str],
        category: str,
        user_intent: str,
        eval_result: dict[str, Any] | None,
        state: Any = None,
    ) -> str:
        # ADR-008 L2 — PII 마스킹
        insights = reattach_pii(state, insights)
        user_intent = reattach_pii(state, user_intent)
        if eval_result:
            eval_result = {**eval_result, "rationale": reattach_pii(state, eval_result.get("rationale"))}

        from pptx import Presentation
        from pptx.util import Inches, Pt

        theme = get_theme(category)
        primary = theme["primary_rgb"]
        accent = theme["accent_rgb"]
        label_ko = theme["label_ko"]

        prs = Presentation()

        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "ADA Auto Report"
        slide.placeholders[1].text = f"Category: {label_ko} ({category})\nIntent: {user_intent or '-'}"
        _set_title_color(slide, primary)

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Insights"
        _set_title_color(slide, primary)
        tf = slide.placeholders[1].text_frame
        tf.text = (insights or "")[:1500]
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.size = Pt(18)

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Best Model"
        _set_title_color(slide, primary)
        body = slide.placeholders[1].text_frame
        bm = best_model or {}
        body.text = f"Model: {bm.get('model_name', '-')}"
        for k, v in (bm.get("metrics") or {}).items():
            p = body.add_paragraph()
            p.text = f"{k}: {v:.4f}" if isinstance(v, float) else f"{k}: {v}"

        # HJ — Method B: 차트 다운로드 병렬화.
        # eda_charts + extras['charts'] 를 한 번에 모아 ThreadPoolExecutor 로 동시 다운로드.
        # 슬라이드 append 자체는 순차 (python-pptx 비스레드 안전).
        eda_paths = list(eda_charts[:4])
        extras = self._call_extras(state, ctx={"output_code": self.output_code, "category": category})
        extras_paths = list(extras.get("charts", [])[:4])
        local_paths = self._download_charts_parallel(eda_paths + extras_paths)
        eda_local = local_paths[: len(eda_paths)]
        extras_local = local_paths[len(eda_paths) :]

        for i, tmp in enumerate(eda_local):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"EDA Chart #{i + 1}"
            _set_title_color(slide, primary)
            if tmp:
                try:
                    slide.shapes.add_picture(tmp, Inches(0.5), Inches(1.5), width=Inches(8))
                except Exception:
                    pass

        for j, tmp in enumerate(extras_local):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"[{label_ko}] Cat Analysis #{j + 1}"
            _set_title_color(slide, primary)
            if tmp:
                try:
                    slide.shapes.add_picture(tmp, Inches(0.5), Inches(1.5), width=Inches(8))
                except Exception:
                    pass

        for tbl in extras.get("tables", [])[:3]:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = str(tbl.get("title", f"[{label_ko}] Table"))
            _set_title_color(slide, primary)
            try:
                rows = list(tbl.get("rows", []))
                cols = list(tbl.get("columns") or (list(rows[0].keys()) if rows else []))
                if cols and rows:
                    n_rows = min(len(rows), 10) + 1
                    n_cols = len(cols)
                    table_shape = slide.shapes.add_table(
                        n_rows, n_cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.5 * n_rows)
                    )
                    tbl_obj = table_shape.table
                    for c_idx, c in enumerate(cols):
                        tbl_obj.cell(0, c_idx).text = str(c)
                    for r_idx, row in enumerate(rows[:10]):
                        for c_idx, c in enumerate(cols):
                            tbl_obj.cell(r_idx + 1, c_idx).text = str(row.get(c, ""))
            except Exception:
                pass

        for text_block in extras.get("text_blocks", [])[:3]:
            safe_block = reattach_pii(state, str(text_block))
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"[{label_ko}] Notes"
            _set_title_color(slide, primary)
            tf = slide.placeholders[1].text_frame
            tf.text = safe_block[:1500]

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Evaluation"
        _set_title_color(slide, primary)
        _set_background(slide, accent)
        body = slide.placeholders[1].text_frame
        ev = eval_result or {}
        body.text = f"passed: {ev.get('passed')}\nrationale: {(ev.get('rationale') or '')[:300]}"

        local = self._tmp()
        prs.save(local)
        return self._upload(local)
