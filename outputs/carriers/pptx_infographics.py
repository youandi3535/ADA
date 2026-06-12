"""outputs.carriers.pptx_infographics - Advanced infographic layouts.

Templated infographic layouts that replace plain text slides:
  - statistics_grid: 4 big numbers + icons + captions
  - timeline_milestones: horizontal nodes + connecting line
  - hex_grid_6: 6 hexagons arranged in honeycomb
  - radial_nodes: center hub + 6 satellite nodes
  - section_team_grid: photo/icon + name + role cards
"""

from __future__ import annotations

from outputs.style.visual_kit import (
    GLYPHS,
    add_chevron,
    add_gradient_rect,
    add_oval,
    add_rect,
    add_rounded_rect,
    add_text_box,
    add_triangle,
)

SLIDE_W = 33.867
SLIDE_H = 19.05


def draw_statistics_grid(slide, items, primary, accent, ink, muted, light_bg):
    """4 large statistics with icon + value + caption.

    items: list of {value, label, caption, icon_glyph?}
    """
    glyphs = [GLYPHS["chart"], GLYPHS["target"], GLYPHS["bulb"], GLYPHS["star"]]
    n = min(len(items), 4)
    margin_x = 1.5
    gap = 0.5
    avail_w = SLIDE_W - 2 * margin_x
    card_w = (avail_w - (n - 1) * gap) / n
    y = 5.5
    card_h = SLIDE_H - y - 2.0

    for i in range(n):
        item = items[i]
        x = margin_x + i * (card_w + gap)
        # Background gradient card (light to lighter)
        add_gradient_rect(slide, x, y, card_w, card_h, light_bg, "#FFFFFF", angle=90)
        # Top color bar
        add_rect(slide, x, y, card_w, 0.4, primary)
        # Large icon circle
        ic_size = 2.2
        add_oval(slide, x + card_w / 2 - ic_size / 2, y + 1.0, ic_size, ic_size, primary)
        add_text_box(
            slide,
            x + card_w / 2 - ic_size / 2,
            y + 1.0,
            ic_size,
            ic_size,
            item.get("icon_glyph", glyphs[i % len(glyphs)]),
            size_pt=32,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Big value
        add_text_box(
            slide,
            x,
            y + 4.0,
            card_w,
            2.2,
            str(item.get("value", "—")),
            size_pt=46,
            bold=True,
            color_hex=primary,
            align="center",
            vcenter=True,
        )
        # Label (small uppercase)
        add_text_box(
            slide,
            x,
            y + 6.4,
            card_w,
            0.7,
            str(item.get("label", "")).upper()[:30],
            size_pt=10,
            bold=True,
            color_hex=muted,
            align="center",
            vcenter=True,
        )
        # Accent line
        add_rect(slide, x + card_w / 2 - 1.2, y + 7.3, 2.4, 0.06, accent)
        # Caption
        add_text_box(
            slide,
            x + 0.3,
            y + 7.6,
            card_w - 0.6,
            card_h - 8.0,
            str(item.get("caption", ""))[:80],
            size_pt=10,
            color_hex=ink,
            align="center",
            vcenter=False,
        )


def draw_timeline_milestones(slide, milestones, primary, accent, ink, muted, light_bg):
    """Horizontal milestone timeline.

    milestones: list of {year, label, caption}
    """
    n = min(len(milestones), 6)
    if n == 0:
        return
    margin_x = 1.8
    avail_w = SLIDE_W - 2 * margin_x
    step = avail_w / max(1, n - 1) if n > 1 else avail_w
    y_line = 10.0

    # Main horizontal line
    add_rect(slide, margin_x, y_line - 0.05, avail_w, 0.1, "#CBD5E1")
    # Colored progress portion (full)
    add_rect(slide, margin_x, y_line - 0.05, avail_w * 0.95, 0.1, primary)

    glyphs = [GLYPHS["data"], GLYPHS["chart"], GLYPHS["settings"], GLYPHS["target"], GLYPHS["bulb"], GLYPHS["ok"]]
    for i in range(n):
        m = milestones[i]
        x = margin_x + i * step if n > 1 else margin_x + avail_w / 2
        # Outer ring
        ring_size = 2.0
        add_oval(slide, x - ring_size / 2, y_line - ring_size / 2, ring_size, ring_size, accent)
        # Inner node
        node_size = 1.5
        add_oval(slide, x - node_size / 2, y_line - node_size / 2, node_size, node_size, primary)
        # Icon glyph
        add_text_box(
            slide,
            x - node_size / 2,
            y_line - node_size / 2,
            node_size,
            node_size,
            glyphs[i % len(glyphs)],
            size_pt=14,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Year/label above node (alternate up/down)
        if i % 2 == 0:
            add_text_box(
                slide,
                x - 3.0,
                y_line - 4.5,
                6.0,
                1.0,
                str(m.get("year", f"M{i + 1}")),
                size_pt=20,
                bold=True,
                color_hex=primary,
                align="center",
                vcenter=True,
            )
            add_rounded_rect(slide, x - 3.0, y_line - 3.3, 6.0, 1.8, light_bg, line_hex=primary)
            add_text_box(
                slide,
                x - 2.8,
                y_line - 3.2,
                5.6,
                1.6,
                str(m.get("label", ""))[:80],
                size_pt=10,
                color_hex=ink,
                align="center",
                vcenter=True,
            )
        else:
            add_text_box(
                slide,
                x - 3.0,
                y_line + 3.5,
                6.0,
                1.0,
                str(m.get("year", f"M{i + 1}")),
                size_pt=20,
                bold=True,
                color_hex=primary,
                align="center",
                vcenter=True,
            )
            add_rounded_rect(slide, x - 3.0, y_line + 1.5, 6.0, 1.8, light_bg, line_hex=primary)
            add_text_box(
                slide,
                x - 2.8,
                y_line + 1.6,
                5.6,
                1.6,
                str(m.get("label", ""))[:80],
                size_pt=10,
                color_hex=ink,
                align="center",
                vcenter=True,
            )


def draw_hex_grid_6(slide, items, primary, accent, ink, muted):
    """6 hexagons arranged in honeycomb pattern.

    items: list of {title, caption, icon_glyph?}
    """
    n = min(len(items), 6)
    glyphs = [GLYPHS["data"], GLYPHS["target"], GLYPHS["bulb"], GLYPHS["model"], GLYPHS["chart"], GLYPHS["star"]]
    # Honeycomb positions (cx, cy in cm)
    hex_size = 4.0
    positions = [
        (10.0, 7.0),
        (15.0, 7.0),
        (20.0, 7.0),
        (12.5, 11.5),
        (17.5, 11.5),
        (22.5, 11.5),
    ]
    for i in range(n):
        item = items[i]
        cx, cy = positions[i]
        # Hex (using DIAMOND as fallback if HEXAGON not available)
        x = cx - hex_size / 2
        y = cy - hex_size / 2
        add_triangle(slide, x, y, hex_size, hex_size, primary, direction="hex")
        # Glyph
        add_text_box(
            slide,
            x,
            y + 0.3,
            hex_size,
            hex_size * 0.5,
            item.get("icon_glyph", glyphs[i % len(glyphs)]),
            size_pt=24,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Title at bottom
        add_text_box(
            slide,
            x,
            y + hex_size * 0.55,
            hex_size,
            hex_size * 0.4,
            str(item.get("title", ""))[:20],
            size_pt=11,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Caption outside hex (below)
        add_text_box(
            slide,
            x - 0.5,
            y + hex_size + 0.2,
            hex_size + 1.0,
            1.2,
            str(item.get("caption", ""))[:60],
            size_pt=9,
            color_hex=ink,
            align="center",
            vcenter=False,
        )


def draw_radial_nodes(slide, center_text, items, primary, accent, ink, muted, light_bg):
    """Hub-and-spoke radial layout with 6 satellite nodes.

    items: list of {label, caption, icon_glyph?}
    """
    import math

    cx, cy = SLIDE_W / 2, 11.0
    hub_size = 3.5
    sat_size = 2.4
    radius = 6.5
    n = min(len(items), 6)
    glyphs = [GLYPHS["data"], GLYPHS["chart"], GLYPHS["settings"], GLYPHS["target"], GLYPHS["bulb"], GLYPHS["ok"]]

    # Connecting lines (drawn first, behind nodes)
    for i in range(n):
        angle = math.radians(-90 + i * (360 / n))
        sx = cx + radius * math.cos(angle)
        sy = cy + radius * math.sin(angle)
        # Approximate line as thin rounded rect
        # Direction vector
        # Use add_rect rotated would be ideal but pptx lib needs rotation
        # Skip connector line - rely on visual proximity of hub and satellite

    # Hub at center
    add_oval(slide, cx - hub_size / 2 - 0.3, cy - hub_size / 2 - 0.3, hub_size + 0.6, hub_size + 0.6, accent)
    add_oval(slide, cx - hub_size / 2, cy - hub_size / 2, hub_size, hub_size, primary)
    add_text_box(
        slide,
        cx - hub_size / 2,
        cy - hub_size / 2,
        hub_size,
        hub_size,
        center_text[:20],
        size_pt=13,
        bold=True,
        color_hex="#FFFFFF",
        align="center",
        vcenter=True,
    )

    # Satellites
    for i in range(n):
        item = items[i]
        angle = math.radians(-90 + i * (360 / n))
        sx = cx + radius * math.cos(angle)
        sy = cy + radius * math.sin(angle)
        # Satellite circle
        add_oval(slide, sx - sat_size / 2, sy - sat_size / 2, sat_size, sat_size, "#FFFFFF")
        add_oval(slide, sx - sat_size / 2 + 0.1, sy - sat_size / 2 + 0.1, sat_size - 0.2, sat_size - 0.2, primary)
        # Icon glyph
        add_text_box(
            slide,
            sx - sat_size / 2,
            sy - sat_size / 2,
            sat_size,
            sat_size,
            item.get("icon_glyph", glyphs[i % len(glyphs)]),
            size_pt=20,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Label outside satellite (radially outward)
        lx, ly = sx + (sat_size / 2 + 0.5) * math.cos(angle), sy + (sat_size / 2 + 0.5) * math.sin(angle)
        add_text_box(
            slide,
            lx - 2.5,
            ly - 0.4,
            5.0,
            0.8,
            str(item.get("label", ""))[:24],
            size_pt=11,
            bold=True,
            color_hex=ink,
            align="center",
            vcenter=True,
        )


def draw_photo_placeholder(slide, x_cm, y_cm, w_cm, h_cm, primary, accent, label="PHOTO"):
    """Gradient placeholder where a real photo would go.

    Designed so user can later swap with actual PNG via add_picture.
    """
    add_gradient_rect(slide, x_cm, y_cm, w_cm, h_cm, primary, accent, angle=135)
    # Big icon in center
    add_text_box(
        slide,
        x_cm,
        y_cm + h_cm / 2 - 1.5,
        w_cm,
        3.0,
        GLYPHS["data"],
        size_pt=72,
        bold=True,
        color_hex="#FFFFFF",
        align="center",
        vcenter=True,
    )
    # Label at bottom
    add_text_box(
        slide,
        x_cm,
        y_cm + h_cm - 1.0,
        w_cm,
        0.7,
        label,
        size_pt=10,
        bold=True,
        color_hex="#FFFFFF",
        align="center",
        vcenter=True,
    )


def draw_percentage_grid_4(slide, items, primary, accent, ink, muted, light_bg):
    """2x2 grid of large percentage cards with icons (user image style)."""
    glyphs = [GLYPHS["target"], GLYPHS["chart"], GLYPHS["bulb"], GLYPHS["star"]]
    margin_x = 2.0
    margin_y = 5.0
    gap = 0.6
    avail_w = SLIDE_W - 2 * margin_x
    avail_h = SLIDE_H - margin_y - 2.0
    card_w = (avail_w - gap) / 2
    card_h = (avail_h - gap) / 2

    for i, item in enumerate(items[:4]):
        r, c = i // 2, i % 2
        x = margin_x + c * (card_w + gap)
        y = margin_y + r * (card_h + gap)
        # Alternate filled vs outlined cards (like the image)
        filled = i % 3 == 0 or i == 3
        if filled:
            add_rounded_rect(slide, x, y, card_w, card_h, primary)
            num_color = "#FFFFFF"
            label_color = "#E2E8F0"
            icon_bg = "#FFFFFF"
            icon_color = primary
        else:
            add_rounded_rect(slide, x, y, card_w, card_h, "#FFFFFF", line_hex=primary)
            num_color = primary
            label_color = muted
            icon_bg = primary
            icon_color = "#FFFFFF"
        # Big percentage
        add_text_box(
            slide,
            x + 0.5,
            y + 0.3,
            card_w * 0.55,
            card_h - 0.6,
            str(item.get("value", "60%")),
            size_pt=64,
            bold=True,
            color_hex=num_color,
            align="left",
            vcenter=True,
        )
        # Icon top-right
        ic_size = 1.8
        add_oval(slide, x + card_w - ic_size - 0.6, y + 0.6, ic_size, ic_size, icon_bg)
        add_text_box(
            slide,
            x + card_w - ic_size - 0.6,
            y + 0.6,
            ic_size,
            ic_size,
            item.get("icon_glyph", glyphs[i % 4]),
            size_pt=22,
            bold=True,
            color_hex=icon_color,
            align="center",
            vcenter=True,
        )
        # Label below
        add_text_box(
            slide,
            x + 0.5,
            y + card_h - 1.8,
            card_w - 1.0,
            1.5,
            str(item.get("label", ""))[:60],
            size_pt=10,
            color_hex=label_color,
            align="left",
            vcenter=False,
        )


def draw_vs_comparison(slide, left, right, points, primary, accent, ink, muted, light_bg):
    """Center 'VS' circle + 3 points on each side (Samsung Pay vs Apple Pay style).

    left/right: {label, icon_glyph?}
    points: list of {left_text, right_text, label_no}
    """
    # Background light tint
    add_rect(slide, 0, 4.5, SLIDE_W, SLIDE_H - 6.0, "#FAFBFC")
    # Center VS circle
    vs_size = 3.5
    cx, cy = SLIDE_W / 2, 11.0
    add_oval(slide, cx - vs_size / 2 - 0.3, cy - vs_size / 2 - 0.3, vs_size + 0.6, vs_size + 0.6, accent)
    add_oval(slide, cx - vs_size / 2, cy - vs_size / 2, vs_size, vs_size, primary)
    add_text_box(
        slide,
        cx - vs_size / 2,
        cy - vs_size / 2,
        vs_size,
        vs_size,
        "VS",
        size_pt=36,
        bold=True,
        color_hex="#FFFFFF",
        align="center",
        vcenter=True,
    )
    # Left side label
    add_text_box(
        slide,
        1.5,
        5.5,
        SLIDE_W / 2 - 3.0,
        1.2,
        left.get("label", "OPTION A"),
        size_pt=18,
        bold=True,
        color_hex=primary,
        align="center",
        vcenter=True,
    )
    # Right side label
    add_text_box(
        slide,
        SLIDE_W / 2 + 1.5,
        5.5,
        SLIDE_W / 2 - 3.0,
        1.2,
        right.get("label", "OPTION B"),
        size_pt=18,
        bold=True,
        color_hex=primary,
        align="center",
        vcenter=True,
    )
    # 3 comparison points
    for i, p in enumerate(points[:3]):
        y = 7.5 + i * 2.6
        # Left point
        add_text_box(
            slide,
            1.5,
            y,
            SLIDE_W / 2 - 3.0,
            2.2,
            f"{i + 1:02d}.  {p.get('left_text', '')}",
            size_pt=11,
            color_hex=ink,
            align="right",
            vcenter=True,
        )
        # Right point
        add_text_box(
            slide,
            SLIDE_W / 2 + 1.5,
            y,
            SLIDE_W / 2 - 3.0,
            2.2,
            f"{i + 1:02d}.  {p.get('right_text', '')}",
            size_pt=11,
            color_hex=ink,
            align="left",
            vcenter=True,
        )


def draw_as_is_to_be(slide, as_is, to_be, primary, accent, ink, muted, light_bg):
    """AS-IS / TO-BE 2-column comparison (consulting classic).

    as_is/to_be: {label, points: list[str]}
    """
    # Dark background panel
    panel_y = 5.0
    panel_h = SLIDE_H - panel_y - 2.0
    add_rounded_rect(slide, 1.5, panel_y, SLIDE_W - 3.0, panel_h, "#334155")
    # AS-IS column
    col_w = (SLIDE_W - 4.5) / 2
    # AS-IS header
    add_rounded_rect(slide, 2.0, panel_y + 0.5, col_w, 1.3, "#475569")
    add_text_box(
        slide,
        2.0,
        panel_y + 0.5,
        col_w,
        1.3,
        as_is.get("label", "AS-IS"),
        size_pt=22,
        bold=True,
        color_hex="#FFFFFF",
        align="center",
        vcenter=True,
    )
    # AS-IS points
    for i, pt in enumerate(as_is.get("points", [])[:4]):
        y = panel_y + 2.2 + i * 1.4
        # Icon circle
        add_oval(slide, 2.3, y, 0.9, 0.9, "#94A3B8")
        add_text_box(
            slide,
            2.3,
            y,
            0.9,
            0.9,
            GLYPHS["no"],
            size_pt=12,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Text
        add_text_box(
            slide, 3.5, y, col_w - 1.7, 1.2, str(pt)[:70], size_pt=14, color_hex="#E2E8F0", align="left", vcenter=True
        )

    # TO-BE column (highlighted)
    to_x = 2.5 + col_w
    add_rounded_rect(slide, to_x, panel_y + 0.5, col_w, 1.3, primary)
    add_text_box(
        slide,
        to_x,
        panel_y + 0.5,
        col_w,
        1.3,
        to_be.get("label", "TO-BE"),
        size_pt=22,
        bold=True,
        color_hex="#FFFFFF",
        align="center",
        vcenter=True,
    )
    # TO-BE points
    for i, pt in enumerate(to_be.get("points", [])[:4]):
        y = panel_y + 2.2 + i * 1.4
        # Icon circle (accent)
        add_oval(slide, to_x + 0.3, y, 0.9, 0.9, accent)
        add_text_box(
            slide,
            to_x + 0.3,
            y,
            0.9,
            0.9,
            GLYPHS["ok"],
            size_pt=12,
            bold=True,
            color_hex=primary,
            align="center",
            vcenter=True,
        )
        add_text_box(
            slide,
            to_x + 1.5,
            y,
            col_w - 1.7,
            1.2,
            str(pt)[:70],
            size_pt=14,
            bold=True,
            color_hex="#FFFFFF",
            align="left",
            vcenter=True,
        )


def draw_5step_alt_callouts(slide, steps, primary, accent, ink, muted, light_bg):
    """5 nodes in a row + alternating top/bottom captions (user image style)."""
    n = min(len(steps), 5)
    if n == 0:
        return
    glyphs = [GLYPHS["data"], GLYPHS["chart"], GLYPHS["settings"], GLYPHS["target"], GLYPHS["bulb"]]
    margin_x = 2.5
    avail_w = SLIDE_W - 2 * margin_x
    step = avail_w / max(1, n - 1) if n > 1 else avail_w
    y_node = 10.0
    # Horizontal connecting line
    add_rect(slide, margin_x, y_node - 0.05, avail_w, 0.08, "#CBD5E1")

    for i in range(n):
        s = steps[i]
        x = margin_x + i * step if n > 1 else margin_x + avail_w / 2
        # Outer ring + inner circle
        ring = 2.4
        add_oval(slide, x - ring / 2, y_node - ring / 2, ring, ring, accent)
        node = 1.8
        add_oval(slide, x - node / 2, y_node - node / 2, node, node, primary)
        add_text_box(
            slide,
            x - node / 2,
            y_node - node / 2,
            node,
            node,
            s.get("icon_glyph", glyphs[i % len(glyphs)]),
            size_pt=18,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Step number badge
        add_oval(slide, x - 0.6, y_node + node / 2 + 0.2, 1.2, 1.2, "#FFFFFF")
        add_oval(slide, x - 0.5, y_node + node / 2 + 0.3, 1.0, 1.0, primary)
        add_text_box(
            slide,
            x - 0.5,
            y_node + node / 2 + 0.3,
            1.0,
            1.0,
            f"{i + 1:02d}",
            size_pt=11,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Alternating caption (above for even, below for odd)
        title = str(s.get("title", f"STEP {i + 1}"))[:30]
        caption = str(s.get("caption", ""))[:120]
        if i % 2 == 0:
            # Above
            add_text_box(
                slide,
                x - 2.5,
                y_node - 4.2,
                5.0,
                1.0,
                title,
                size_pt=13,
                bold=True,
                color_hex=primary,
                align="center",
                vcenter=True,
            )
            add_text_box(
                slide, x - 2.5, y_node - 3.2, 5.0, 2.0, caption, size_pt=11, color_hex=ink, align="center", vcenter=False
            )
        else:
            # Below (with offset for badge)
            add_text_box(
                slide,
                x - 2.5,
                y_node + 3.0,
                5.0,
                1.0,
                title,
                size_pt=13,
                bold=True,
                color_hex=primary,
                align="center",
                vcenter=True,
            )
            add_text_box(
                slide, x - 2.5, y_node + 4.0, 5.0, 2.0, caption, size_pt=11, color_hex=ink, align="center", vcenter=False
            )


def draw_thank_you_closing(slide, ctx, primary, accent, ink, muted):
    """Full-bleed THANK YOU closing slide (teal template style)."""
    add_gradient_rect(slide, 0, 0, SLIDE_W, SLIDE_H, primary, "#0E7490", angle=135)
    # Big curved decoration on left
    add_oval(slide, -8.0, -8.0, 16.0, 16.0, accent)
    add_oval(slide, -4.0, SLIDE_H - 8.0, 12.0, 12.0, "#FFFFFF")
    add_oval(slide, -3.0, SLIDE_H - 7.0, 10.0, 10.0, primary)
    # Big "Thank You" title
    add_text_box(
        slide,
        SLIDE_W * 0.35,
        SLIDE_H * 0.35,
        SLIDE_W * 0.55,
        4.0,
        "Thank You",
        size_pt=96,
        bold=True,
        color_hex="#FFFFFF",
        align="left",
        vcenter=True,
    )
    # Subtitle
    add_text_box(
        slide,
        SLIDE_W * 0.35,
        SLIDE_H * 0.55,
        SLIDE_W * 0.55,
        1.5,
        "감사합니다. 질문 받겠습니다.",
        size_pt=20,
        color_hex="#FFFFFF",
        align="left",
        vcenter=True,
    )
    # Accent line
    add_rect(slide, SLIDE_W * 0.35, SLIDE_H * 0.66, 6.0, 0.15, accent)
    # Footer info — jh 2026-06-12: 게이트 태그 suffix 제거 (S20 푸터가
    # "(주제: ... 데이터로 " 처럼 중간 절단되던 결함, cover 와 동일 규칙)
    import re as _re

    intent = (ctx.meta.user_intent or "분석 보고서").strip()
    _m = _re.match(r"^(?:분석 방향|주제|방법론|모델 전략):\s*(.+)$", intent)
    if _m:
        intent = _m.group(1).strip()
    _cut = _re.search(r"\s*\((?:분석 방향|주제|방법론|모델 전략):", intent)
    if _cut:
        intent = intent[: _cut.start()].strip()
    if len(intent) > 50:
        intent = intent[:49].rstrip() + "…"
    add_text_box(
        slide,
        SLIDE_W * 0.35,
        SLIDE_H * 0.72,
        SLIDE_W * 0.55,
        1.0,
        f"본 보고서  ·  {intent}",
        size_pt=12,
        color_hex="#FFFFFF",
        align="left",
        vcenter=True,
    )
    add_text_box(
        slide,
        SLIDE_W * 0.35,
        SLIDE_H * 0.77,
        SLIDE_W * 0.55,
        1.0,
        f"생성  ·  {(ctx.meta.generated_at or '')[:10]}  ·  ADA",
        size_pt=11,
        color_hex="#E2E8F0",
        align="left",
        vcenter=True,
    )


def draw_checklist_4(slide, items, primary, accent, ink, muted, light_bg):
    """4 horizontal checklist cards with big check circles (teal template style)."""
    n = min(len(items), 4)
    if n == 0:
        return
    margin_x = 2.0
    avail_w = SLIDE_W - 2 * margin_x
    card_w = avail_w / n - 0.6
    gap = 0.6
    y = 5.5

    for i in range(n):
        item = items[i]
        x = margin_x + i * (card_w + gap)
        # Big check circle (top)
        circle_size = 3.0
        cx = x + card_w / 2
        add_oval(slide, cx - circle_size / 2 - 0.25, y - 0.25, circle_size + 0.5, circle_size + 0.5, accent)
        add_oval(slide, cx - circle_size / 2, y, circle_size, circle_size, primary)
        # Big check mark
        add_text_box(
            slide,
            cx - circle_size / 2,
            y,
            circle_size,
            circle_size,
            GLYPHS["ok"],
            size_pt=44,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Title below
        add_text_box(
            slide,
            x,
            y + circle_size + 0.6,
            card_w,
            1.2,
            str(item.get("title", f"항목 {i + 1}"))[:30],
            size_pt=14,
            bold=True,
            color_hex=primary,
            align="center",
            vcenter=True,
        )
        # Accent line
        add_rect(slide, cx - 1.0, y + circle_size + 1.9, 2.0, 0.05, accent)
        # Caption
        add_text_box(
            slide,
            x,
            y + circle_size + 2.2,
            card_w,
            SLIDE_H - y - circle_size - 4.5,
            str(item.get("caption", ""))[:120],
            size_pt=10,
            color_hex=ink,
            align="center",
            vcenter=False,
        )


def draw_step_cards_vertical(slide, steps, primary, accent, ink, muted, light_bg):
    """Vertical STEP 01/02/03/04 cards (teal template style)."""
    n = min(len(steps), 4)
    if n == 0:
        return
    margin_x = 1.5
    avail_w = SLIDE_W - 2 * margin_x
    card_w = avail_w / n - 0.4
    gap = 0.4
    y = 5.0
    card_h = SLIDE_H - y - 2.0

    for i in range(n):
        s = steps[i]
        x = margin_x + i * (card_w + gap)
        # Card with gradient
        add_gradient_rect(slide, x, y, card_w, card_h, primary, "#0E7490" if i % 2 == 0 else "#0F766E", angle=180)
        # STEP label top
        add_text_box(
            slide,
            x,
            y + 0.8,
            card_w,
            1.2,
            "STEP",
            size_pt=12,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Big number
        add_text_box(
            slide,
            x,
            y + 1.5,
            card_w,
            3.0,
            f"{i + 1:02d}",
            size_pt=72,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Accent line
        add_rect(slide, x + card_w / 2 - 1.5, y + 4.8, 3.0, 0.1, accent)
        # Title
        add_text_box(
            slide,
            x + 0.4,
            y + 5.5,
            card_w - 0.8,
            1.5,
            str(s.get("title", f"단계 {i + 1}"))[:30],
            size_pt=14,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Caption
        add_text_box(
            slide,
            x + 0.4,
            y + 7.2,
            card_w - 0.8,
            card_h - 8.0,
            str(s.get("caption", ""))[:120],
            size_pt=10,
            color_hex="#E2E8F0",
            align="center",
            vcenter=False,
        )


def draw_qna_slide(slide, ctx, primary, accent, ink, muted, light_bg):
    """Q&A slide with left big Q icon + right question/answer boxes."""
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, "#FFFFFF")
    # Left side - big Q&A circle
    add_oval(slide, 2.0, 4.0, 12.0, 12.0, light_bg)
    add_oval(slide, 3.0, 5.0, 10.0, 10.0, primary)
    add_text_box(
        slide, 3.0, 5.0, 10.0, 10.0, "Q&A", size_pt=96, bold=True, color_hex="#FFFFFF", align="center", vcenter=True
    )
    # Right side - 3 question boxes
    rx = 16.0
    rw = SLIDE_W - rx - 1.5
    titles = ["Question 01", "Question 02", "Question 03"]
    for i, t in enumerate(titles):
        y = 5.0 + i * 3.2
        add_rounded_rect(slide, rx, y, rw, 2.6, light_bg, line_hex=primary)
        # Number badge
        add_oval(slide, rx - 0.8, y + 0.8, 1.6, 1.6, primary)
        add_text_box(
            slide,
            rx - 0.8,
            y + 0.8,
            1.6,
            1.6,
            f"{i + 1:02d}",
            size_pt=14,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Title
        add_text_box(
            slide,
            rx + 1.5,
            y + 0.4,
            rw - 2.0,
            1.0,
            t,
            size_pt=14,
            bold=True,
            color_hex=primary,
            align="left",
            vcenter=True,
        )
        # Body placeholder
        add_text_box(
            slide,
            rx + 1.5,
            y + 1.3,
            rw - 2.0,
            1.2,
            "본 보고서 결론·근거·실행안에 대한 추가 질문 환영",
            size_pt=10,
            color_hex=ink,
            align="left",
            vcenter=True,
        )


def draw_chevron_strategies(slide, items, primary, accent, ink, muted, light_bg):
    """5 chevron arrows in sequence (Strategy 1->5 style)."""
    n = min(len(items), 5)
    if n == 0:
        return
    margin_x = 1.8
    avail_w = SLIDE_W - 2 * margin_x
    chevron_w = avail_w / n
    y_chev = 7.0
    chevron_h = 3.0
    overlap = 0.8

    for i in range(n):
        item = items[i]
        x = margin_x + i * (chevron_w - overlap)
        color = primary if i % 2 == 0 else accent
        text_color = "#FFFFFF" if i % 2 == 0 else primary
        # Chevron arrow
        add_chevron(slide, x, y_chev, chevron_w, chevron_h, color)
        # Strategy title
        add_text_box(
            slide,
            x + 0.3,
            y_chev + 0.6,
            chevron_w - 0.8,
            1.0,
            f"STRATEGY {i + 1:02d}",
            size_pt=11,
            bold=True,
            color_hex=text_color,
            align="center",
            vcenter=True,
        )
        # Body
        add_text_box(
            slide,
            x + 0.3,
            y_chev + 1.6,
            chevron_w - 0.8,
            1.0,
            str(item.get("title", "")).strip()[:24],
            size_pt=16,
            bold=True,
            color_hex=text_color,
            align="center",
            vcenter=True,
        )
        # Caption below chevron
        add_text_box(
            slide,
            x + 0.3,
            y_chev + chevron_h + 0.5,
            chevron_w - 1.0,
            4.0,
            str(item.get("caption", ""))[:120],
            size_pt=12,
            color_hex=ink,
            align="center",
            vcenter=False,
        )


def draw_index_cards(slide, items, primary, accent, ink, muted, light_bg):
    """Big number Index cards (Index 01/02/03/04 style)."""
    n = min(len(items), 4)
    if n == 0:
        return
    margin_x = 2.0
    avail_w = SLIDE_W - 2 * margin_x
    gap = 0.6
    card_w = (avail_w - (n - 1) * gap) / n
    y = 5.5
    card_h = SLIDE_H - y - 2.0

    for i in range(n):
        item = items[i]
        x = margin_x + i * (card_w + gap)
        # Top "INDEX" label band
        add_rect(slide, x, y, card_w, 1.0, primary)
        add_text_box(
            slide, x, y, card_w, 1.0, "INDEX", size_pt=11, bold=True, color_hex="#FFFFFF", align="center", vcenter=True
        )
        # White body card
        add_rect(slide, x, y + 1.0, card_w, card_h - 1.0, "#FFFFFF")
        add_rect(slide, x, y + 1.0, card_w, card_h - 1.0, light_bg, line=False)
        add_rounded_rect(slide, x + 0.1, y + 1.1, card_w - 0.2, card_h - 1.2, "#FFFFFF", line_hex="#CBD5E1")
        # Big number
        add_text_box(
            slide,
            x,
            y + 2.5,
            card_w,
            4.0,
            f"{i + 1:02d}",
            size_pt=96,
            bold=True,
            color_hex=primary,
            align="center",
            vcenter=True,
        )
        # Accent divider
        add_rect(slide, x + card_w / 2 - 1.0, y + 7.0, 2.0, 0.08, accent)
        # Title
        add_text_box(
            slide,
            x + 0.4,
            y + 7.5,
            card_w - 0.8,
            1.5,
            str(item.get("title", f"항목 {i + 1}"))[:24],
            size_pt=14,
            bold=True,
            color_hex=ink,
            align="center",
            vcenter=True,
        )
        # Caption
        add_text_box(
            slide,
            x + 0.4,
            y + 9.2,
            card_w - 0.8,
            card_h - 9.8,
            str(item.get("caption", ""))[:100],
            size_pt=9,
            color_hex=muted,
            align="center",
            vcenter=False,
        )


def draw_big_stats(slide, items, primary, accent, ink, muted, light_bg):
    """2-3 BIG numbers split layout (18,000 / 17,500 style)."""
    n = min(len(items), 3)
    if n == 0:
        return
    margin_x = 1.5
    avail_w = SLIDE_W - 2 * margin_x
    card_w = avail_w / n
    y = 4.5
    card_h = SLIDE_H - y - 2.0

    for i in range(n):
        item = items[i]
        x = margin_x + i * card_w
        # Alternate solid teal / light bg
        is_solid = i % 2 == 1
        bg = primary if is_solid else light_bg
        num_color = "#FFFFFF" if is_solid else primary
        label_color = "#E2E8F0" if is_solid else muted
        add_rect(slide, x, y, card_w, card_h, bg)
        # Top thin line
        add_rect(slide, x + 1.0, y + 1.5, card_w - 2.0, 0.06, accent if is_solid else primary)
        # HUGE number — jh 2026-06-11: 길이별 자동 폰트.
        # 고정 128pt 는 "79%" 같은 3자도 카드 폭 초과 줄바꿈 (probe 실측) →
        # 카드 폭 (avail/n) 과 글자 수에서 안전 폰트 도출.
        _val = str(item.get("value", "0"))
        # jh 2026-06-12 — 견본 실측 120pt(3자) 기준 상향.
        _size = {1: 128, 2: 128, 3: 120, 4: 96, 5: 80, 6: 66}.get(len(_val), 56)
        add_text_box(
            slide,
            x,
            y + 2.5,
            card_w,
            6.0,
            _val,
            size_pt=_size,
            bold=True,
            color_hex=num_color,
            align="center",
            vcenter=True,
        )
        # Accent line under
        add_rect(slide, x + card_w / 2 - 1.5, y + 9.0, 3.0, 0.1, accent if is_solid else primary)
        # Label
        add_text_box(
            slide,
            x + 0.5,
            y + 9.5,
            card_w - 1.0,
            1.5,
            str(item.get("label", "")).upper()[:40],
            size_pt=16,
            bold=True,
            color_hex=label_color,
            align="center",
            vcenter=True,
        )
        # Caption
        add_text_box(
            slide,
            x + 0.5,
            y + 11.0,
            card_w - 1.0,
            card_h - 11.5,
            str(item.get("caption", ""))[:120],
            size_pt=13,
            color_hex=label_color,
            align="center",
            vcenter=False,
        )


def draw_strategy_4_circular(slide, items, primary, accent, ink, muted, light_bg):
    """4 strategy circles around central arrow product (Strategy 1-4 style)."""
    cx, cy = SLIDE_W / 2, 11.0
    # Central product circle
    add_oval(slide, cx - 1.6, cy - 1.6, 3.2, 3.2, accent)
    add_oval(slide, cx - 1.4, cy - 1.4, 2.8, 2.8, primary)
    add_text_box(
        slide,
        cx - 1.4,
        cy - 1.4,
        2.8,
        2.8,
        "PRODUCT",
        size_pt=12,
        bold=True,
        color_hex="#FFFFFF",
        align="center",
        vcenter=True,
    )
    # 4 strategy circles at NW, NE, SW, SE
    positions = [(-6.0, -4.5), (6.0, -4.5), (-6.0, 4.5), (6.0, 4.5)]
    glyphs = [GLYPHS["data"], GLYPHS["chart"], GLYPHS["target"], GLYPHS["bulb"]]
    for i, (dx, dy) in enumerate(positions[: len(items)]):
        item = items[i]
        sx = cx + dx
        sy = cy + dy
        # Circle
        sat_size = 3.0
        add_oval(slide, sx - sat_size / 2 - 0.2, sy - sat_size / 2 - 0.2, sat_size + 0.4, sat_size + 0.4, accent)
        add_oval(slide, sx - sat_size / 2, sy - sat_size / 2, sat_size, sat_size, primary)
        # Icon
        add_text_box(
            slide,
            sx - sat_size / 2,
            sy - sat_size / 2 - 0.3,
            sat_size,
            sat_size * 0.5,
            glyphs[i % 4],
            size_pt=18,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Label
        add_text_box(
            slide,
            sx - sat_size / 2,
            sy + 0.0,
            sat_size,
            sat_size * 0.5,
            f"STRATEGY {i + 1}",
            size_pt=10,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Caption outside (further from center)
        cap_offset_x = sat_size / 2 + 0.3 if dx > 0 else -(sat_size / 2 + 4.5)
        add_text_box(
            slide,
            sx + cap_offset_x,
            sy - 0.8,
            4.0,
            1.6,
            str(item.get("title", ""))[:24],
            size_pt=12,
            bold=True,
            color_hex=primary,
            align="left" if dx > 0 else "right",
            vcenter=True,
        )


def _shade(hex_color: str, factor: float) -> str:
    """Darken (factor<1) or lighten (factor>1) a hex color."""
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    if factor < 1:
        r, g, b = int(r * factor), int(g * factor), int(b * factor)
    else:
        r = min(255, int(r + (255 - r) * (factor - 1)))
        g = min(255, int(g + (255 - g) * (factor - 1)))
        b = min(255, int(b + (255 - b) * (factor - 1)))
    return f"#{r:02x}{g:02x}{b:02x}"


def draw_cube_3d(slide, items, primary, accent, ink, muted, light_bg):
    """3-4 fake-3D cubes (3 faces: front + top + side)."""
    n = min(len(items), 4)
    if n == 0:
        return
    margin_x = 2.0
    avail_w = SLIDE_W - 2 * margin_x
    cube_w = avail_w / n - 0.8
    cube_size = min(cube_w, 5.0)
    y = 6.0

    primary_dark = _shade(primary, 0.7)
    primary_light = _shade(primary, 1.3)

    for i in range(n):
        item = items[i]
        x = margin_x + i * (cube_w + 0.8)
        # Cube faces - depth offset
        depth = cube_size * 0.25
        # Top face (parallelogram - approximated with triangle + rect)
        add_triangle(slide, x + depth * 0.5, y - depth * 0.5, cube_size, depth, primary_light, direction="diamond")
        # Right side face
        add_rect(slide, x + cube_size, y, depth, cube_size, primary_dark)
        # Front face (main)
        add_rect(slide, x, y, cube_size, cube_size, primary)
        # Number/icon on front face
        add_text_box(
            slide,
            x,
            y,
            cube_size,
            cube_size * 0.5,
            f"{i + 1:02d}",
            size_pt=42,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Label inside front face (bottom)
        add_text_box(
            slide,
            x + 0.3,
            y + cube_size * 0.55,
            cube_size - 0.6,
            cube_size * 0.4,
            str(item.get("title", ""))[:24],
            size_pt=11,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Caption below cube
        add_text_box(
            slide,
            x - 0.4,
            y + cube_size + depth + 0.4,
            cube_size + 0.8,
            4.0,
            str(item.get("caption", ""))[:100],
            size_pt=10,
            color_hex=ink,
            align="center",
            vcenter=False,
        )


def draw_donut_metric(slide, value_pct, label, caption, primary, accent, ink, muted, light_bg, side_items=None):
    """Big donut chart with center % + side bullet list.

    Renders donut via matplotlib then embeds. side_items: list[str].
    """
    import tempfile

    try:
        import matplotlib

        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception:
        return

    # Donut PNG via matplotlib
    pct = max(0.0, min(100.0, float(value_pct)))
    fig, ax = plt.subplots(figsize=(6, 6), dpi=130, subplot_kw={"aspect": "equal"})
    fig.patch.set_alpha(0)
    ax.set_facecolor("none")
    wedges, _ = ax.pie(
        [pct, 100 - pct],
        colors=[primary, "#E2E8F0"],
        startangle=90,
        wedgeprops={"width": 0.32, "edgecolor": "white", "linewidth": 3},
    )
    ax.text(0, 0.05, f"{pct:.0f}%", ha="center", va="center", fontsize=58, fontweight="bold", color=primary)
    ax.text(0, -0.25, label.upper()[:24], ha="center", va="center", fontsize=12, color="#64748B")
    ax.axis("off")
    out = tempfile.NamedTemporaryFile(suffix=".png", delete=False).name
    plt.savefig(out, dpi=140, bbox_inches="tight", transparent=True)
    plt.close(fig)

    # Embed in slide left side
    from pptx.util import Cm

    donut_size = 10.0
    donut_x = 3.0
    donut_y = 5.0
    try:
        slide.shapes.add_picture(out, Cm(donut_x), Cm(donut_y), width=Cm(donut_size), height=Cm(donut_size))
    except Exception:
        pass

    # Right side - title + bullets
    text_x = donut_x + donut_size + 1.5
    text_w = SLIDE_W - text_x - 1.5
    add_text_box(
        slide,
        text_x,
        5.5,
        text_w,
        1.5,
        caption.upper()[:60] if caption else "KEY METRIC",
        size_pt=13,
        bold=True,
        color_hex=primary,
        align="left",
        vcenter=False,
    )
    add_rect(slide, text_x, 7.0, 3.0, 0.1, accent)
    if side_items:
        for i, item in enumerate(side_items[:5]):
            y = 7.5 + i * 1.5
            # Bullet circle
            add_oval(slide, text_x, y + 0.2, 0.7, 0.7, primary)
            add_text_box(
                slide,
                text_x,
                y + 0.2,
                0.7,
                0.7,
                f"{i + 1:02d}",
                size_pt=9,
                bold=True,
                color_hex="#FFFFFF",
                align="center",
                vcenter=True,
            )
            # Text
            add_text_box(
                slide,
                text_x + 1.0,
                y,
                text_w - 1.0,
                1.2,
                str(item)[:90],
                size_pt=11,
                color_hex=ink,
                align="left",
                vcenter=True,
            )


def draw_numbered_rows(slide, items, primary, accent, ink, muted, light_bg):
    """Horizontal numbered rows: '01  Title  ----  Caption'."""
    n = min(len(items), 5)
    if n == 0:
        return
    margin_x = 1.5
    margin_y = 5.0
    avail_h = SLIDE_H - margin_y - 2.0
    row_h = avail_h / n - 0.3

    for i in range(n):
        item = items[i]
        y = margin_y + i * (row_h + 0.3)
        # Background card (alternate)
        bg = light_bg if i % 2 == 0 else "#FFFFFF"
        add_rounded_rect(slide, margin_x, y, SLIDE_W - 2 * margin_x, row_h, bg)
        # Big number on left
        add_text_box(
            slide,
            margin_x + 0.5,
            y,
            3.5,
            row_h,
            f"{i + 1:02d}",
            size_pt=44,
            bold=True,
            color_hex=primary,
            align="center",
            vcenter=True,
        )
        # Vertical accent line
        add_rect(slide, margin_x + 4.5, y + 0.5, 0.1, row_h - 1.0, accent)
        # Title
        add_text_box(
            slide,
            margin_x + 5.0,
            y + 0.3,
            9.0,
            row_h - 0.6,
            str(item.get("title", ""))[:30],
            size_pt=18,
            bold=True,
            color_hex=ink,
            align="left",
            vcenter=True,
        )
        # Caption (right)
        add_text_box(
            slide,
            margin_x + 14.5,
            y + 0.3,
            SLIDE_W - margin_x - 15.5,
            row_h - 0.6,
            str(item.get("caption", ""))[:140],
            size_pt=11,
            color_hex=muted,
            align="left",
            vcenter=True,
        )


def draw_research_cover(slide, ctx, primary, accent, ink, muted):
    """Research-style cover with big year + contact box (alternative cover)."""
    # Top primary band
    add_rect(slide, 0, 0, SLIDE_W, 1.0, primary)
    # Right side decorative wave (oval mass)
    add_oval(slide, SLIDE_W - 14.0, -6.0, 22.0, 22.0, _shade(primary, 1.4))
    add_oval(slide, SLIDE_W - 10.0, -3.0, 16.0, 16.0, _shade(primary, 1.2))
    # Big year
    year = (ctx.meta.generated_at or "")[:4] or "2026"
    add_text_box(
        slide, 2.0, 4.0, 12.0, 4.0, year, size_pt=120, bold=True, color_hex=primary, align="left", vcenter=False
    )
    # Subtitle
    add_text_box(
        slide,
        2.0,
        8.5,
        18.0,
        2.0,
        (ctx.meta.user_intent or "Analysis Report")[:60],
        size_pt=28,
        bold=True,
        color_hex=ink,
        align="left",
        vcenter=False,
    )
    # Accent line
    add_rect(slide, 2.0, 11.0, 6.0, 0.15, accent)
    # Contact box bottom
    add_rect(slide, 2.0, 14.0, 8.0, 3.0, primary)
    add_text_box(
        slide, 2.3, 14.2, 7.4, 0.6, "CONTACT", size_pt=10, bold=True, color_hex="#FFFFFF", align="left", vcenter=True
    )
    chosen = (ctx.model_selection.chosen or {}).get("name", "-")
    add_text_box(
        slide, 2.3, 15.0, 7.4, 0.8, f"Model · {chosen}", size_pt=11, color_hex="#FFFFFF", align="left", vcenter=True
    )
    add_text_box(
        slide,
        2.3,
        15.8,
        7.4,
        0.8,
        f"Classification · {ctx.meta.classification}",
        size_pt=11,
        color_hex="#FFFFFF",
        align="left",
        vcenter=True,
    )


def draw_linked_circles_4(slide, items, primary, accent, ink, muted, light_bg):
    """4 overlapping circles in a row with shared center hub."""
    n = min(len(items), 4)
    if n == 0:
        return
    cy = 10.0
    circle_d = 5.0
    overlap = 1.2
    total_w = n * circle_d - (n - 1) * overlap
    start_x = (SLIDE_W - total_w) / 2
    glyphs = [GLYPHS["data"], GLYPHS["chart"], GLYPHS["target"], GLYPHS["bulb"]]

    for i in range(n):
        item = items[i]
        cx = start_x + i * (circle_d - overlap) + circle_d / 2
        # Outer ring
        add_oval(slide, cx - circle_d / 2 - 0.2, cy - circle_d / 2 - 0.2, circle_d + 0.4, circle_d + 0.4, accent)
        # Inner color (alternate primary / lighter)
        c = primary if i % 2 == 0 else _shade(primary, 1.3)
        add_oval(slide, cx - circle_d / 2, cy - circle_d / 2, circle_d, circle_d, c)
        # Number badge top
        add_oval(slide, cx - 0.8, cy - circle_d / 2 - 1.0, 1.6, 1.6, primary)
        add_text_box(
            slide,
            cx - 0.8,
            cy - circle_d / 2 - 1.0,
            1.6,
            1.6,
            f"{i + 1:02d}",
            size_pt=14,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Icon
        add_text_box(
            slide,
            cx - circle_d / 2,
            cy - 0.5,
            circle_d,
            1.2,
            glyphs[i % len(glyphs)],
            size_pt=22,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Title
        add_text_box(
            slide,
            cx - circle_d / 2,
            cy + 0.5,
            circle_d,
            1.2,
            str(item.get("title", ""))[:24],
            size_pt=14,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Caption below
        add_text_box(
            slide,
            cx - circle_d / 2,
            cy + circle_d / 2 + 0.6,
            circle_d,
            2.5,
            str(item.get("caption", ""))[:90],
            size_pt=11,
            color_hex=ink,
            align="center",
            vcenter=False,
        )


def draw_vertical_arrow_steps(slide, items, primary, accent, ink, muted, light_bg):
    """Big upward arrow on left + numbered step cards on right."""
    n = min(len(items), 5)
    if n == 0:
        return
    # Big up arrow on left (chevron pointing up = isosceles triangle on top of rectangle)
    arrow_x = 2.5
    arrow_w = 3.0
    arrow_y_top = 4.5
    arrow_h_total = SLIDE_H - 7.0
    # Triangle head
    add_triangle(slide, arrow_x, arrow_y_top, arrow_w, arrow_w, primary, direction="up")
    # Stem rectangle
    add_rect(
        slide,
        arrow_x + arrow_w * 0.25,
        arrow_y_top + arrow_w * 0.7,
        arrow_w * 0.5,
        arrow_h_total - arrow_w * 0.7,
        primary,
    )

    # Step cards on right
    card_x = 8.0
    card_w = SLIDE_W - card_x - 2.0
    avail_h = SLIDE_H - 6.0
    card_h = avail_h / n - 0.4

    for i in range(n):
        item = items[i]
        y = 5.0 + i * (card_h + 0.4)
        # Card background
        add_rounded_rect(slide, card_x, y, card_w, card_h, _shade(primary, 1.5 - i * 0.08))
        # Number on left
        add_text_box(
            slide,
            card_x + 0.5,
            y,
            2.5,
            card_h,
            f"{i + 1:02d}",
            size_pt=32,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Vertical divider
        add_rect(slide, card_x + 3.2, y + 0.4, 0.06, card_h - 0.8, "#FFFFFF")
        # Title + caption
        add_text_box(
            slide,
            card_x + 3.7,
            y + 0.2,
            card_w - 4.0,
            card_h * 0.4,
            str(item.get("title", f"Step {i + 1}"))[:30],
            size_pt=14,
            bold=True,
            color_hex="#FFFFFF",
            align="left",
            vcenter=True,
        )
        add_text_box(
            slide,
            card_x + 3.7,
            y + card_h * 0.45,
            card_w - 4.0,
            card_h * 0.5,
            str(item.get("caption", ""))[:120],
            size_pt=10,
            color_hex="#FFFFFF",
            align="left",
            vcenter=False,
        )


def draw_clock_diagram(slide, items, primary, accent, ink, muted, light_bg):
    """Clock face with hour markers as 4 phase labels."""
    import math

    n = min(len(items), 4)
    if n == 0:
        return
    cx, cy = SLIDE_W / 2, 11.0
    clock_size = 9.0
    # Outer ring
    add_oval(slide, cx - clock_size / 2 - 0.3, cy - clock_size / 2 - 0.3, clock_size + 0.6, clock_size + 0.6, accent)
    # Clock face
    add_oval(slide, cx - clock_size / 2, cy - clock_size / 2, clock_size, clock_size, "#FFFFFF")
    # Inner ring
    add_oval(slide, cx - clock_size / 2 + 0.4, cy - clock_size / 2 + 0.4, clock_size - 0.8, clock_size - 0.8, light_bg)
    # Hour markers (4 at 12/3/6/9)
    glyphs = [GLYPHS["bulb"], GLYPHS["chart"], GLYPHS["target"], GLYPHS["ok"]]
    angles = [-90, 0, 90, 180]  # 12 / 3 / 6 / 9
    for i in range(n):
        a = math.radians(angles[i])
        mx = cx + (clock_size / 2 - 1.0) * math.cos(a)
        my = cy + (clock_size / 2 - 1.0) * math.sin(a)
        # Marker circle
        add_oval(slide, mx - 0.9, my - 0.9, 1.8, 1.8, primary)
        add_text_box(
            slide,
            mx - 0.9,
            my - 0.9,
            1.8,
            1.8,
            glyphs[i],
            size_pt=16,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
    # Center hands (two lines as rectangles)
    add_rect(slide, cx - 0.06, cy - clock_size * 0.3, 0.12, clock_size * 0.3, primary)
    add_rect(slide, cx, cy - 0.06, clock_size * 0.25, 0.12, primary)
    add_oval(slide, cx - 0.5, cy - 0.5, 1.0, 1.0, primary)
    # Phase labels around clock (outside)
    label_positions = [
        (cx - 2.5, cy - clock_size / 2 - 2.0),  # 12 top
        (cx + clock_size / 2 + 0.5, cy - 0.6),  # 3 right
        (cx - 2.5, cy + clock_size / 2 + 0.5),  # 6 bottom
        (cx - clock_size / 2 - 5.5, cy - 0.6),  # 9 left
    ]
    for i in range(n):
        item = items[i]
        lx, ly = label_positions[i]
        add_text_box(
            slide,
            lx,
            ly,
            5.0,
            1.2,
            str(item.get("title", f"P{i + 1}"))[:24],
            size_pt=12,
            bold=True,
            color_hex=primary,
            align="center",
            vcenter=True,
        )


def draw_gear_infographic(slide, items, primary, accent, ink, muted, light_bg):
    """Big central gear + surrounding bullets (rotational gear visualization)."""
    import math

    n = min(len(items), 6)
    if n == 0:
        return
    cx, cy = SLIDE_W * 0.35, 11.0
    # Outer gear teeth (8 small rectangles around)
    teeth_r = 5.5
    for t in range(12):
        ang = math.radians(t * 30)
        tx = cx + teeth_r * math.cos(ang) - 0.4
        ty = cy + teeth_r * math.sin(ang) - 0.4
        add_rect(slide, tx, ty, 0.8, 0.8, primary)
    # Main gear circle
    main_r = 4.8
    add_oval(slide, cx - main_r, cy - main_r, main_r * 2, main_r * 2, primary)
    # Inner ring
    inner_r = 3.5
    add_oval(slide, cx - inner_r, cy - inner_r, inner_r * 2, inner_r * 2, "#FFFFFF")
    # Center hub
    hub_r = 2.0
    add_oval(slide, cx - hub_r, cy - hub_r, hub_r * 2, hub_r * 2, accent)
    add_oval(slide, cx - hub_r + 0.2, cy - hub_r + 0.2, (hub_r - 0.2) * 2, (hub_r - 0.2) * 2, primary)
    # Center label
    add_text_box(
        slide,
        cx - hub_r,
        cy - hub_r,
        hub_r * 2,
        hub_r * 2,
        GLYPHS["settings"],
        size_pt=28,
        bold=True,
        color_hex="#FFFFFF",
        align="center",
        vcenter=True,
    )

    # Right side - numbered bullets
    bx = SLIDE_W * 0.62
    bw = SLIDE_W - bx - 1.5
    for i in range(n):
        item = items[i]
        y = 5.5 + i * (1.8 if n > 4 else 2.2)
        # Number circle
        add_oval(slide, bx, y, 1.4, 1.4, primary)
        add_text_box(
            slide,
            bx,
            y,
            1.4,
            1.4,
            f"{i + 1:02d}",
            size_pt=12,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Title
        add_text_box(
            slide,
            bx + 1.7,
            y + 0.1,
            bw - 1.7,
            0.8,
            str(item.get("title", f"항목 {i + 1}"))[:30],
            size_pt=13,
            bold=True,
            color_hex=primary,
            align="left",
            vcenter=True,
        )
        # Caption
        add_text_box(
            slide,
            bx + 1.7,
            y + 0.9,
            bw - 1.7,
            0.8,
            str(item.get("caption", ""))[:80],
            size_pt=10,
            color_hex=ink,
            align="left",
            vcenter=True,
        )


def draw_photo_overlay_opener(slide, ctx, title, subtitle, primary, accent, ink, muted, light_bg):
    """Full-bleed gradient (as photo placeholder) + overlay text + corner accent."""
    # Full-bleed dark gradient (acts as photo bg substitute)
    add_gradient_rect(slide, 0, 0, SLIDE_W, SLIDE_H, _shade(primary, 0.6), primary, angle=135)
    # Decorative corner triangles
    add_triangle(slide, -2, -2, 8, 8, accent, direction="diamond")
    add_triangle(slide, SLIDE_W - 6, SLIDE_H - 6, 8, 8, _shade(primary, 0.4), direction="diamond")
    # Dark overlay for text legibility
    add_rect(slide, 0, SLIDE_H * 0.35, SLIDE_W, SLIDE_H * 0.4, "#000000")
    # Title
    add_text_box(
        slide,
        2.0,
        SLIDE_H * 0.4,
        SLIDE_W - 4.0,
        4.0,
        title[:50],
        size_pt=72,
        bold=True,
        color_hex="#FFFFFF",
        align="left",
        vcenter=True,
    )
    # Accent line
    add_rect(slide, 2.0, SLIDE_H * 0.58, 6.0, 0.18, accent)
    # Subtitle
    add_text_box(
        slide,
        2.0,
        SLIDE_H * 0.62,
        SLIDE_W - 4.0,
        2.0,
        subtitle[:100],
        size_pt=18,
        color_hex="#E2E8F0",
        align="left",
        vcenter=False,
    )
    # Category tag bottom
    add_rounded_rect(slide, 2.0, SLIDE_H - 3.0, 6.0, 1.0, accent)
    add_text_box(
        slide,
        2.0,
        SLIDE_H - 3.0,
        6.0,
        1.0,
        ctx.meta.category.upper(),
        size_pt=11,
        bold=True,
        color_hex=primary,
        align="center",
        vcenter=True,
    )


def draw_mini_stats_row(slide, items, primary, accent, ink, muted, light_bg):
    """4 small horizontal stat cards (BOOK PPT style)."""
    n = min(len(items), 4)
    if n == 0:
        return
    margin_x = 2.0
    avail_w = SLIDE_W - 2 * margin_x
    gap = 0.4
    card_w = (avail_w - (n - 1) * gap) / n
    y = 6.0
    card_h = 5.0
    for i in range(n):
        item = items[i]
        x = margin_x + i * (card_w + gap)
        # Card with thin top border
        add_rounded_rect(slide, x, y, card_w, card_h, "#FFFFFF", line_hex="#E2E8F0")
        add_rect(slide, x, y, card_w, 0.15, primary)
        # Big number
        add_text_box(
            slide,
            x,
            y + 0.5,
            card_w,
            2.5,
            str(item.get("value", "—")),
            size_pt=42,
            bold=True,
            color_hex=primary,
            align="center",
            vcenter=True,
        )
        # Accent line
        add_rect(slide, x + card_w / 2 - 1.0, y + 3.0, 2.0, 0.05, accent)
        # Label
        add_text_box(
            slide,
            x + 0.3,
            y + 3.3,
            card_w - 0.6,
            1.5,
            str(item.get("label", ""))[:24],
            size_pt=10,
            color_hex=muted,
            align="center",
            vcenter=False,
        )


def draw_swot_matrix(slide, items, primary, accent, ink, muted, light_bg):
    """SWOT 2x2 matrix - S/W/O/T quadrants with colored panels."""
    colors = ["#16A34A", "#DC2626", "#2563EB", "#D97706"]
    labels = ["S", "W", "O", "T"]
    names = ["Strengths", "Weaknesses", "Opportunities", "Threats"]
    glyphs = [GLYPHS["ok"], GLYPHS["no"], GLYPHS["bulb"], GLYPHS["warn"]]
    cell_w = (SLIDE_W - 4.0) / 2
    cell_h = (SLIDE_H - 7.0) / 2
    y0 = 5.0
    x0 = 2.0
    gap = 0.3
    for i in range(4):
        r, c = i // 2, i % 2
        x = x0 + c * (cell_w + gap)
        y = y0 + r * (cell_h + gap)
        add_rounded_rect(slide, x, y, cell_w, cell_h, colors[i])
        # Big letter top-left
        add_text_box(
            slide,
            x + 0.5,
            y + 0.3,
            3.0,
            3.0,
            labels[i],
            size_pt=72,
            bold=True,
            color_hex="#FFFFFF",
            align="left",
            vcenter=False,
        )
        # Name
        add_text_box(
            slide,
            x + 3.5,
            y + 0.8,
            cell_w - 4.5,
            1.2,
            names[i].upper(),
            size_pt=18,
            bold=True,
            color_hex="#FFFFFF",
            align="left",
            vcenter=True,
        )
        # Glyph corner
        add_text_box(
            slide,
            x + cell_w - 1.5,
            y + 0.5,
            1.2,
            1.2,
            glyphs[i],
            size_pt=24,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Bullets — jh 2026-06-12 견본 구조: 첫 항목은 16B 헤드라인, 나머지는 11pt 세부
        import re as _re

        # 라벨 접두 제거 — "S ·", "S(강점) ·", "W · 약점(W) ·" 류 중복 방지 (2회 적용)
        _lbl_pat = (
            r"^[-\s]*(?:[SWOT]|강점|약점|기회|위협)\s*"
            r"(?:\((?:[SWOT]|강점|약점|기회|위협)\))?\s*[·:]\s*"
        )
        pts = [
            _re.sub(_lbl_pat, "", _re.sub(_lbl_pat, "", str(p))).strip()
            for p in (items[i].get("points", []) if i < len(items) else [])
        ][:3]
        # 페어 bullet 1개뿐이면 " — " 에서 분할해 헤드라인+세부 2층 구성
        if len(pts) == 1 and " — " in pts[0]:
            _head, _det = pts[0].split(" — ", 1)
            pts = [_head.strip(), _det.strip()]
        if pts:
            add_text_box(
                slide,
                x + 0.5,
                y + 3.2,
                cell_w - 1.0,
                1.1,
                f"{labels[i]} · {str(pts[0])[:46]}",
                size_pt=16,
                bold=True,
                color_hex="#FFFFFF",
                align="left",
                vcenter=False,
            )
        detail = "\n".join(f"·  {str(p)[:70]}" for p in pts[1:]) if len(pts) > 1 else ""
        if not pts:
            detail = "(추가 분석 필요)"
        if detail:
            add_text_box(
                slide,
                x + 0.5,
                y + 4.4,
                cell_w - 1.0,
                cell_h - 4.7,
                detail,
                size_pt=11,
                color_hex="#FFFFFF",
                align="left",
                vcenter=False,
            )


def draw_funnel_chart(slide, stages, primary, accent, ink, muted, light_bg):
    """Funnel from wide top to narrow bottom (trapezoid layers)."""
    n = min(len(stages), 5)
    if n == 0:
        return
    cx = SLIDE_W * 0.4
    margin_y = 5.0
    funnel_h = SLIDE_H - margin_y - 2.5
    layer_h = funnel_h / n
    top_w = 14.0
    bot_w = 4.0
    for i in range(n):
        item = stages[i]
        w_top = top_w - (top_w - bot_w) * (i / n)
        w_bot = top_w - (top_w - bot_w) * ((i + 1) / n)
        y = margin_y + i * layer_h
        # Use rectangle approximation (trapezoid via shape would need OXML)
        w_avg = (w_top + w_bot) / 2
        c = primary if i % 2 == 0 else accent
        text_c = "#FFFFFF" if i % 2 == 0 else primary
        add_rect(slide, cx - w_avg / 2, y, w_avg, layer_h - 0.15, c)
        # Layer value/name
        add_text_box(
            slide,
            cx - w_avg / 2,
            y,
            w_avg,
            layer_h - 0.15,
            f"{item.get('value', '')}  ·  {item.get('label', '')[:24]}",
            size_pt=14,
            bold=True,
            color_hex=text_c,
            align="center",
            vcenter=True,
        )
        # Right side caption
        add_text_box(
            slide,
            cx + 8.0,
            y + 0.3,
            SLIDE_W - cx - 9.0,
            layer_h - 0.6,
            str(item.get("caption", ""))[:120],
            size_pt=10,
            color_hex=ink,
            align="left",
            vcenter=True,
        )


def draw_two_col_stats(slide, items, primary, accent, ink, muted, light_bg):
    """Left big number + right description rows (4 rows)."""
    n = min(len(items), 4)
    if n == 0:
        return
    margin_y = 5.0
    row_h = (SLIDE_H - margin_y - 2.5) / n - 0.3
    for i in range(n):
        item = items[i]
        y = margin_y + i * (row_h + 0.3)
        add_rect(slide, 1.5, y, SLIDE_W - 3.0, row_h, "#FFFFFF")
        add_rect(slide, 1.5, y, 0.15, row_h, primary)
        # Big number left
        add_text_box(
            slide,
            2.0,
            y,
            6.0,
            row_h,
            str(item.get("value", "0")),
            size_pt=44,
            bold=True,
            color_hex=primary,
            align="left",
            vcenter=True,
        )
        # Right text
        add_text_box(
            slide,
            8.5,
            y + 0.3,
            SLIDE_W - 10.5,
            1.5,
            str(item.get("label", "")).upper()[:40],
            size_pt=11,
            bold=True,
            color_hex=muted,
            align="left",
            vcenter=True,
        )
        add_text_box(
            slide,
            8.5,
            y + 1.7,
            SLIDE_W - 10.5,
            row_h - 2.0,
            str(item.get("caption", ""))[:160],
            size_pt=11,
            color_hex=ink,
            align="left",
            vcenter=False,
        )


def draw_price_compare_3(slide, items, primary, accent, ink, muted, light_bg):
    """3 pricing/option compare cards (middle highlighted)."""
    n = min(len(items), 3)
    if n == 0:
        return
    margin_x = 3.0
    gap = 0.8
    avail_w = SLIDE_W - 2 * margin_x
    card_w = (avail_w - 2 * gap) / 3
    y = 5.5
    for i in range(n):
        item = items[i]
        x = margin_x + i * (card_w + gap)
        is_highlight = i == 1
        h = SLIDE_H - y - 2.5
        if is_highlight:
            h += 1.5
            y_card = y - 0.5
        else:
            y_card = y
        bg = primary if is_highlight else "#FFFFFF"
        line = None if is_highlight else "#E2E8F0"
        add_rounded_rect(slide, x, y_card, card_w, h, bg, line_hex=line)
        # Top label
        add_text_box(
            slide,
            x,
            y_card + 0.5,
            card_w,
            1.0,
            str(item.get("tier", f"Plan {i + 1}")).upper(),
            size_pt=12,
            bold=True,
            color_hex="#FFFFFF" if is_highlight else primary,
            align="center",
            vcenter=True,
        )
        # Big value
        add_text_box(
            slide,
            x,
            y_card + 2.0,
            card_w,
            3.0,
            str(item.get("value", "—")),
            size_pt=42,
            bold=True,
            color_hex="#FFFFFF" if is_highlight else primary,
            align="center",
            vcenter=True,
        )
        # Accent line
        add_rect(slide, x + card_w / 2 - 1.5, y_card + 5.2, 3.0, 0.08, accent if is_highlight else primary)
        # Caption lines
        for j, ln in enumerate(item.get("lines", [])[:4]):
            add_text_box(
                slide,
                x + 0.5,
                y_card + 5.8 + j * 1.0,
                card_w - 1.0,
                0.9,
                f"- {ln[:44]}",
                size_pt=12,
                color_hex="#E2E8F0" if is_highlight else ink,
                align="left",
                vcenter=True,
            )


def draw_team_cards_4(slide, members, primary, accent, ink, muted, light_bg):
    """4 team member cards: photo placeholder circle + name + role."""
    n = min(len(members), 4)
    if n == 0:
        return
    margin_x = 2.0
    avail_w = SLIDE_W - 2 * margin_x
    gap = 0.5
    card_w = (avail_w - 3 * gap) / 4
    y = 5.5
    card_h = SLIDE_H - y - 2.5
    glyphs = [GLYPHS["data"], GLYPHS["model"], GLYPHS["chart"], GLYPHS["bulb"]]
    for i in range(n):
        m = members[i]
        x = margin_x + i * (card_w + gap)
        add_rounded_rect(slide, x, y, card_w, card_h, "#FFFFFF", line_hex="#E2E8F0")
        # Photo placeholder (circle with glyph)
        psize = min(card_w - 1.5, 5.0)
        add_oval(slide, x + card_w / 2 - psize / 2 - 0.2, y + 0.8 - 0.2, psize + 0.4, psize + 0.4, accent)
        add_oval(slide, x + card_w / 2 - psize / 2, y + 0.8, psize, psize, primary)
        add_text_box(
            slide,
            x + card_w / 2 - psize / 2,
            y + 0.8,
            psize,
            psize,
            glyphs[i % 4],
            size_pt=42,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Name
        add_text_box(
            slide,
            x + 0.3,
            y + psize + 1.5,
            card_w - 0.6,
            1.2,
            str(m.get("name", f"Member {i + 1}"))[:24],
            size_pt=15,
            bold=True,
            color_hex=ink,
            align="center",
            vcenter=True,
        )
        # Accent line
        add_rect(slide, x + card_w / 2 - 1.0, y + psize + 2.7, 2.0, 0.05, accent)
        # Role
        add_text_box(
            slide,
            x + 0.3,
            y + psize + 3.0,
            card_w - 0.6,
            1.0,
            str(m.get("role", "")).upper()[:30],
            size_pt=10,
            bold=True,
            color_hex=primary,
            align="center",
            vcenter=True,
        )
        # Bio
        add_text_box(
            slide,
            x + 0.5,
            y + psize + 4.0,
            card_w - 1.0,
            card_h - psize - 4.5,
            str(m.get("bio", ""))[:100],
            size_pt=10,
            color_hex=muted,
            align="center",
            vcenter=False,
        )


def draw_split_compare(slide, left, right, primary, accent, ink, muted, light_bg):
    """Two-column split layout - left option vs right option side by side."""
    half = SLIDE_W / 2
    # Left side
    add_gradient_rect(slide, 0, 4.5, half - 0.05, SLIDE_H - 6.5, _shade(primary, 1.2), primary, angle=90)
    add_text_box(
        slide,
        1.0,
        5.5,
        half - 1.5,
        2.0,
        str(left.get("label", "OPTION A")).upper(),
        size_pt=24,
        bold=True,
        color_hex="#FFFFFF",
        align="left",
        vcenter=False,
    )
    add_rect(slide, 1.0, 7.5, 4.0, 0.1, accent)
    add_text_box(
        slide,
        1.0,
        8.0,
        half - 1.5,
        3.0,
        str(left.get("headline", ""))[:80],
        size_pt=16,
        color_hex="#E2E8F0",
        align="left",
        vcenter=False,
    )
    for i, pt in enumerate(left.get("points", [])[:4]):
        y = 11.5 + i * 1.3
        add_oval(slide, 1.0, y, 0.7, 0.7, accent)
        add_text_box(
            slide,
            1.0,
            y,
            0.7,
            0.7,
            GLYPHS["ok"],
            size_pt=10,
            bold=True,
            color_hex=primary,
            align="center",
            vcenter=True,
        )
        add_text_box(
            slide, 2.0, y, half - 2.5, 1.0, str(pt)[:70], size_pt=14, color_hex="#FFFFFF", align="left", vcenter=True
        )
    # Center divider with VS
    add_oval(slide, half - 1.2, SLIDE_H / 2 - 0.5, 2.4, 2.4, "#FFFFFF")
    add_oval(slide, half - 1.0, SLIDE_H / 2 - 0.3, 2.0, 2.0, accent)
    add_text_box(
        slide,
        half - 1.0,
        SLIDE_H / 2 - 0.3,
        2.0,
        2.0,
        "VS",
        size_pt=18,
        bold=True,
        color_hex=primary,
        align="center",
        vcenter=True,
    )
    # Right side (light bg)
    add_rect(slide, half + 0.05, 4.5, half - 0.05, SLIDE_H - 6.5, light_bg)
    add_text_box(
        slide,
        half + 1.0,
        5.5,
        half - 2.0,
        2.0,
        str(right.get("label", "OPTION B")).upper(),
        size_pt=24,
        bold=True,
        color_hex=primary,
        align="left",
        vcenter=False,
    )
    add_rect(slide, half + 1.0, 7.5, 4.0, 0.1, accent)
    add_text_box(
        slide,
        half + 1.0,
        8.0,
        half - 2.0,
        3.0,
        str(right.get("headline", ""))[:80],
        size_pt=16,
        color_hex=ink,
        align="left",
        vcenter=False,
    )
    for i, pt in enumerate(right.get("points", [])[:4]):
        y = 11.5 + i * 1.3
        add_oval(slide, half + 1.0, y, 0.7, 0.7, primary)
        add_text_box(
            slide,
            half + 1.0,
            y,
            0.7,
            0.7,
            GLYPHS["ok"],
            size_pt=10,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        add_text_box(
            slide, half + 2.0, y, half - 2.5, 1.0, str(pt)[:70], size_pt=14, color_hex=ink, align="left", vcenter=True
        )


def draw_big_index_chart(slide, big_index, chart_items, primary, accent, ink, muted, light_bg):
    """Left huge index number + right bar mini-chart."""
    # Left big number panel
    add_gradient_rect(slide, 0, 4.5, 10.0, SLIDE_H - 6.5, primary, _shade(primary, 0.7), angle=180)
    add_text_box(
        slide,
        1.0,
        6.0,
        8.0,
        6.0,
        str(big_index.get("number", "01")),
        size_pt=180,
        bold=True,
        color_hex="#FFFFFF",
        align="left",
        vcenter=True,
    )
    add_rect(slide, 1.0, 12.0, 5.0, 0.12, accent)
    add_text_box(
        slide,
        1.0,
        12.5,
        8.0,
        2.0,
        str(big_index.get("title", "INDEX")).upper()[:30],
        size_pt=16,
        bold=True,
        color_hex="#FFFFFF",
        align="left",
        vcenter=False,
    )
    add_text_box(
        slide,
        1.0,
        14.0,
        8.0,
        3.0,
        str(big_index.get("caption", ""))[:120],
        size_pt=11,
        color_hex="#E2E8F0",
        align="left",
        vcenter=False,
    )
    # Right area - mini bar items
    bx = 11.5
    bw = SLIDE_W - bx - 2.0
    for i, it in enumerate(chart_items[:5]):
        y = 5.5 + i * 2.2
        label = str(it.get("label", f"M{i + 1}"))[:24]
        val_pct = float(it.get("pct", 50))
        add_text_box(slide, bx, y, bw, 0.8, label, size_pt=11, bold=True, color_hex=ink, align="left", vcenter=True)
        add_text_box(
            slide,
            bx + bw - 3.0,
            y,
            3.0,
            0.8,
            f"{val_pct:.0f}%",
            size_pt=11,
            bold=True,
            color_hex=primary,
            align="right",
            vcenter=True,
        )
        # Bar bg
        add_rect(slide, bx, y + 0.9, bw, 0.4, "#E2E8F0")
        # Bar fill
        add_rect(slide, bx, y + 0.9, bw * val_pct / 100, 0.4, primary)


def draw_process_4_chevron(slide, items, primary, accent, ink, muted, light_bg):
    """4 chevron arrows pointing right with numbered headers."""
    n = min(len(items), 4)
    if n == 0:
        return
    margin_x = 2.0
    avail_w = SLIDE_W - 2 * margin_x
    chev_w = avail_w / n - 0.4
    chev_h = 4.5
    y = 7.0
    for i in range(n):
        item = items[i]
        x = margin_x + i * (chev_w + 0.4)
        c = primary if i % 2 == 0 else accent
        text_c = "#FFFFFF" if i % 2 == 0 else primary
        # Chevron
        add_chevron(slide, x, y, chev_w, chev_h, c)
        # Number circle at top-left of chevron
        add_oval(slide, x + 0.3, y - 1.5, 2.0, 2.0, primary)
        add_text_box(
            slide,
            x + 0.3,
            y - 1.5,
            2.0,
            2.0,
            f"{i + 1:02d}",
            size_pt=18,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Title inside chevron
        add_text_box(
            slide,
            x + 0.5,
            y + 0.5,
            chev_w - 1.5,
            1.5,
            str(item.get("title", ""))[:24],
            size_pt=14,
            bold=True,
            color_hex=text_c,
            align="center",
            vcenter=True,
        )
        # Caption inside
        add_text_box(
            slide,
            x + 0.5,
            y + 2.2,
            chev_w - 1.5,
            2.0,
            str(item.get("caption", ""))[:80],
            size_pt=10,
            color_hex=text_c,
            align="center",
            vcenter=False,
        )


def draw_dashboard_chart(slide, ctx, primary, accent, ink, muted, light_bg, render_fn):
    """Big chart left + 3 mini-stats stacked right."""
    # Big chart left 60%
    from outputs.architect.plan import VisualSpec

    if render_fn:
        try:
            vs = VisualSpec(
                type="chart_annotated_bar",
                title="핵심 지표 추이",
                spec={
                    "items": [
                        (k, m.get("value", 0))
                        for k, m in list(ctx.evaluation.metrics.items())[:5]
                        if isinstance(m.get("value"), (int, float))
                    ]
                },
            )
            png = render_fn(vs, ctx)
            if png:
                from pptx.util import Cm

                slide.shapes.add_picture(png, Cm(1.5), Cm(5.0), width=Cm(18.5), height=Cm(SLIDE_H - 7.0))
        except Exception:
            pass
    # Right - 3 stat cards stacked
    rx = 21.0
    rw = SLIDE_W - rx - 1.5
    card_h = (SLIDE_H - 7.0) / 3 - 0.3
    stats = []
    for k, m in list(ctx.evaluation.metrics.items())[:3]:
        v = m.get("value")
        stats.append((k, f"{v:.3f}" if isinstance(v, float) and v < 1 else str(v)))
    while len(stats) < 3:
        stats.append(("metric", "—"))
    for i, (lab, val) in enumerate(stats[:3]):
        y = 5.0 + i * (card_h + 0.3)
        add_rounded_rect(
            slide, rx, y, rw, card_h, primary if i == 0 else light_bg, line_hex=None if i == 0 else primary
        )
        add_text_box(
            slide,
            rx + 0.3,
            y + 0.3,
            rw - 0.6,
            1.0,
            lab.upper()[:24],
            size_pt=10,
            bold=True,
            color_hex="#E2E8F0" if i == 0 else muted,
            align="left",
            vcenter=True,
        )
        add_text_box(
            slide,
            rx + 0.3,
            y + 1.3,
            rw - 0.6,
            card_h - 1.5,
            val,
            size_pt=32,
            bold=True,
            color_hex="#FFFFFF" if i == 0 else primary,
            align="left",
            vcenter=True,
        )


def draw_circular_progress_4(slide, items, primary, accent, ink, muted, light_bg):
    """4 circular progress rings (with percent inside)."""
    n = min(len(items), 4)
    if n == 0:
        return
    margin_x = 2.0
    avail_w = SLIDE_W - 2 * margin_x
    gap = 0.6
    card_w = (avail_w - (n - 1) * gap) / n
    y = 5.5
    ring_size = min(card_w * 0.7, 7.0)
    for i in range(n):
        item = items[i]
        x = margin_x + i * (card_w + gap)
        cx = x + card_w / 2
        cy = y + ring_size / 2 + 0.5
        # Outer ring (full circle = bg)
        add_oval(slide, cx - ring_size / 2, cy - ring_size / 2, ring_size, ring_size, "#E2E8F0")
        # Filled inner ring (approximation using single inner circle with color = arc)
        # Since pptx doesn't have arc easily, use solid colored ring
        add_oval(slide, cx - ring_size / 2 + 0.5, cy - ring_size / 2 + 0.5, ring_size - 1.0, ring_size - 1.0, primary)
        # White inner
        inner = ring_size - 2.5
        add_oval(slide, cx - inner / 2, cy - inner / 2, inner, inner, "#FFFFFF")
        # Percent value
        add_text_box(
            slide,
            cx - inner / 2,
            cy - inner / 2,
            inner,
            inner * 0.6,
            str(item.get("value", "0%")),
            size_pt=28,
            bold=True,
            color_hex=primary,
            align="center",
            vcenter=True,
        )
        # Sub label
        add_text_box(
            slide,
            cx - inner / 2,
            cy,
            inner,
            inner * 0.4,
            str(item.get("sub", ""))[:12],
            size_pt=10,
            color_hex=muted,
            align="center",
            vcenter=True,
        )
        # Title below
        add_text_box(
            slide,
            x + 0.3,
            cy + ring_size / 2 + 0.5,
            card_w - 0.6,
            1.5,
            str(item.get("title", ""))[:30],
            size_pt=14,
            bold=True,
            color_hex=ink,
            align="center",
            vcenter=True,
        )
        # Caption
        add_text_box(
            slide,
            x + 0.3,
            cy + ring_size / 2 + 2.2,
            card_w - 0.6,
            SLIDE_H - cy - ring_size / 2 - 4.5,
            str(item.get("caption", ""))[:80],
            size_pt=9,
            color_hex=muted,
            align="center",
            vcenter=False,
        )


def draw_five_why(slide, why_chain, primary, accent, ink, muted, light_bg):
    """5 Why analysis - vertical cascade: 관찰 → Why1 → Why2 → ... → 근본원인."""
    n = min(len(why_chain), 5)
    if n == 0:
        return
    margin_x = 2.5
    avail_w = SLIDE_W - 2 * margin_x
    y0 = 5.0
    avail_h = SLIDE_H - y0 - 2.0
    step_h = avail_h / n - 0.3
    for i in range(n):
        item = why_chain[i]
        y = y0 + i * (step_h + 0.3)
        # Number badge on left
        add_oval(slide, margin_x, y + step_h / 2 - 1.0, 2.0, 2.0, primary)
        add_text_box(
            slide,
            margin_x,
            y + step_h / 2 - 1.0,
            2.0,
            2.0,
            f"WHY {i + 1}" if i > 0 else "관찰",
            size_pt=11,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Connecting arrow down (except last)
        if i < n - 1:
            add_text_box(
                slide,
                margin_x + 0.5,
                y + step_h + 0.05,
                1.0,
                0.4,
                GLYPHS["arrow_d"],
                size_pt=16,
                color_hex=accent,
                align="center",
                vcenter=True,
            )
        # Card on right
        card_x = margin_x + 2.5
        card_w = avail_w - 2.5
        # Last one (root cause) highlighted
        if i == n - 1:
            add_rounded_rect(slide, card_x, y, card_w, step_h, primary)
            text_c = "#FFFFFF"
            label_c = accent
            label = "근본 원인 (ROOT CAUSE)"
        else:
            add_rounded_rect(slide, card_x, y, card_w, step_h, "#FFFFFF", line_hex=primary)
            text_c = ink
            label_c = primary
            label = item.get("question", f"왜 {i + 1}?")
        # Label
        add_text_box(
            slide,
            card_x + 0.5,
            y + 0.2,
            card_w - 1.0,
            0.8,
            label.upper()[:50],
            size_pt=10,
            bold=True,
            color_hex=label_c,
            align="left",
            vcenter=True,
        )
        # Content
        add_text_box(
            slide,
            card_x + 0.5,
            y + 1.0,
            card_w - 1.0,
            step_h - 1.3,
            str(item.get("answer", ""))[:160],
            size_pt=13,
            bold=(i == n - 1),
            color_hex=text_c,
            align="left",
            vcenter=True,
        )


def draw_hypothesis_evidence_insight(slide, items, primary, accent, ink, muted, light_bg):
    """3-column flow: 가설 (Hypothesis) → 증거 (Evidence) → 인사이트 (Insight).

    items: list of {hypothesis, evidence, insight}.
    """
    n = min(len(items), 4)
    if n == 0:
        return
    # 3 column headers
    margin_x = 1.5
    avail_w = SLIDE_W - 2 * margin_x
    col_w = (avail_w - 2 * 0.6) / 3

    headers = [
        ("HYPOTHESIS", "가설", GLYPHS["bulb"]),
        ("EVIDENCE", "증거", GLYPHS["chart"]),
        ("INSIGHT", "인사이트", GLYPHS["star"]),
    ]
    h_y = 4.8
    for i, (lab_en, lab_ko, gly) in enumerate(headers):
        cx = margin_x + i * (col_w + 0.6)
        add_rect(slide, cx, h_y, col_w, 1.5, primary)
        add_text_box(slide, cx + 0.5, h_y, 1.5, 1.5, gly, size_pt=22, color_hex="#FFFFFF", align="center", vcenter=True)
        add_text_box(
            slide,
            cx + 2.0,
            h_y,
            col_w - 2.5,
            1.5,
            f"{lab_en}\n{lab_ko}",
            size_pt=15,
            bold=True,
            color_hex="#FFFFFF",
            align="left",
            vcenter=True,
        )

    # Rows
    y0 = h_y + 1.8
    avail_h = SLIDE_H - y0 - 2.0
    row_h = avail_h / n - 0.3
    for r in range(n):
        item = items[r]
        y = y0 + r * (row_h + 0.3)
        # Row number circle (left margin)
        add_text_box(
            slide,
            0.3,
            y + row_h / 2 - 0.4,
            margin_x - 0.5,
            0.8,
            f"{r + 1:02d}",
            size_pt=22,
            bold=True,
            color_hex=primary,
            align="right",
            vcenter=True,
        )
        # 3 cells
        cells = [
            (item.get("hypothesis", ""), "#FFFFFF", primary),
            (item.get("evidence", ""), light_bg, ink),
            (item.get("insight", ""), primary, "#FFFFFF"),
        ]
        for i, (text, bg, tc) in enumerate(cells):
            cx = margin_x + i * (col_w + 0.6)
            add_rounded_rect(slide, cx, y, col_w, row_h, bg, line_hex=primary if bg == "#FFFFFF" else None)
            add_text_box(
                slide,
                cx + 0.4,
                y + 0.3,
                col_w - 0.8,
                row_h - 0.6,
                str(text)[:200],
                size_pt=16,
                color_hex=tc,
                align="left",
                vcenter=True,
            )
        # Arrows between cells
        for i in range(2):
            cx = margin_x + (i + 1) * col_w + i * 0.6 + 0.05
            add_text_box(
                slide,
                cx,
                y + row_h / 2 - 0.3,
                0.5,
                0.6,
                GLYPHS["arrow_r"],
                size_pt=18,
                bold=True,
                color_hex=accent,
                align="center",
                vcenter=True,
            )


def draw_insight_funnel_4(slide, items, primary, accent, ink, muted, light_bg):
    """4-stage insight derivation funnel: 데이터 → 패턴 → 인사이트 → 액션."""
    n = min(len(items), 4)
    if n == 0:
        return
    labels = [
        ("DATA", "데이터", GLYPHS["data"]),
        ("PATTERN", "패턴", GLYPHS["chart"]),
        ("INSIGHT", "인사이트", GLYPHS["bulb"]),
        ("ACTION", "액션", GLYPHS["target"]),
    ]
    margin_x = 1.5
    avail_w = SLIDE_W - 2 * margin_x
    card_w = (avail_w - 3 * 0.5) / 4
    y = 5.5
    card_h = SLIDE_H - y - 2.5
    for i in range(n):
        item = items[i] if i < len(items) else {}
        lab_en, lab_ko, gly = labels[i]
        x = margin_x + i * (card_w + 0.5)
        # Gradient deepens through stages
        bg_start = _shade(primary, 1.5 - i * 0.15)
        bg_end = _shade(primary, 1.3 - i * 0.18)
        add_gradient_rect(slide, x, y, card_w, card_h, bg_start, bg_end, angle=180)
        # Icon at top
        add_text_box(
            slide,
            x,
            y + 0.5,
            card_w,
            2.0,
            gly,
            size_pt=36,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # English label
        add_text_box(
            slide,
            x,
            y + 2.5,
            card_w,
            0.8,
            lab_en,
            size_pt=10,
            bold=True,
            color_hex=accent,
            align="center",
            vcenter=True,
        )
        # Korean label
        add_text_box(
            slide,
            x,
            y + 3.3,
            card_w,
            1.5,
            lab_ko,
            size_pt=20,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Accent line
        add_rect(slide, x + card_w / 2 - 1.0, y + 5.0, 2.0, 0.08, accent)
        # Body
        add_text_box(
            slide,
            x + 0.4,
            y + 5.5,
            card_w - 0.8,
            card_h - 6.0,
            str(item.get("text", ""))[:160],
            size_pt=11,
            color_hex="#FFFFFF",
            align="left",
            vcenter=False,
        )
        # Arrow to next
        if i < n - 1:
            ax = x + card_w + 0.05
            add_text_box(
                slide,
                ax,
                y + card_h / 2 - 0.5,
                0.4,
                1.0,
                GLYPHS["arrow_r"],
                size_pt=18,
                bold=True,
                color_hex=primary,
                align="center",
                vcenter=True,
            )


def draw_horizontal_progress(slide, items, primary, accent, ink, muted, light_bg):
    """4-5 horizontal progress bars with percentages on right."""
    n = min(len(items), 5)
    if n == 0:
        return
    margin_x = 2.0
    y_start = 5.5
    avail_h = SLIDE_H - y_start - 2.5
    row_h = avail_h / n - 0.4
    bar_w = SLIDE_W - 2 * margin_x - 6.0
    for i in range(n):
        item = items[i]
        y = y_start + i * (row_h + 0.4)
        # Label
        add_text_box(
            slide,
            margin_x,
            y,
            8.0,
            row_h * 0.4,
            str(item.get("label", f"Item {i + 1}"))[:40],
            size_pt=14,
            bold=True,
            color_hex=ink,
            align="left",
            vcenter=False,
        )
        # Caption
        add_text_box(
            slide,
            margin_x,
            y + row_h * 0.4,
            8.0,
            row_h * 0.3,
            str(item.get("caption", ""))[:60],
            size_pt=10,
            color_hex=muted,
            align="left",
            vcenter=False,
        )
        # Progress track
        track_y = y + row_h * 0.75
        add_rect(slide, margin_x, track_y, bar_w, 0.6, "#E2E8F0")
        # Fill
        pct = float(item.get("pct", 50))
        fill_w = bar_w * pct / 100.0
        add_rect(slide, margin_x, track_y, fill_w, 0.6, primary)
        # Percentage on right
        add_text_box(
            slide,
            margin_x + bar_w + 0.5,
            y,
            5.0,
            row_h,
            f"{pct:.0f}%",
            size_pt=24,
            bold=True,
            color_hex=primary,
            align="left",
            vcenter=True,
        )


def draw_vertical_timeline(slide, milestones, primary, accent, ink, muted, light_bg):
    """Vertical timeline with side cards (alternating L/R)."""
    n = min(len(milestones), 5)
    if n == 0:
        return
    cx = SLIDE_W / 2
    y_start = 5.5
    y_end = SLIDE_H - 2.5
    avail = y_end - y_start
    spacing = avail / max(1, n - 1) if n > 1 else avail
    # Center vertical line
    add_rect(slide, cx - 0.06, y_start, 0.12, avail, primary)
    for i in range(n):
        m = milestones[i]
        ny = y_start + i * spacing if n > 1 else y_start + avail / 2
        # Node circle
        add_oval(slide, cx - 1.0, ny - 1.0, 2.0, 2.0, accent)
        add_oval(slide, cx - 0.7, ny - 0.7, 1.4, 1.4, primary)
        add_text_box(
            slide,
            cx - 0.7,
            ny - 0.7,
            1.4,
            1.4,
            f"{i + 1:02d}",
            size_pt=12,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Alternating cards
        if i % 2 == 0:
            cx_card = cx - 11.0
            align = "right"
        else:
            cx_card = cx + 2.0
            align = "left"
        add_rounded_rect(slide, cx_card, ny - 1.5, 9.0, 3.0, light_bg, line_hex=primary)
        add_text_box(
            slide,
            cx_card + 0.5,
            ny - 1.3,
            8.0,
            1.0,
            str(m.get("title", f"Phase {i + 1}"))[:30],
            size_pt=13,
            bold=True,
            color_hex=primary,
            align=align,
            vcenter=True,
        )
        add_text_box(
            slide,
            cx_card + 0.5,
            ny - 0.2,
            8.0,
            1.5,
            str(m.get("caption", ""))[:100],
            size_pt=10,
            color_hex=ink,
            align=align,
            vcenter=True,
        )


def draw_org_hierarchy(slide, root, children, primary, accent, ink, muted, light_bg):
    """3-tier org-style hierarchy: 1 root → up to 4 children."""
    cx = SLIDE_W / 2
    # Root box
    root_w, root_h = 7.0, 2.5
    root_y = 5.5
    add_rounded_rect(slide, cx - root_w / 2, root_y, root_w, root_h, primary)
    add_text_box(
        slide,
        cx - root_w / 2,
        root_y,
        root_w,
        root_h,
        str(root)[:30],
        size_pt=18,
        bold=True,
        color_hex="#FFFFFF",
        align="center",
        vcenter=True,
    )
    # Children row
    n = min(len(children), 4)
    if n == 0:
        return
    child_y = root_y + root_h + 3.0
    margin_x = 2.0
    avail = SLIDE_W - 2 * margin_x
    gap = 0.4
    child_w = (avail - (n - 1) * gap) / n
    child_h = 5.0
    for i in range(n):
        child = children[i]
        x = margin_x + i * (child_w + gap)
        # Connector line from root center to child center top
        cx_child = x + child_w / 2
        # Vertical drop from root bottom
        add_rect(slide, cx - 0.05, root_y + root_h, 0.1, 1.5, accent)
        # Horizontal bar
        if i == 0:
            add_rect(slide, min(cx_child, cx), root_y + root_h + 1.4, abs(cx_child - cx) + 0.1, 0.1, accent)
        else:
            add_rect(slide, min(cx_child, cx), root_y + root_h + 1.4, abs(cx_child - cx) + 0.1, 0.1, accent)
        # Vertical drop to child
        add_rect(slide, cx_child - 0.05, root_y + root_h + 1.5, 0.1, 1.5, accent)
        # Child card
        add_rounded_rect(slide, x, child_y, child_w, child_h, "#FFFFFF", line_hex=primary)
        # Number badge
        add_oval(slide, x + child_w / 2 - 0.8, child_y - 0.8, 1.6, 1.6, primary)
        add_text_box(
            slide,
            x + child_w / 2 - 0.8,
            child_y - 0.8,
            1.6,
            1.6,
            f"{i + 1:02d}",
            size_pt=11,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        add_text_box(
            slide,
            x + 0.3,
            child_y + 1.0,
            child_w - 0.6,
            1.2,
            str(child.get("title", ""))[:30],
            size_pt=13,
            bold=True,
            color_hex=primary,
            align="center",
            vcenter=True,
        )
        add_text_box(
            slide,
            x + 0.3,
            child_y + 2.3,
            child_w - 0.6,
            child_h - 2.6,
            str(child.get("caption", ""))[:100],
            size_pt=10,
            color_hex=ink,
            align="center",
            vcenter=False,
        )


def draw_stats_with_delta(slide, items, primary, accent, ink, muted, light_bg):
    """4 stats with ▲▼ delta indicators."""
    n = min(len(items), 4)
    if n == 0:
        return
    margin_x = 1.5
    avail_w = SLIDE_W - 2 * margin_x
    gap = 0.5
    card_w = (avail_w - (n - 1) * gap) / n
    y = 5.5
    card_h = SLIDE_H - y - 2.5
    for i in range(n):
        item = items[i]
        x = margin_x + i * (card_w + gap)
        add_rounded_rect(slide, x, y, card_w, card_h, "#FFFFFF", line_hex="#E2E8F0")
        # Top color bar
        add_rect(slide, x, y, card_w, 0.15, primary)
        # Label top
        add_text_box(
            slide,
            x + 0.4,
            y + 0.6,
            card_w - 0.8,
            0.8,
            str(item.get("label", "")).upper()[:30],
            size_pt=10,
            bold=True,
            color_hex=muted,
            align="left",
            vcenter=True,
        )
        # Big value
        add_text_box(
            slide,
            x + 0.4,
            y + 1.8,
            card_w - 0.8,
            3.0,
            str(item.get("value", "—")),
            size_pt=42,
            bold=True,
            color_hex=primary,
            align="left",
            vcenter=True,
        )
        # Delta row
        delta = float(item.get("delta", 0))
        is_up = delta >= 0
        delta_color = "#16A34A" if is_up else "#DC2626"
        glyph = GLYPHS["arrow_u"] if is_up else GLYPHS["arrow_d"]
        add_text_box(
            slide,
            x + 0.4,
            y + 5.5,
            1.0,
            1.0,
            glyph,
            size_pt=18,
            bold=True,
            color_hex=delta_color,
            align="left",
            vcenter=True,
        )
        add_text_box(
            slide,
            x + 1.4,
            y + 5.5,
            card_w - 1.8,
            1.0,
            f"{abs(delta):.1f}%  vs baseline",
            size_pt=11,
            bold=True,
            color_hex=delta_color,
            align="left",
            vcenter=True,
        )
        # Caption
        add_text_box(
            slide,
            x + 0.4,
            y + 7.0,
            card_w - 0.8,
            card_h - 7.5,
            str(item.get("caption", ""))[:100],
            size_pt=10,
            color_hex=ink,
            align="left",
            vcenter=False,
        )


def draw_bento_grid(slide, items, primary, accent, ink, muted, light_bg):
    """Mixed-size bento grid (1 big + 4 small)."""
    margin_x = 1.5
    margin_y = 5.0
    avail_w = SLIDE_W - 2 * margin_x
    avail_h = SLIDE_H - margin_y - 2.0
    big_w = avail_w * 0.55 - 0.3
    big_h = avail_h
    small_w = avail_w * 0.45 - 0.3
    small_h = (avail_h - 0.3) / 2

    # Big card left (item 0)
    if items:
        big = items[0]
        x = margin_x
        y = margin_y
        add_gradient_rect(slide, x, y, big_w, big_h, primary, _shade(primary, 0.7), angle=135)
        # Decorative circle
        add_oval(slide, x + big_w - 3.0, y + big_h - 3.0, 6.0, 6.0, accent)
        # Icon
        add_text_box(
            slide,
            x + 1.0,
            y + 1.0,
            3.0,
            3.0,
            big.get("icon", GLYPHS["target"]),
            size_pt=42,
            bold=True,
            color_hex="#FFFFFF",
            align="left",
            vcenter=True,
        )
        # Title
        add_text_box(
            slide,
            x + 1.0,
            y + 4.5,
            big_w - 2.0,
            3.0,
            str(big.get("title", "Main"))[:40],
            size_pt=28,
            bold=True,
            color_hex="#FFFFFF",
            align="left",
            vcenter=False,
        )
        # Accent line
        add_rect(slide, x + 1.0, y + 8.0, 4.0, 0.15, accent)
        # Body
        add_text_box(
            slide,
            x + 1.0,
            y + 8.5,
            big_w - 2.0,
            big_h - 9.5,
            str(big.get("caption", ""))[:200],
            size_pt=12,
            color_hex="#FFFFFF",
            align="left",
            vcenter=False,
        )

    # 4 small cards right
    small_x = margin_x + big_w + 0.3
    glyphs = [GLYPHS["chart"], GLYPHS["bulb"], GLYPHS["star"], GLYPHS["data"]]
    for i, item in enumerate(items[1:5]):
        r, c = i // 2, i % 2
        x = small_x + c * (small_w / 2 + 0.15)
        y = margin_y + r * (small_h + 0.3)
        sw = small_w / 2 - 0.15
        bg = "#FFFFFF" if i % 2 == 0 else light_bg
        add_rounded_rect(slide, x, y, sw, small_h, bg, line_hex=primary)
        add_text_box(
            slide,
            x + 0.3,
            y + 0.3,
            1.5,
            1.5,
            glyphs[i % 4],
            size_pt=22,
            bold=True,
            color_hex=primary,
            align="left",
            vcenter=True,
        )
        add_text_box(
            slide,
            x + 0.3,
            y + 2.0,
            sw - 0.6,
            1.5,
            str(item.get("title", f"Item {i + 2}"))[:18],
            size_pt=14,
            bold=True,
            color_hex=ink,
            align="left",
            vcenter=True,
        )
        add_text_box(
            slide,
            x + 0.3,
            y + 3.5,
            sw - 0.6,
            small_h - 4.0,
            str(item.get("caption", ""))[:80],
            size_pt=9,
            color_hex=muted,
            align="left",
            vcenter=False,
        )


def draw_roadmap_with_upgrades(slide, phases, upgrades, primary, accent, ink, muted, light_bg):
    """Top: 3-phase horizontal chevron timeline.
    Bottom: 3 upgrade cards grid (future development).
    """
    # ===== Top section: Phase 1 -> 2 -> 3 horizontal flow =====
    add_text_box(
        slide,
        1.5,
        4.7,
        14.0,
        0.8,
        "단계별 실행 · ROLLOUT PHASES",
        size_pt=11,
        bold=True,
        color_hex=primary,
        align="left",
        vcenter=True,
    )
    add_rect(slide, 1.5, 5.5, 3.0, 0.08, accent)

    phase_y = 6.0
    phase_h = 3.5
    margin_x = 1.5
    avail_w = SLIDE_W - 2 * margin_x
    n_phases = min(len(phases), 3)
    chev_w = avail_w / n_phases - 0.4
    for i in range(n_phases):
        item = phases[i]
        x = margin_x + i * (chev_w + 0.4)
        c = primary if i % 2 == 0 else accent
        text_c = "#FFFFFF" if i % 2 == 0 else primary
        # Chevron arrow
        add_chevron(slide, x, phase_y, chev_w, phase_h, c)
        # Phase label
        add_text_box(
            slide,
            x + 0.5,
            phase_y + 0.3,
            chev_w - 1.5,
            0.9,
            f"PHASE {i + 1:02d}",
            size_pt=11,
            bold=True,
            color_hex=text_c,
            align="left",
            vcenter=True,
        )
        # Title
        add_text_box(
            slide,
            x + 0.5,
            phase_y + 1.2,
            chev_w - 1.5,
            1.0,
            str(item.get("title", "")).strip()[:30],
            size_pt=14,
            bold=True,
            color_hex=text_c,
            align="left",
            vcenter=True,
        )
        # Caption
        add_text_box(
            slide,
            x + 0.5,
            phase_y + 2.2,
            chev_w - 1.5,
            1.0,
            str(item.get("caption", ""))[:80],
            size_pt=10,
            color_hex=text_c,
            align="left",
            vcenter=False,
        )

    # ===== Divider =====
    div_y = phase_y + phase_h + 0.8
    add_rect(slide, 1.5, div_y, SLIDE_W - 3.0, 0.04, "#E2E8F0")

    # ===== Bottom section: 3 upgrade cards =====
    add_text_box(
        slide,
        1.5,
        div_y + 0.4,
        14.0,
        0.8,
        "향후 고도화 · ENHANCEMENT ROADMAP",
        size_pt=11,
        bold=True,
        color_hex=primary,
        align="left",
        vcenter=True,
    )
    add_rect(slide, 1.5, div_y + 1.2, 3.0, 0.08, accent)

    n_up = min(len(upgrades), 3)
    if n_up == 0:
        return
    up_card_w = (avail_w - (n_up - 1) * 0.5) / n_up
    up_y = div_y + 1.8
    up_h = SLIDE_H - up_y - 2.0
    glyphs = [GLYPHS["bulb"], GLYPHS["chart"], GLYPHS["target"]]
    for i in range(n_up):
        item = upgrades[i]
        x = margin_x + i * (up_card_w + 0.5)
        add_rounded_rect(slide, x, up_y, up_card_w, up_h, "#FFFFFF", line_hex=primary)
        # Top accent strip
        add_rect(slide, x, up_y, up_card_w, 0.15, accent)
        # Glyph + number
        add_oval(slide, x + 0.5, up_y + 0.5, 1.6, 1.6, primary)
        add_text_box(
            slide,
            x + 0.5,
            up_y + 0.5,
            1.6,
            1.6,
            glyphs[i % len(glyphs)],
            size_pt=18,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Upgrade number tag
        add_text_box(
            slide,
            x + 2.3,
            up_y + 0.5,
            up_card_w - 2.8,
            0.8,
            f"UPGRADE  {i + 1:02d}",
            size_pt=11,
            bold=True,
            color_hex=muted,
            align="left",
            vcenter=True,
        )
        # Title
        add_text_box(
            slide,
            x + 2.3,
            up_y + 1.3,
            up_card_w - 2.8,
            0.9,
            str(item.get("title", "")).strip()[:30],
            size_pt=15,
            bold=True,
            color_hex=primary,
            align="left",
            vcenter=True,
        )
        # Caption
        add_text_box(
            slide,
            x + 0.5,
            up_y + 2.5,
            up_card_w - 1.0,
            up_h - 3.0,
            str(item.get("caption", ""))[:120],
            size_pt=12,
            color_hex=ink,
            align="left",
            vcenter=False,
        )


def draw_exec_summary_v32(slide, findings, boxes, primary, accent, ink, muted, light_bg):
    """Executive Summary v32 — 목표치 덱 S3 실측 이식 (jh 2026-06-12).

    상단: 3 FINDING 카드 (라벨 10B / 빅수치 46B / 이름 14B+밑줄 / 세부 10pt 3줄)
    하단: METHOD·PERFORMANCE·LIMITATION 3박스 (라벨 10B / 핵심값 18B / 세부 11pt 3줄)
    PERFORMANCE 박스만 primary 솔리드 강조.

    Args:
        findings: [{label, big, feature(이름), sub(세부 줄들 str|list)}] ×3
        boxes: [(label, headline, [lines])] ×3 — index 1 이 강조 박스
    """
    margin_x = 1.5
    gap = 0.5
    card_w = (SLIDE_W - 2 * margin_x - 2 * gap) / 3

    # ── Row 1: FINDING 카드 3개 ──
    y1, h1 = 4.4, 6.8
    for i, f in enumerate((findings or [])[:3]):
        x = margin_x + i * (card_w + gap)
        add_rounded_rect(slide, x, y1, card_w, h1, "#FFFFFF", line_hex=primary)
        add_rect(slide, x + 0.3, y1, card_w - 0.6, 0.12, accent)
        add_text_box(
            slide, x + 0.6, y1 + 0.4, card_w - 1.2, 0.7,
            str(f.get("label", f"FINDING {i + 1:02d}")),
            size_pt=10, bold=True, color_hex=muted, align="left", vcenter=True,
        )
        add_text_box(
            slide, x + 0.6, y1 + 1.1, card_w - 1.2, 2.3,
            str(f.get("big", "-"))[:12],
            size_pt=46, bold=True, color_hex=primary, align="left", vcenter=True,
        )
        add_text_box(
            slide, x + 0.6, y1 + 3.5, card_w - 1.2, 0.9,
            str(f.get("feature", ""))[:20],
            size_pt=14, bold=True, color_hex=ink, align="left", vcenter=True,
        )
        add_rect(slide, x + 0.6, y1 + 4.4, 2.5, 0.06, accent)
        sub = f.get("sub", "")
        sub_lines = sub if isinstance(sub, list) else [s for s in str(sub).split("\n") if s]
        add_text_box(
            slide, x + 0.6, y1 + 4.7, card_w - 1.2, h1 - 5.0,
            "\n".join(str(s)[:40] for s in sub_lines[:3]),
            size_pt=10, color_hex=muted, align="left", vcenter=False,
        )

    # ── Row 2: METHOD / PERFORMANCE / LIMITATION ──
    y2, h2 = 11.7, 5.0
    for i, (label, headline, lines) in enumerate((boxes or [])[:3]):
        x = margin_x + i * (card_w + gap)
        solid = i == 1
        bg = primary if solid else "#FFFFFF"
        add_rounded_rect(slide, x, y2, card_w, h2, bg, line_hex=None if solid else primary)
        add_text_box(
            slide, x + 0.6, y2 + 0.35, card_w - 1.2, 0.7,
            str(label).upper(),
            size_pt=10, bold=True,
            color_hex="#BFDBFE" if solid else muted, align="left", vcenter=True,
        )
        add_text_box(
            slide, x + 0.6, y2 + 1.0, card_w - 1.2, 1.0,
            str(headline)[:24],
            size_pt=18, bold=True,
            color_hex="#FFFFFF" if solid else ink, align="left", vcenter=True,
        )
        add_text_box(
            slide, x + 0.6, y2 + 2.1, card_w - 1.2, h2 - 2.5,
            "\n".join(str(ln)[:42] for ln in (lines or [])[:3]),
            size_pt=11,
            color_hex="#E2E8F0" if solid else ink, align="left", vcenter=False,
        )


# ==============================================================
# 2026-06-11 jh — 수작업 견본 덱 레이아웃 이식 3종 (카탈로그 보강)
# 기하는 견본 (2026-06-10 수정본) 슬라이드에서 실측한 좌표.
# ==============================================================


def draw_solution_overview(slide, caption, steps, specs, primary, accent, ink, muted, light_bg):
    """견본 S8 — 좌: 파이프라인 3단계 카드 / 우: label-value 스펙 4쌍.

    Args:
        caption: 좌측 패널 상단 캡션 (예: "Solution Pipeline — 입력 → 모델 → 출력")
        steps:   [{"label": "INPUT", "text": "..."}] 최대 3개
        specs:   [{"label": "모델", "text": "..."}] 최대 4개
    """
    # 좌측 패널 (실측: 1.0,4.5 / 14.5×12.6)
    add_rounded_rect(slide, 1.0, 4.5, 14.5, 12.6, light_bg, line_hex=primary)
    add_text_box(slide, 1.4, 4.8, 13.8, 0.9, str(caption)[:60], size_pt=12, bold=True,
                 color_hex=primary, align="left", vcenter=True)

    n = min(len(steps), 3) or 1
    card_w, card_h, gap = 4.2, 4.6, 0.55
    total = n * card_w + (n - 1) * gap
    x0 = 1.0 + (14.5 - total) / 2
    for i, st in enumerate(steps[:3]):
        x = x0 + i * (card_w + gap)
        add_rounded_rect(slide, x, 8.7, card_w, card_h, "#FFFFFF", line_hex=accent)
        add_text_box(slide, x, 9.0, card_w, 0.8, str(st.get("label", ""))[:12].upper(),
                     size_pt=12, bold=True, color_hex=accent, align="center", vcenter=True)
        add_text_box(slide, x + 0.3, 10.0, card_w - 0.6, card_h - 1.6, str(st.get("text", ""))[:70],
                     size_pt=10, color_hex=ink, align="center", vcenter=False)
        if i < n - 1:
            add_text_box(slide, x + card_w, 10.3, gap, 1.2, "→", size_pt=16, bold=True,
                         color_hex=primary, align="center", vcenter=True)

    # 우측 스펙 (실측: 19.5, 5.1 시작 / label 1.0 + value 1.8 / step 3.0)
    y = 5.1
    for sp in specs[:4]:
        add_text_box(slide, 19.5, y, 13.0, 0.9, str(sp.get("label", ""))[:20], size_pt=13,
                     bold=True, color_hex=primary, align="left", vcenter=True)
        add_text_box(slide, 19.5, y + 1.0, 13.0, 1.9, str(sp.get("text", ""))[:110], size_pt=11,
                     color_hex=ink, align="left", vcenter=False)
        y += 3.0


def draw_lineage_2col(slide, bars, cards, primary, accent, ink, muted, light_bg):
    """견본 S9 — 좌: lineage 단계 바 4개 / 우: 보충 설명 카드 3개.

    Args:
        bars:  [{"label": "원본 데이터", "text": "..."}] 최대 4개
        cards: [{"label": "...", "text": "..."}] 최대 3개
    """
    # 좌측 바 (실측: 1.5, 4.5 + i*3.16 / 19.1×2.6)
    for i, b in enumerate(bars[:4]):
        y = 4.5 + i * 3.16
        add_rounded_rect(slide, 1.5, y, 19.1, 2.6, light_bg if i % 2 == 0 else "#FFFFFF",
                         line_hex=primary)
        add_text_box(slide, 1.9, y + 0.3, 1.6, 2.0, f"{i + 1:02d}", size_pt=22, bold=True,
                     color_hex=accent, align="left", vcenter=True)
        add_text_box(slide, 3.8, y + 0.25, 16.4, 0.9, str(b.get("label", ""))[:30], size_pt=12,
                     bold=True, color_hex=ink, align="left", vcenter=True)
        add_text_box(slide, 3.8, y + 1.2, 16.4, 1.2, str(b.get("text", ""))[:90], size_pt=10,
                     color_hex=muted, align="left", vcenter=False)

    # 우측 카드 (실측: 21.2, 4.5 + i*4.15 / 11.1×3.8)
    for i, c in enumerate(cards[:3]):
        y = 4.5 + i * 4.15
        add_rounded_rect(slide, 21.2, y, 11.1, 3.8, light_bg, line_hex=accent)
        add_text_box(slide, 21.6, y + 0.3, 10.3, 0.8, str(c.get("label", ""))[:24], size_pt=11,
                     bold=True, color_hex=primary, align="left", vcenter=True)
        add_text_box(slide, 21.6, y + 1.3, 10.3, 2.2, str(c.get("text", ""))[:110], size_pt=10,
                     color_hex=ink, align="left", vcenter=False)


def draw_icon_columns(slide, items, primary, accent, ink, muted, light_bg):
    """견본 S15 — N컬럼 (TAG + 라벨 + 글리프 + 본문), 2~4열 자동.

    Args:
        items: [{"tag": "DATA", "label": "데이터", "glyph": "data", "text": "..."}]
    """
    n = max(2, min(len(items), 4))
    gap = 0.5
    col_w = (SLIDE_W - 3.0 - (n - 1) * gap) / n
    for i, it in enumerate(items[:4]):
        x = 1.5 + i * (col_w + gap)
        add_rounded_rect(slide, x, 5.0, col_w, 12.1, light_bg, line_hex=primary)
        add_text_box(slide, x, 5.5, col_w, 0.9, str(it.get("tag", ""))[:14].upper(), size_pt=11,
                     bold=True, color_hex=accent, align="center", vcenter=True)
        add_text_box(slide, x, 6.5, col_w, 1.0, str(it.get("label", ""))[:16], size_pt=16,
                     bold=True, color_hex=ink, align="center", vcenter=True)
        glyph = GLYPHS.get(str(it.get("glyph", "")), GLYPHS.get("data", "")) or ""
        add_text_box(slide, x, 8.0, col_w, 1.6, glyph, size_pt=28, color_hex=primary,
                     align="center", vcenter=True)
        # jh 2026-06-12 — 본문 10→13pt + 세로 중앙 (인사이트 종합 글씨 작고 짧다 지적)
        add_text_box(slide, x + 0.4, 10.0, col_w - 0.8, 6.7, str(it.get("text", ""))[:220],
                     size_pt=13, color_hex=ink, align="left", vcenter=True)
