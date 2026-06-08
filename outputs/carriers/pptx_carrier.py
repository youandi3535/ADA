"""outputs.carriers.pptx_carrier - PPTX with KO font + visual PNG embed."""

from __future__ import annotations

from pathlib import Path

from outputs.architect.plan import ReportPlan
from outputs.context.schema import ReportContext

KO_FONT = "Malgun Gothic"


def generate_pptx(plan: ReportPlan, ctx: ReportContext, output_path) -> str:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Cm, Pt
    except Exception:
        return _fb(plan, ctx, output_path)

    from outputs.layouts.pptx_layouts import shows_footer, slide_box_specs
    from outputs.localization.korean import format_date_ko
    from outputs.style.classification import classification_treatment
    from outputs.style.palette import get_palette, hex_to_rgb
    from outputs.visuals.render import render_visual_to_png

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    pal = get_palette(ctx.meta.category)
    treat = classification_treatment(ctx.meta.classification)
    primary = RGBColor(*hex_to_rgb(pal["primary"]))
    ink = RGBColor(0x0F, 0x17, 0x2A)
    muted = RGBColor(0x64, 0x74, 0x8B)

    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)

    backup_ids = {s.id for sec in plan.sections if sec.id == "backup" for s in sec.slides}
    slides_flat = [sl for sec in plan.sections for sl in sec.slides]

    for idx, sl in enumerate(slides_flat):
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
        slide = prs.slides.add_slide(layout)

        if sl.layout in ("section_divider",):
            try:
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = primary
            except Exception:
                pass

        boxes = slide_box_specs(sl.layout)
        for box in boxes:
            _slot(slide, sl, box, ctx, primary, ink, Cm, Pt, RGBColor, render_visual_to_png)

        if shows_footer(sl.layout):
            _footer(slide, ctx, treat, idx + 1, len(slides_flat), Cm, Pt, muted, RGBColor, format_date_ko)

        if sl.speaker_notes_hint:
            try:
                slide.notes_slide.notes_text_frame.text = sl.speaker_notes_hint
            except Exception:
                pass

        if sl.id in backup_ids:
            try:
                slide.element.set("show", "0")
            except Exception:
                pass

    prs.save(str(out))
    return str(out)


def _slot(slide, sl, box, ctx, primary, ink, Cm, Pt, RGBColor, render_fn) -> None:
    """Render single slot — image for visual slots, text for the rest."""
    t = box["type"]
    x, y, w, h = Cm(box["x_cm"]), Cm(box["y_cm"]), Cm(box["w_cm"]), Cm(box["h_cm"])

    # Visual slots — try PNG embed first
    if t in ("chart", "diagram", "table", "kpi_card") and sl.visual_spec:
        try:
            png = render_fn(sl.visual_spec, ctx, slide=sl)
            if png:
                slide.shapes.add_picture(png, x, y, width=w, height=h)
                return
        except Exception:
            pass

    # Text slots
    try:
        tx = slide.shapes.add_textbox(x, y, w, h)
        tf = tx.text_frame
        tf.word_wrap = True

        if t == "so_what":
            tf.text = sl.so_what or ""
            _font(tf, 18, True, primary, Pt)
        elif t == "title":
            tf.text = sl.title_ko or sl.id
            _font(tf, 28, True, primary, Pt)
        elif t == "body":
            tf.text = ""
            for i, b in enumerate(sl.body_outline):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.text = "- " + b
            _font(tf, 14, False, ink, Pt)
        elif t == "cover_block":
            tf.text = sl.title_ko or "ADA 보고서"
            _font(tf, 44, True, primary, Pt)
            if sl.body_outline:
                p = tf.add_paragraph()
                p.text = "\n".join(sl.body_outline[:3])
                for r in p.runs:
                    r.font.size = Pt(14)
                    _set_face(r, KO_FONT)
        elif t == "closing_block":
            tf.text = sl.so_what or "Summary"
            _font(tf, 28, True, primary, Pt)
            for b in sl.body_outline[:3]:
                p = tf.add_paragraph()
                p.text = b
                for r in p.runs:
                    r.font.size = Pt(14)
                    _set_face(r, KO_FONT)
        elif t == "agenda_list":
            tf.text = ""
            for i, item in enumerate(sl.body_outline):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"{i + 1:>2}.  {item}"
            _font(tf, 16, False, ink, Pt)
        elif t == "quote_block":
            tf.text = '"' + (sl.so_what or "") + '"'
            _font(tf, 24, False, ink, Pt)
        elif t == "kpi_card":
            txt = (sl.body_outline[0] if sl.body_outline else sl.so_what) or ""
            tf.text = txt[:40]
            _font(tf, 20, True, primary, Pt)
    except Exception:
        pass


def _font(tf, size_pt, bold, color, Pt):
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(size_pt)
            run.font.bold = bold
            if color is not None:
                run.font.color.rgb = color
            _set_face(run, KO_FONT)


def _set_face(run, family):
    """Set both latin and east-asian font face for Korean."""
    try:
        run.font.name = family
        from pptx.oxml.ns import qn

        rPr = run._r.get_or_add_rPr()
        for ea in rPr.findall(qn("a:ea")):
            rPr.remove(ea)
        from lxml import etree

        ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", family)
    except Exception:
        pass


def _footer(slide, ctx, treat, page, total, Cm, Pt, muted, RGBColor, fmt_date):
    try:
        from outputs.style.palette import hex_to_rgb

        title_short = (ctx.meta.user_intent or "보고서")[:30]
        date_str = fmt_date(ctx.meta.generated_at, style="iso")

        left = slide.shapes.add_textbox(Cm(1.0), Cm(18.3), Cm(13.0), Cm(0.5))
        left.text_frame.text = f"ADA · {title_short}"
        _font(left.text_frame, 9, False, muted, Pt)

        mid = slide.shapes.add_textbox(Cm(14.5), Cm(18.3), Cm(7.0), Cm(0.5))
        mid.text_frame.text = f"{date_str}   {page}/{total}"
        _font(mid.text_frame, 9, False, muted, Pt)

        right = slide.shapes.add_textbox(Cm(22.5), Cm(18.3), Cm(10.0), Cm(0.5))
        right.text_frame.text = treat.get("footer_text", "INTERNAL")
        _font(right.text_frame, 9, True, RGBColor(*hex_to_rgb(treat.get("footer_color", "#334155"))), Pt)
    except Exception:
        pass


def _fb(plan, ctx, output_path):
    out = Path(str(output_path) + ".txt")
    out.parent.mkdir(parents=True, exist_ok=True)
    lines = [f"[PPT Fallback {plan.skeleton}]"]
    for sec in plan.sections:
        lines.append(f"# {sec.title}")
        for sl in sec.slides:
            lines.append(f"  - [{sl.layout}] {sl.title_ko}: {sl.so_what}")
    out.write_text("\n".join(lines), encoding="utf-8")
    return str(out)
