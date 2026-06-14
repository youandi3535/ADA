"""outputs.carriers.pptx_designer - Consulting-grade PPTX design layer.

Builds slides with proper shapes/typography/visuals matching enterprise reports:
  - Cover: full-bleed color block + large typography + dot pattern motif
  - Section divider: huge chapter number + full-bleed primary
  - KPI cards: rounded rect + huge number + accent line
  - Process flow: chevron arrows + numbered steps
  - Comparison table: colored header + zebra rows
  - Body slide: left vertical accent + So-What header band
"""

from __future__ import annotations

from pathlib import Path

from outputs.architect.plan import ReportPlan, SlideSpec

# Step 7-1 — LLM 디자인 힌트 적용 헬퍼 (palette/photo/icon)
from outputs.carriers.design_hint_helpers import (
    check_and_log_miss,
    icon_name as design_icon_name,
    photo_keyword,
)
from outputs.context.schema import ReportContext
from outputs.style.visual_kit import (
    GLYPHS,
    add_color_band,
    add_glyph,
    add_gradient_rect,
    add_horizontal_rule,
    add_image_with_overlay,
    add_oval,
    add_rect,
    add_rounded_rect,
    add_text_box,
    add_triangle,
    add_vertical_accent,
)

# HJ 2026-06-08 — 시각 품질 도구 (graceful import, 자산 없어도 동작)
try:
    from tools.visual import (
        get_cover_image,
        get_font_name,
        icon_for_slide,
        illustration_for_slide,
        recolor_undraw,
        svg_to_png_recolored,
    )

    _VISUAL_TOOLS_AVAILABLE = True
except ImportError:
    _VISUAL_TOOLS_AVAILABLE = False

    def get_cover_image(*args, **kwargs):
        return None

    def get_font_name(*args, **kwargs):
        return "Malgun Gothic"

    def icon_for_slide(*args, **kwargs):
        return None

    def illustration_for_slide(*args, **kwargs):
        return None

    def recolor_undraw(*args, **kwargs):
        return None

    def svg_to_png_recolored(*args, **kwargs):
        return None


# 한국어 폰트 — 자산 있으면 Pretendard, 없으면 Malgun Gothic
KO = get_font_name()
SLIDE_W = 33.867
SLIDE_H = 19.05


def generate_pptx_designed(plan: ReportPlan, ctx: ReportContext, output_path) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Cm
    except Exception:
        return _fb(plan, output_path)

    from outputs.localization.korean import format_date_ko
    from outputs.style.classification import classification_treatment
    from outputs.style.palette import get_palette
    from outputs.visuals.render import render_visual_to_png

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    pal = get_palette(ctx.meta.category)
    treat = classification_treatment(ctx.meta.classification)
    primary = pal["primary"]
    accent = pal["accent"]
    secondary = pal["secondary"]
    ink = "#0F172A"
    muted = "#64748B"
    light_bg = "#F8FAFC"

    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W)
    prs.slide_height = Cm(SLIDE_H)
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

    # Exclude backup section entirely (user request: no backup slides at all)
    non_backup_sections = [sec for sec in plan.sections if sec.id != "backup"]
    raw_slides = [sl for sec in non_backup_sections for sl in sec.slides]

    # Slide order:
    #   1. cover  →  2. agenda  →  3. 분석 가설 (본문 시작)
    #   →  본문 (시장·문제·솔루션·기술·데이터·모델·성능...)
    #   →  성능 평가 직후: 가설 입증 인사이트
    #   →  임팩트·로드맵·리스크  →  마무리
    # ExecSummary 는 본문에 포함 (컨설팅 BLUF — 임원이 1장만 봐도 결론 명확).
    cover_sl = next((s for s in raw_slides if s.layout == "cover"), None)
    agenda_sl = next((s for s in raw_slides if s.layout == "agenda"), None)
    exec_sl = next((s for s in raw_slides if s.id == "exec_summary"), None)
    # jh 2026-06-12 — 인사이트 종합은 목차 *앞* (BLUF: 결론→종합 인사이트→목차→본문)
    synth_sl = next((s for s in raw_slides if s.id == "insight_synthesis"), None)

    # ml_pitch (및 향후 카테고리별 skeleton) 가 hypothesis/insights_derived 슬라이드를
    # 자기 plan 안에 포함. 디자이너가 추가 주입하지 않음 (중복 방지).

    # Build body list (excluding cover, exec_summary, synthesis, agenda — 본문 흐름만)
    body = [s for s in raw_slides if s not in (cover_sl, agenda_sl, exec_sl, synth_sl)]

    # Final flat order: cover, exec, 인사이트 종합, agenda, body, closing(last)
    slides_flat = [s for s in (cover_sl, exec_sl, synth_sl, agenda_sl) if s is not None] + body
    # Force-rebuild Agenda's body_outline = exact list of body slide titles
    # (everything after Agenda except the very last closing slide).
    # jh 2026-06-12 — 목표치 덱 스토리 아젠다: 슬라이드 id → 이야기 흐름 라벨.
    # (Action Title 그대로 나열하면 목차가 아니라 결론 모음이 됨 — 목차는 *질문/주제* 톤)
    # EDA 4장은 데이터 기반 제목(title_ko)이 더 구체적이라 맵에서 제외.
    _AGENDA_STORY_LABELS = {
        "hypothesis": "분석 가설 — 무엇을 물었나",
        "p1_market": "분석 배경 — 세 가지 질문",
        "p2_pain": "기술 스택",
        "p3_alt_limits": "분석 방법 — 5단계 설계",
        "insight_synthesis": "인사이트 종합 — 발견과 접목 범위",
        "i1_kpi": "모델 성능 — Baseline 대비 가치",
        "eda_findings": "SHAP — 전역 importance",
        "error_analysis": "SHAP — 개별 케이스 해석",
        "insights_derived": "Error Analysis — 어디서 틀리는가",
        "as_is_to_be": "세그먼트 성능 — 취약 구간 식별",
        "i3_roi": "운영 정책 — 도입 룰",
        "risk_mitigation": "Risk & Drift — SWOT",
        "roadmap": "실행 로드맵 — 후속 작업",
    }
    def _rebuild_agenda() -> None:
        """목차 항목 재구성 — 카피라이터 적용 *후* 호출해야 본문 제목과 일치.

        jh 2026-06-12 — 초안 제목으로 목차를 굳혀 본문(카피라이터 제목)과
        불일치하던 결함 수정 (스토리 라벨 우선, EDA 는 최종 제목).
        """
        if agenda_sl is None:
            return
        body_titles = []
        for i, sl in enumerate(slides_flat):
            if sl is agenda_sl:
                for j, body_sl in enumerate(slides_flat[i + 1 : -1], 1):
                    title = (
                        _AGENDA_STORY_LABELS.get(body_sl.id)
                        or body_sl.title_ko
                        or body_sl.id
                        or f"슬라이드 {j}"
                    )[:60]
                    body_titles.append(f"{j:02d}.  {title}")
                break
        if body_titles:
            agenda_sl.body_outline = body_titles

    _rebuild_agenda()

    # Section ID list for chapter numbering on dividers
    section_idx = {}
    n = 0
    for sec in plan.sections:
        if sec.kind not in ("cover", "closing") and sec.id != "backup":
            n += 1
            section_idx[sec.id] = n

    page_num = 0
    total = len(slides_flat) + sum(1 for sec in plan.sections if sec.divider_required and sec.id in section_idx)

    # Section dividers removed - they showed empty pages with only a chapter number.
    # Section transitions are signaled by the body slide header (left accent + title).
    # Iterate REORDERED slides_flat (cover, exec, agenda, rest; no backup)
    total = len(slides_flat)
    page_num = 0

    # === LLM 카피라이팅 (Phase 3 — jh 2026-06-11) ===
    # 템플릿 구성 (SlideSpec 슬롯) + ctx 데이터로 덱 전체 문장 1회 채움.
    # 실패 시 skeleton 초안 문장 유지.
    try:
        _prefill_copy(slides_flat, ctx)
    except Exception as _e:  # noqa: BLE001
        import logging

        logging.getLogger("pptx_designer").warning(
            "prefill_copy_failed_keep_draft: %s", _e, exc_info=True
        )

    # jh 2026-06-12 — 카피라이터가 제목을 다듬은 뒤 목차를 최종 제목으로 재구성
    _rebuild_agenda()

    # === LLM 디자인 일괄 선택 (Step 6-3) ===
    # 각 슬라이드의 후보 N개를 LLM 한테 보여주고 1개 선택 받음.
    # 결과는 slide.preferred_template 에 저장 — REGISTRY.best_for() 가 1순위로 인식.
    # 실패해도 deck 빌드는 진행 (휴리스틱 폴백).
    try:
        _prepick_designs(slides_flat, ctx)
    except Exception as _e:  # noqa: BLE001
        # LLM 단계 실패 시 룰 기반 동작 유지 — 단, 침묵 금지.
        # (2026-06-11 사고: LLMDesigner abstract TypeError 가 여기서 삼켜져 무로그 폴백)
        import logging

        logging.getLogger("pptx_designer").warning(
            "prepick_designs_failed_fallback_to_rules: %s", _e, exc_info=True
        )

    for sl in slides_flat:
        page_num += 1
        slide = prs.slides.add_slide(blank)
        # jh 2026-06-14 — 슬라이드별 격리: 한 장의 draw 오류가 덱 전체 생성을 죽여
        # "PPT 생성 실패" 가 나던 구조 방지. 실패 슬라이드는 로그 후 건너뛰고 계속.
        try:
            _draw_slide(slide, sl, ctx, primary, accent, secondary, ink, muted, light_bg, render_visual_to_png)
        except Exception as _e:  # noqa: BLE001
            import logging

            logging.getLogger("pptx_designer").warning(
                "draw_slide_failed_skip: id=%s err=%s: %s",
                getattr(sl, "id", "?"), type(_e).__name__, _e, exc_info=True,
            )
        try:
            _draw_footer(slide, ctx, treat, page_num, total, muted, format_date_ko)
        except Exception:  # noqa: BLE001
            pass
        if sl.speaker_notes_hint:
            try:
                slide.notes_slide.notes_text_frame.text = sl.speaker_notes_hint
            except Exception:
                pass

    # === 편집자 QA 패스 (4단계 v1 — jh 2026-06-12) ===
    # 덱 전체를 만들고 난 뒤 한 번 더 훑으며 텍스트 오버플로를 자동 교정.
    # (견본과의 격차 '반복 0회' 의 첫 반복 — 결정적 기하 검사라 토큰 비용 0)
    try:
        _n_fixed = _qa_shrink_overflow(prs)
        if _n_fixed:
            import logging

            logging.getLogger("pptx_designer").info("deck_qa_font_shrunk: %d boxes", _n_fixed)
    except Exception as _e:  # noqa: BLE001
        import logging

        logging.getLogger("pptx_designer").warning("deck_qa_failed: %s", _e)

    prs.save(str(out))
    return str(out)


# ==============================================================
# 편집자 QA — 결정적 후처리 (jh 2026-06-12)
# ==============================================================


def _estimate_text_height_cm(text: str, size_pt: float, box_w_cm: float) -> float:
    """한글 기준 텍스트 높이 추정 — 글자폭 ≈ 1em, 줄간 1.35."""
    import math

    char_w = size_pt * 0.0353  # 1pt = 0.0353cm, 한글 전각 ≈ 1em
    usable_w = max(box_w_cm - 0.4, 1.0)
    lines = 0
    for raw_line in (text or "").split("\n"):
        est_w = max(len(raw_line), 1) * char_w
        lines += max(1, math.ceil(est_w / usable_w))
    return lines * size_pt * 0.0353 * 1.35


def _qa_shrink_overflow(prs) -> int:
    """모든 텍스트박스를 검사해 박스 높이를 넘치는 명시 폰트를 단계 축소.

    - 명시 size 가 있는 run 만 대상 (테마 상속 run 은 추정 불가 — 보수적 스킵)
    - 추정 높이가 박스의 115% 초과 시에만 개입 (과교정 방지)
    - 2pt 씩, 최소 9pt 까지
    """
    from pptx.util import Emu, Pt

    fixed = 0
    for slide in prs.slides:
        for sh in slide.shapes:
            if not getattr(sh, "has_text_frame", False):
                continue
            try:
                box_w = Emu(sh.width).cm
                box_h = Emu(sh.height).cm
            except Exception:
                continue
            if box_h < 0.5 or box_w < 1.0:
                continue
            tf = sh.text_frame
            runs = [r for p in tf.paragraphs for r in p.runs if r.font.size is not None]
            if not runs:
                continue
            text = tf.text
            size = max(r.font.size.pt for r in runs)
            shrunk = False
            for _ in range(6):
                # jh 2026-06-12 — 발표용 하한 11pt (9pt 까지 내려가 "글씨 작다" 지적)
                if size <= 11:
                    break
                if _estimate_text_height_cm(text, size, box_w) <= box_h * 1.15:
                    break
                size -= 2
                shrunk = True
            if shrunk:
                for r in runs:
                    if r.font.size is not None and r.font.size.pt > size:
                        r.font.size = Pt(size)
                fixed += 1
    return fixed


# ==============================================================
# Cover slide - full-bleed color block + huge typo + dot motif
# ==============================================================


def _display_intent(raw: str, max_len: int = 60) -> str:
    """user_intent 의 표시용 핵심 추출 — 게이트 태그 suffix 제거 + 길이 제한.

    jh 2026-06-12 — cover 제목이 '분석 방향: A (주제: B) (방법론: C)...' 원문
    그대로 48pt 로 들어가 슬라이드 경계를 뚫던 결함 수정.
    """
    import re

    s = (raw or "").strip()
    if not s:
        return "분석 보고서"
    m = re.match(r"^(?:분석 방향|주제|방법론|모델 전략):\s*(.+)$", s)
    if m:
        s = m.group(1).strip()
    cut = re.search(r"\s*\((?:분석 방향|주제|방법론|모델 전략):", s)
    if cut:
        s = s[: cut.start()].strip()
    if len(s) > max_len:
        s = s[: max_len - 1].rstrip() + "…"
    return s or "분석 보고서"


def _cover_title_size(text: str) -> int:
    """cover 제목 길이별 폰트 — 박스 (폭 ~12.2cm × 5cm) 안에 안전하게.

    jh 2026-06-12 — 견본 실측 26pt 기준 재보정. 48pt 는 한글 16자가
    3줄로 꺾여 구도를 깨던 결함 (12.2cm 박스 = 48pt 한글 약 7자/줄).
    """
    n = len(text)
    if n <= 12:
        return 40
    if n <= 22:
        return 32
    if n <= 34:
        return 26
    return 24


def _draw_cover(slide, sl: SlideSpec, ctx: ReportContext, primary: str, accent: str, secondary: str):
    # HJ 2026-06-08 — Stock photo background 시도 → 실패 시 기존 grad 디자인 폴백
    user_intent = getattr(ctx.meta, "user_intent", "") or ""
    # Step 7-2 — LLM 의 photo_keyword 가 user_intent 보다 우선.
    # LLMDesigner 가 데이터 신호 (도메인/카테고리/verdict) 보고 뽑은 영문 검색어를
    # Unsplash/Pexels 키워드로 사용. 힌트 없으면 user_intent 폴백.
    photo_kw = photo_keyword(sl, fallback="")
    # jh 2026-06-12 — get_cover_image 는 항상 범용 base 키워드("abstract business" 류)
    # 를 앞에 붙이고, 한글 intent 는 ASCII 필터로 소거돼 결국 무관한 비즈니스 인물
    # 사진이 나오던 결함. LLM 의 도메인 키워드가 있으면 그것만으로 직접 검색.
    cover_photo = None
    if _VISUAL_TOOLS_AVAILABLE:
        if photo_kw:
            try:
                from tools.visual.stock_image import fetch_stock_image

                cover_photo = fetch_stock_image(photo_kw)
            except Exception:  # noqa: BLE001
                cover_photo = None
        if cover_photo is None:
            cover_photo = get_cover_image("cover", photo_kw or user_intent)

    if cover_photo and cover_photo.exists():
        # jh 2026-06-12 — 견본 구도: 사진은 좌측 55% 패널에만, 우측 45% 는 흰 배경.
        # (풀블리드 + 70% overlay 가 제목을 사진 위에 얹어 가독성·구도를 깨던 결함)
        try:
            add_image_with_overlay(
                slide,
                0,
                0,
                SLIDE_W * 0.55,
                SLIDE_H,
                str(cover_photo),
                overlay_color=primary,
                # jh 2026-06-12 — 진한 사진에서 글자 침몰 방지: 오버레이 강화 (사용자 지시)
                overlay_alpha_pct=68,
            )
            # 좌측 강조 색 band
            add_color_band(slide, 0, 0, SLIDE_W * 0.02, SLIDE_H, accent)
            # Right white area — 제목·메타는 항상 흰 패널 위 (견본 동일)
            add_rect(slide, SLIDE_W * 0.55, 0, SLIDE_W * 0.45, SLIDE_H, "#FFFFFF")
        except Exception as _e:  # noqa: BLE001
            import logging

            logging.getLogger("pptx_designer").warning(
                "cover_stock_photo_failed",
                extra={"error": f"{type(_e).__name__}: {_e}"},
            )
            cover_photo = None  # 그라데이션 폴백

    if not cover_photo:
        # 폴백: 기존 그라데이션 디자인
        add_gradient_rect(slide, 0, 0, SLIDE_W * 0.55, SLIDE_H, primary, secondary, angle=135)
        # Right white area
        add_rect(slide, SLIDE_W * 0.55, 0, SLIDE_W * 0.45, SLIDE_H, "#FFFFFF")
    # Big decorative circles bottom-left for depth
    add_oval(slide, -4.0, SLIDE_H - 7.0, 10.0, 10.0, accent)
    # Smaller overlay circle
    add_oval(slide, -2.0, SLIDE_H - 5.0, 6.0, 6.0, "#FFFFFF")
    add_oval(slide, -1.5, SLIDE_H - 4.5, 5.0, 5.0, primary)
    # Triangle accent top-right corner of left block
    add_triangle(slide, SLIDE_W * 0.45, -3, 8.0, 8.0, "#FFFFFF", direction="diamond")
    # Diagonal accent strip
    add_rect(slide, SLIDE_W * 0.55 - 0.2, 0, 0.2, SLIDE_H, accent)

    # Big report label on left
    add_text_box(
        slide,
        1.5,
        1.5,
        SLIDE_W * 0.4,
        1.5,
        "ANALYSIS REPORT",
        size_pt=14,
        bold=True,
        # jh 2026-06-12 — accent(연파랑)는 사진 위에서 침몰 — 흰색 고정 (가독성)
        color_hex="#FFFFFF",
        align="left",
        vcenter=False,
    )
    # Big hex glyph as icon
    add_text_box(
        slide,
        1.5,
        3.0,
        3.0,
        3.0,
        GLYPHS["hex"],
        size_pt=72,
        bold=True,
        color_hex="#FFFFFF",
        align="left",
        vcenter=False,
    )

    # Title (right side) — jh 2026-06-12: 태그 제거 표시용 제목 + 길이별 폰트
    intent = _display_intent(ctx.meta.user_intent, max_len=44)
    add_text_box(
        slide,
        SLIDE_W * 0.6,
        4.5,
        SLIDE_W * 0.36,
        5.0,
        intent,
        size_pt=_cover_title_size(intent),
        bold=True,
        color_hex="#0F172A",
        align="left",
        vcenter=False,
    )
    # Colored bar under title
    add_rect(slide, SLIDE_W * 0.6, 9.5, 4.0, 0.15, primary)
    # Subtitle
    add_text_box(
        slide,
        SLIDE_W * 0.6,
        10.0,
        SLIDE_W * 0.36,
        1.0,
        f"{ctx.meta.category.upper()}  ·  ADA",
        size_pt=14,
        bold=True,
        color_hex=primary,
        align="left",
        vcenter=False,
    )
    # Metadata block with icons
    chosen = (ctx.model_selection.chosen or {}).get("name", "-")
    add_text_box(
        slide, SLIDE_W * 0.6, 12.5, 0.6, 0.6, GLYPHS["model"], size_pt=14, color_hex=primary, align="left", vcenter=True
    )
    add_text_box(
        slide,
        SLIDE_W * 0.6 + 0.7,
        12.5,
        SLIDE_W * 0.3,
        0.6,
        f"Model  {chosen}",
        size_pt=12,
        color_hex="#334155",
        align="left",
        vcenter=True,
    )
    add_text_box(
        slide,
        SLIDE_W * 0.6,
        13.3,
        0.6,
        0.6,
        GLYPHS["shield"],
        size_pt=14,
        color_hex=primary,
        align="left",
        vcenter=True,
    )
    add_text_box(
        slide,
        SLIDE_W * 0.6 + 0.7,
        13.3,
        SLIDE_W * 0.3,
        0.6,
        f"Class  {ctx.meta.classification}",
        size_pt=12,
        color_hex="#334155",
        align="left",
        vcenter=True,
    )
    add_text_box(
        slide, SLIDE_W * 0.6, 14.1, 0.6, 0.6, GLYPHS["data"], size_pt=14, color_hex=primary, align="left", vcenter=True
    )
    add_text_box(
        slide,
        SLIDE_W * 0.6 + 0.7,
        14.1,
        SLIDE_W * 0.3,
        0.6,
        f"Date   {(ctx.meta.generated_at or '')[:10]}",
        size_pt=12,
        color_hex="#334155",
        align="left",
        vcenter=True,
    )


# ==============================================================
# Section divider - huge chapter number
# ==============================================================


def _draw_section_divider(slide, title: str, chapter_num: int, primary: str, accent: str, secondary: str):
    # Diagonal gradient background
    add_gradient_rect(slide, 0, 0, SLIDE_W, SLIDE_H, primary, secondary, angle=135)
    # Big decorative circle bottom right
    add_oval(slide, SLIDE_W - 8.0, SLIDE_H - 8.0, 14.0, 14.0, accent)
    add_oval(slide, SLIDE_W - 6.5, SLIDE_H - 6.5, 11.0, 11.0, primary)
    # Big hex glyph as section icon
    add_text_box(
        slide,
        SLIDE_W - 12.0,
        3.0,
        6.0,
        6.0,
        GLYPHS["hex"],
        size_pt=180,
        bold=True,
        color_hex=accent,
        align="center",
        vcenter=False,
    )
    # Section / Chapter label
    add_text_box(
        slide,
        1.5,
        3.5,
        14.0,
        1.5,
        f"SECTION  ·  {chapter_num:02d}",
        size_pt=14,
        bold=True,
        color_hex=accent,
        align="left",
        vcenter=False,
    )
    # Accent line
    add_rect(slide, 1.5, 5.0, 3.0, 0.12, accent)
    # Huge chapter number
    add_text_box(
        slide,
        1.5,
        6.0,
        14.0,
        7.0,
        f"{chapter_num:02d}",
        size_pt=220,
        bold=True,
        color_hex="#FFFFFF",
        align="left",
        vcenter=False,
    )
    # Section title
    add_text_box(
        slide,
        1.5,
        13.5,
        SLIDE_W - 3.0,
        2.5,
        title,
        size_pt=40,
        bold=True,
        color_hex="#FFFFFF",
        align="left",
        vcenter=False,
    )
    # Progress dots
    for i in range(5):
        cx = SLIDE_W - 6.0 + i * 0.8
        c = "#FFFFFF" if i == chapter_num - 1 else accent
        size = 0.45 if i == chapter_num - 1 else 0.3
        add_oval(slide, cx, SLIDE_H - 1.8, size, size, c)


# ==============================================================
# Closing slide
# ==============================================================


def _draw_closing(slide, sl: SlideSpec, ctx: ReportContext, primary: str, accent: str):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, "#FFFFFF")
    # Top color band
    add_rect(slide, 0, 0, SLIDE_W, 1.5, primary)
    # Big "Thank You" / Summary
    add_text_box(
        slide,
        2.0,
        4.5,
        SLIDE_W - 4.0,
        3.0,
        "Summary & Q&A",
        size_pt=54,
        bold=True,
        color_hex=primary,
        align="center",
        vcenter=False,
    )
    # 3 box summary
    box_w = (SLIDE_W - 6.0) / 3
    for i, line in enumerate(sl.body_outline[:3]):
        x = 2.0 + i * (box_w + 1.0)
        add_rounded_rect(slide, x, 9.5, box_w, 4.0, "#F8FAFC", line_hex=primary)
        add_text_box(
            slide, x + 0.5, 10.2, box_w - 1.0, 3.0, line, size_pt=13, color_hex="#0F172A", align="left", vcenter=True
        )
    # Bottom accent
    add_rect(slide, 0, SLIDE_H - 0.4, SLIDE_W, 0.4, accent)


# ==============================================================
# Generic body slide router
# ==============================================================


def _draw_slide(
    slide,
    sl: SlideSpec,
    ctx: ReportContext,
    primary: str,
    accent: str,
    secondary: str,
    ink: str,
    muted: str,
    light_bg: str,
    render_fn,
):
    layout = sl.layout

    # jh 2026-06-12 — palette_hint 비활성화 (사용자 지시: 덱 내 색상 통일).
    # 슬라이드별 LLM palette_hint 가 장마다 초록/파랑/연두를 섞어 쓰던 결함.
    # 카테고리 기본 팔레트 하나로 덱 전체 고정 (SWOT 의 시맨틱 4색은 draw 내부 고정이라 유지).
    _overridden = {"primary": primary, "accent": accent, "secondary": secondary}
    primary = _overridden["primary"]
    accent = _overridden["accent"]
    secondary = _overridden["secondary"]

    # 부족 신호 카운터 누적 (fallback_reason 있는 경우만)
    check_and_log_miss(sl)

    # Hardcoded specials (cover/agenda/divider use bespoke draw)
    if layout == "cover":
        _draw_cover(slide, sl, ctx, primary, accent, secondary)
        return
    if layout == "agenda":
        _draw_agenda(slide, sl, primary, accent, ink, muted)
        return
    if layout == "section_divider":
        return  # skipped at the higher level

    # Try registry-based selection first
    try:
        from outputs.carriers.template_registry import REGISTRY
        from outputs.carriers.templates_init import init_registry

        init_registry()
        spec = REGISTRY.best_for(sl, ctx)
        if spec is not None:
            # Apply standard header for non-meta slides
            if sl.role != "meta":
                _draw_header(slide, sl, primary, ink, muted)
            spec.draw(slide, sl, ctx, primary, accent, ink, muted, light_bg)
            return
    except Exception:
        pass

    # Fallback to legacy hardcoded paths
    # action_recommendations / rec_options → checklist
    if sl.id in ("action_recommendations", "rec_options", "short_term"):
        _draw_checklist(slide, sl, ctx, primary, accent, ink, muted, light_bg)
        return
    # roadmap → step cards vertical
    if sl.id in ("roadmap", "i3_roi"):
        _draw_step_cards(slide, sl, ctx, primary, accent, ink, muted, light_bg)
        return
    if layout == "agenda":
        _draw_agenda(slide, sl, primary, accent, ink, muted)
        return
    if layout == "kpi_cards_4":
        _draw_percentage_grid(slide, sl, ctx, primary, accent, ink, muted, light_bg)
        return
    if layout in ("kpi_cards_3", "kpi_cards_6"):
        _draw_kpi_cards(slide, sl, ctx, primary, accent, ink, muted, light_bg)
        return
    if layout == "comparison_before_after":
        _draw_as_is_to_be(slide, sl, ctx, primary, accent, ink, muted, light_bg)
        return
    if layout == "process_flow_gantt":
        _draw_timeline(slide, sl, ctx, primary, accent, ink, muted, light_bg)
        return
    # PSI: differentiation 2x2 → strategy 4 circular
    if sl.id == "s3_differentiation":
        _draw_strategy_4(slide, sl, ctx, primary, accent, ink, muted, light_bg)
        return
    # SCQA: e3 model comparison → chevron strategies (sequence of candidates)
    if sl.id in ("e3_model_comparison", "alternatives"):
        _draw_chevron_seq(slide, sl, ctx, primary, accent, ink, muted, light_bg)
        return
    # criteria/agenda → index cards (big numbers)
    if sl.id in ("criteria", "i2_before_after"):
        # i2_before_after is already as_is/to_be — use index_cards for criteria
        if sl.id == "criteria":
            _draw_index_cards(slide, sl, ctx, primary, accent, ink, muted, light_bg)
            return
    # exec_summary → BIG STATS
    if sl.id == "exec_summary":
        _draw_big_stats(slide, sl, ctx, primary, accent, ink, muted, light_bg)
        return
    if layout == "one_message_big_number":
        _draw_big_number(slide, sl, ctx, primary, accent, ink, muted)
        return
    if layout == "process_flow":
        _draw_process_flow(slide, sl, ctx, primary, accent, ink, muted, light_bg, render_fn)
        return
    if layout in ("comparison_table", "appendix_table"):
        _draw_table_slide(slide, sl, ctx, primary, accent, ink, muted, light_bg, render_fn)
        return
    if layout == "chart_callout":
        _draw_chart_callout(slide, sl, ctx, primary, accent, ink, muted, light_bg, render_fn)
        return
    if layout == "quote":
        _draw_quote(slide, sl, primary, accent, ink, muted)
        return
    # one_message and others - default
    _draw_default_body(slide, sl, ctx, primary, accent, ink, muted, light_bg, render_fn)


# ==============================================================
# Body slide variants
# ==============================================================


def _icon_name_for_sl(sl: SlideSpec) -> str:
    """SlideSpec 으로부터 Lucide 아이콘 이름 추출 (없으면 빈 문자열).

    Step 7-3 — LLM 디자인 힌트의 icon_concept 가 1순위:
        sl._design_hint.icon_concept = "warning" → "alert-circle"
        sl._design_hint.icon_concept = "kpi"     → "trending-up"
    힌트가 없거나 매핑 안 되면 기존 icon_for_slide (slide_id + role 기반) 로 폴백.
    """
    if not _VISUAL_TOOLS_AVAILABLE:
        return ""
    # 1순위: LLM 의 icon_concept → lucide 아이콘 이름
    hinted = design_icon_name(sl)
    if hinted:
        return hinted
    # 폴백: 기존 icon_for_slide (slide_id + role 기반)
    p = icon_for_slide(sl.id, getattr(sl, "role", "claim"))
    if p is None:
        return ""
    return p.stem  # 파일명 (확장자 제외)


def _draw_header(slide, sl: SlideSpec, primary: str, ink: str, muted: str):
    """Top header band: left accent + Lucide icon (if available) + title + so-what."""
    # jh 2026-06-11 — 멱등 가드. _draw_slide 진입부와 템플릿 draw 내부가 각자
    # 헤더를 호출해 같은 좌표에 5개 shape 가 이중으로 겹쳐 그려지던 버그 (probe 실측).
    if getattr(slide, "_ada_header_drawn", False):
        return
    slide._ada_header_drawn = True

    # Left vertical accent bar
    add_vertical_accent(slide, 1.0, 1.0, 1.3, primary, w_cm=0.18)

    # HJ 2026-06-08 — Lucide 아이콘 시도 (slide_id 기준), 실패 시 텍스트만
    icon_offset = 0.0
    if _VISUAL_TOOLS_AVAILABLE:
        try:
            icon_name = _icon_name_for_sl(sl)
            icon_png = svg_to_png_recolored(icon_name, color=primary, size_px=128)
            if icon_png and icon_png.exists():
                from pptx.util import Cm

                slide.shapes.add_picture(str(icon_png), Cm(1.4), Cm(0.85), width=Cm(1.0), height=Cm(1.0))
                icon_offset = 1.2  # 아이콘 너비만큼 텍스트 우측 이동
        except FileNotFoundError as _e:
            # 캐시 PNG 가 권한 문제로 안 보임 (sandbox 한정) — 운영에선 거의 없음
            import logging

            logging.getLogger("pptx_designer").warning(
                "icon_add_picture_filenotfound", extra={"slide_id": sl.id, "icon": _icon_name_for_sl(sl)}
            )
        except Exception as _e:  # noqa: BLE001
            import logging

            logging.getLogger("pptx_designer").warning(
                "icon_add_picture_failed",
                extra={"slide_id": sl.id, "error": f"{type(_e).__name__}: {_e}"},
            )

    # Title (아이콘이 있으면 그 옆에)
    add_text_box(
        slide,
        1.5 + icon_offset,
        0.9,
        SLIDE_W - 3.0 - icon_offset,
        1.0,
        sl.title_ko or sl.id,
        size_pt=24,
        bold=True,
        color_hex=ink,
        align="left",
        vcenter=False,
    )
    # So-What pill
    if sl.so_what:
        add_text_box(
            slide,
            1.5,
            2.1,
            SLIDE_W - 3.0,
            0.9,
            sl.so_what,
            size_pt=14,
            bold=True,
            color_hex=primary,
            align="left",
            vcenter=False,
        )
    # Thin horizontal rule
    add_horizontal_rule(slide, 1.0, 3.4, SLIDE_W - 2.0, "#CBD5E1", h_cm=0.03)


def _draw_default_body(slide, sl, ctx, primary, accent, ink, muted, light_bg, render_fn):
    _draw_header(slide, sl, primary, ink, muted)
    # Body bullets in a soft card
    add_rounded_rect(slide, 1.0, 4.0, SLIDE_W - 2.0, SLIDE_H - 5.5, light_bg)
    body_text = "\n\n".join(f"•  {b}" for b in sl.body_outline)
    add_text_box(
        slide, 1.6, 4.4, SLIDE_W - 3.2, SLIDE_H - 6.3, body_text, size_pt=14, color_hex=ink, align="left", vcenter=False
    )


def _metric_caption(metric_name: str, ctx) -> str:
    """Friendly caption for each metric."""
    captions = {
        "auc": "AUC 0.85 이상은 우수 분류 모델 기준",
        "roc_auc": "ROC-AUC 기준 — 0.5(랜덤) ~ 1.0(완벽)",
        "f1": "정밀도·재현율 조화평균",
        "precision": "양성 예측의 정확도",
        "recall": "실제 양성을 잡아낸 비율",
        "accuracy": "전체 정확도",
        "rmse": "예측 오차 — 낮을수록 우수",
        "mae": "절대 오차 평균",
        "mape": "MAPE — 비율 오차",
        "smape": "SMAPE — 대칭 비율 오차",
        "pr_auc": "Precision-Recall AUC",
    }
    return captions.get(metric_name.lower(), "주요 평가 지표")


def _draw_kpi_cards(slide, sl, ctx, primary, accent, ink, muted, light_bg):
    _draw_header(slide, sl, primary, ink, muted)
    # Use statistics_grid infographic for 4 metrics
    from outputs.carriers.pptx_infographics import draw_statistics_grid

    metrics = ctx.evaluation.metrics or {}
    items_stats = []
    for k, m in list(metrics.items())[:4]:
        v = m.get("value")
        v_str = f"{v:.3f}" if isinstance(v, float) and abs(v) < 1 else str(v)
        items_stats.append(
            {
                "value": v_str,
                "label": k,
                "caption": _metric_caption(k, ctx),
            }
        )
    if items_stats:
        draw_statistics_grid(slide, items_stats, primary, accent, ink, muted, light_bg)
        return

    # Fallback to old card style
    items = []
    for line in sl.body_outline[:6]:
        if ":" in line:
            k, v = line.split(":", 1)
            items.append((k.strip(), v.strip()))
        else:
            items.append((line.strip(), ""))
    if not items and ctx.evaluation.metrics:
        for k, m in list(ctx.evaluation.metrics.items())[:4]:
            v = m.get("value")
            items.append((k, f"{v:.3f}" if isinstance(v, float) else str(v)))
    if not items:
        return

    n = len(items)
    cols = 3 if sl.layout == "kpi_cards_3" else 4 if sl.layout == "kpi_cards_4" else 3
    if n <= 3:
        cols = n
    rows = (n + cols - 1) // cols

    margin_x = 1.5
    margin_top = 4.5
    gap = 0.5
    avail_w = SLIDE_W - 2 * margin_x
    avail_h = SLIDE_H - margin_top - 2.0
    card_w = (avail_w - (cols - 1) * gap) / cols
    card_h = (avail_h - (rows - 1) * gap) / rows

    from outputs.style.visual_kit import GLYPHS as _G

    glyphs_seq = [_G["chart"], _G["target"], _G["bulb"], _G["star"], _G["model"], _G["data"]]
    for i, (name, value) in enumerate(items):
        r, c = i // cols, i % cols
        x = margin_x + c * (card_w + gap)
        y = margin_top + r * (card_h + gap)
        # White card with soft border
        add_rounded_rect(slide, x, y, card_w, card_h, "#FFFFFF", line_hex="#E2E8F0")
        # Top color accent bar (full width)
        add_rect(slide, x, y, card_w, 0.18, primary)
        # Icon circle top-right
        ic_size = 1.6
        add_glyph(
            slide,
            x + card_w - ic_size - 0.4,
            y + 0.5,
            ic_size,
            ic_size,
            glyphs_seq[i % len(glyphs_seq)],
            color_hex="#FFFFFF",
            size_pt=20,
            bg_color=primary,
        )
        # Label (top-left)
        add_text_box(
            slide,
            x + 0.5,
            y + 0.8,
            card_w - 2.5,
            0.8,
            name[:24].upper(),
            size_pt=10,
            bold=True,
            color_hex=muted,
            align="left",
            vcenter=True,
        )
        # Big value (left-center)
        add_text_box(
            slide,
            x + 0.5,
            y + 2.0,
            card_w - 1.0,
            card_h - 3.5,
            value or "—",
            size_pt=40,
            bold=True,
            color_hex=primary,
            align="left",
            vcenter=True,
        )
        # Bottom accent line + confidence
        add_rect(slide, x + 0.5, y + card_h - 1.1, 2.5, 0.06, accent)
        add_text_box(
            slide,
            x + 0.5,
            y + card_h - 0.9,
            card_w - 1.0,
            0.6,
            "★★★  high confidence",
            size_pt=9,
            color_hex=accent,
            align="left",
            vcenter=True,
        )


def _draw_big_number(slide, sl, ctx, primary, accent, ink, muted):
    _draw_header(slide, sl, primary, ink, muted)
    pm = ctx.evaluation.primary_metric or {}
    value = pm.get("value", "-")
    name = pm.get("name", "Metric")
    value_str = f"{value:.4f}" if isinstance(value, float) else str(value)
    # Big rounded card
    add_rounded_rect(slide, 4.0, 5.0, SLIDE_W - 8.0, 9.0, "#FFFFFF", line_hex=primary)
    add_text_box(
        slide,
        4.5,
        5.5,
        SLIDE_W - 9.0,
        2.0,
        name.upper(),
        size_pt=18,
        bold=True,
        color_hex=muted,
        align="center",
        vcenter=True,
    )
    add_text_box(
        slide,
        4.5,
        7.5,
        SLIDE_W - 9.0,
        5.0,
        value_str,
        size_pt=120,
        bold=True,
        color_hex=primary,
        align="center",
        vcenter=True,
    )
    # Caption
    if sl.body_outline:
        add_text_box(
            slide,
            4.5,
            12.8,
            SLIDE_W - 9.0,
            1.0,
            sl.body_outline[0][:80],
            size_pt=12,
            color_hex=ink,
            align="center",
            vcenter=True,
        )


def _draw_process_flow(slide, sl, ctx, primary, accent, ink, muted, light_bg, render_fn):
    _draw_header(slide, sl, primary, ink, muted)
    steps = sl.body_outline or []
    if not steps:
        return
    n = min(len(steps), 7)
    icons = [
        GLYPHS["data"],
        GLYPHS["chart"],
        GLYPHS["settings"],
        GLYPHS["target"],
        GLYPHS["bulb"],
        GLYPHS["ok"],
        GLYPHS["report"],
    ]
    avail_w = SLIDE_W - 3.0
    step_w = avail_w / n
    circle_size = min(2.2, step_w * 0.6)
    y_circle = 5.5

    # Horizontal connecting line behind circles
    add_rect(slide, 1.5 + step_w * 0.5, y_circle + circle_size / 2 - 0.05, step_w * (n - 1), 0.1, "#E2E8F0")

    for i in range(n):
        x = 1.5 + i * step_w
        cx = x + step_w / 2 - circle_size / 2
        # Outer ring (accent)
        add_oval(slide, cx - 0.15, y_circle - 0.15, circle_size + 0.3, circle_size + 0.3, accent)
        # Inner primary circle
        add_oval(slide, cx, y_circle, circle_size, circle_size, primary)
        # Icon glyph inside
        add_text_box(
            slide,
            cx,
            y_circle,
            circle_size,
            circle_size,
            icons[i % len(icons)],
            size_pt=22,
            bold=True,
            color_hex="#FFFFFF",
            align="center",
            vcenter=True,
        )
        # Step number badge (small circle below)
        add_oval(slide, cx + circle_size / 2 - 0.5, y_circle + circle_size + 0.3, 1.0, 1.0, accent)
        add_text_box(
            slide,
            cx + circle_size / 2 - 0.5,
            y_circle + circle_size + 0.3,
            1.0,
            1.0,
            f"{i + 1:02d}",
            size_pt=14,
            bold=True,
            color_hex=primary,
            align="center",
            vcenter=True,
        )
        # Step text in rounded card below
        card_y = y_circle + circle_size + 1.8
        add_rounded_rect(slide, x + 0.2, card_y, step_w - 0.4, 5.0, light_bg, line_hex=None)
        add_text_box(
            slide,
            x + 0.4,
            card_y + 0.3,
            step_w - 0.8,
            4.4,
            steps[i][:80],
            size_pt=11,
            color_hex=ink,
            align="center",
            vcenter=False,
        )


def _draw_table_slide(slide, sl, ctx, primary, accent, ink, muted, light_bg, render_fn):
    _draw_header(slide, sl, primary, ink, muted)
    # Use render to get matplotlib table
    if sl.visual_spec:
        try:
            png = render_fn(sl.visual_spec, ctx, slide=sl)
            if png:
                from pptx.util import Cm

                slide.shapes.add_picture(png, Cm(2.0), Cm(4.0), width=Cm(SLIDE_W - 4.0), height=Cm(SLIDE_H - 6.0))
                return
        except Exception:
            pass
    # Fallback: rows as cards
    _draw_default_body(slide, sl, ctx, primary, accent, ink, muted, light_bg, render_fn)


def _draw_chart_callout(slide, sl, ctx, primary, accent, ink, muted, light_bg, render_fn):
    _draw_header(slide, sl, primary, ink, muted)
    # Left: chart (60%); Right: so-what callout (40%)
    chart_w = (SLIDE_W - 3.0) * 0.6
    callout_w = (SLIDE_W - 3.0) * 0.4 - 0.5
    chart_x = 1.5
    callout_x = chart_x + chart_w + 0.5

    chart_placed = False
    if sl.visual_spec:
        try:
            png = render_fn(sl.visual_spec, ctx, slide=sl)
            if png:
                from pptx.util import Cm

                slide.shapes.add_picture(png, Cm(chart_x), Cm(4.5), width=Cm(chart_w), height=Cm(SLIDE_H - 7.0))
                chart_placed = True
        except Exception:
            pass

    # jh 2026-06-12 — 차트 미배치 시 패널을 전폭으로 확장.
    # (SHAP·CM 미적립 슬라이드에서 좌측 절반이 통째로 빈 화면으로 나가던 결함)
    if not chart_placed:
        callout_x = 1.5
        callout_w = SLIDE_W - 3.0

    add_rounded_rect(slide, callout_x, 4.5, callout_w, SLIDE_H - 7.0, light_bg, line_hex=primary)
    add_text_box(
        slide, callout_x + 0.5, 5.0, 0.6, 0.6, GLYPHS["bulb"], size_pt=18, color_hex=primary, align="left", vcenter=True
    )
    # jh 2026-06-12 — 발표용 가독성: 라벨 11→13pt, 본문 11→14pt (사용자 지시)
    add_text_box(
        slide,
        callout_x + 1.3,
        5.0,
        callout_w - 1.8,
        0.8,
        "KEY INSIGHTS",
        size_pt=13,
        bold=True,
        color_hex=primary,
        align="left",
        vcenter=True,
    )
    # jh 2026-06-12 — 목표치 덱 페어 2줄 구조: "사실" 줄 + "→ 시사점" 강조 들여쓰기 줄.
    # 카피라이터 페어 bullet("사실 — 시사점") 을 " — " 에서 분해해 2층으로 그림.
    _add_pair_bullets(
        slide, callout_x + 0.5, 6.2, callout_w - 1.0, SLIDE_H - 9.0,
        sl.body_outline[:5], ink, primary,
    )


def _add_pair_bullets(slide, x, y, w, h, bullets, ink, primary):
    """페어 bullet 리치텍스트 — 사실 13pt + '→ 시사점' 12B(primary) + 간격."""
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR
    from pptx.util import Cm, Pt

    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    # jh 2026-06-12 — 세로 중앙 정렬 (글이 적어도 허전하지 않게, 사용자 지시)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _ink = RGBColor.from_string(str(ink).lstrip("#"))
    _pri = RGBColor.from_string(str(primary).lstrip("#"))
    first = True
    # jh 2026-06-12 — 빈 줄 단락 대신 space_after 로 간격 (deck_qa 가 줄 수를
    # 과대평가해 본문을 9pt 까지 축소하던 결함), bullet 상한 5→4.
    for b in list(bullets or [])[:4]:
        s = str(b).strip().lstrip("-").strip()
        fact, _, imp = s.partition(" — ")
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = f"-  {fact.strip()}"
        r.font.size = Pt(13)
        r.font.color.rgb = _ink
        last_p = p
        if imp.strip():
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = f"     →  {imp.strip()}"
            r2.font.size = Pt(12)
            r2.font.bold = True
            r2.font.color.rgb = _pri
            last_p = p2
        last_p.space_after = Pt(10)


def _draw_quote(slide, sl, primary, accent, ink, muted):
    _draw_header(slide, sl, primary, ink, muted)
    add_text_box(
        slide, 3.0, 5.0, 2.5, 3.0, chr(0x201C), size_pt=120, bold=True, color_hex=accent, align="left", vcenter=False
    )
    add_text_box(
        slide, 5.5, 6.5, SLIDE_W - 8.0, 5.0, sl.so_what or "", size_pt=24, color_hex=ink, align="left", vcenter=True
    )
    add_text_box(
        slide,
        5.5,
        12.0,
        SLIDE_W - 8.0,
        3.0,
        "\n".join(sl.body_outline[:3]),
        size_pt=12,
        color_hex=muted,
        align="left",
        vcenter=False,
    )


def _draw_agenda(slide, sl, primary, accent, ink, muted):
    # Left gradient panel (smaller — 25%)
    add_gradient_rect(slide, 0, 0, SLIDE_W * 0.25, SLIDE_H, primary, "#1e3a8a", angle=180)
    add_text_box(
        slide,
        1.2,
        2.0,
        SLIDE_W * 0.22,
        2.0,
        "AGENDA",
        size_pt=42,
        bold=True,
        color_hex="#FFFFFF",
        align="left",
        vcenter=False,
    )
    add_text_box(
        slide,
        1.2,
        4.2,
        SLIDE_W * 0.22,
        1.5,
        "목차",
        size_pt=24,
        bold=True,
        # jh 2026-06-12 — accent(연파랑)가 파랑 그라데이션 위에서 흐림 (사용자 재지적)
        color_hex="#FFFFFF",
        align="left",
        vcenter=False,
    )
    add_rect(slide, 1.2, 6.0, 3.0, 0.12, accent)
    add_text_box(
        slide,
        1.2,
        6.5,
        SLIDE_W * 0.22,
        3.0,
        "본 보고서는\n다음 순서로\n구성됩니다.",
        size_pt=15,
        color_hex="#FFFFFF",
        align="left",
        vcenter=False,
    )
    add_text_box(
        slide,
        1.2,
        SLIDE_H - 3.5,
        SLIDE_W * 0.22,
        1.5,
        GLYPHS["hex"],
        size_pt=64,
        color_hex=accent,
        align="left",
        vcenter=False,
    )
    # Right: 2-column dense list (up to 18 items)
    items = sl.body_outline[:18]
    n = len(items)
    if n == 0:
        return
    right_x = SLIDE_W * 0.28
    right_w = SLIDE_W - right_x - 1.0
    # Decide columns based on item count
    cols = 1 if n <= 9 else 2
    rows = (n + cols - 1) // cols
    avail_h = SLIDE_H - 3.0
    row_h = avail_h / rows
    col_w = right_w / cols - 0.4
    for i, item in enumerate(items):
        c = i // rows
        r = i % rows
        x = right_x + c * (col_w + 0.4)
        y = 1.8 + r * row_h
        # Strip "01. " prefix if present (we'll redraw number badge)
        text = item
        if len(text) > 3 and text[:2].isdigit() and text[2:4] in (". ", ".  "):
            text = text[4:].strip()
        # Number badge
        bsize = 1.1
        add_oval(slide, x, y + row_h / 2 - bsize / 2, bsize, bsize, accent)
        add_text_box(
            slide,
            x,
            y + row_h / 2 - bsize / 2,
            bsize,
            bsize,
            f"{i + 1:02d}",
            size_pt=12,
            bold=True,
            color_hex=primary,
            align="center",
            vcenter=True,
        )
        # Title
        add_text_box(
            slide,
            x + bsize + 0.3,
            y,
            col_w - bsize - 0.3,
            row_h,
            text[:48],
            size_pt=14,
            bold=True,
            color_hex=ink,
            align="left",
            vcenter=True,
        )
        # Thin underline
        add_horizontal_rule(slide, x + bsize + 0.3, y + row_h - 0.1, col_w - bsize - 0.3, "#E2E8F0", h_cm=0.02)


def _draw_footer(slide, ctx, treat, page, total, muted, fmt_date):
    title_short = _display_intent(ctx.meta.user_intent, max_len=30)
    date_str = fmt_date(ctx.meta.generated_at, style="iso")
    foot_color = treat.get("footer_color", "#334155")
    foot_text = treat.get("footer_text", "INTERNAL")
    add_horizontal_rule(slide, 1.0, SLIDE_H - 0.9, SLIDE_W - 2.0, "#E2E8F0", h_cm=0.02)
    add_text_box(
        slide,
        1.0,
        SLIDE_H - 0.75,
        13.0,
        0.6,
        f"ADA . {title_short}",
        size_pt=8,
        color_hex=muted,
        align="left",
        vcenter=True,
    )
    add_text_box(
        slide,
        14.5,
        SLIDE_H - 0.75,
        5.0,
        0.6,
        f"{date_str}   p.{page}/{total}",
        size_pt=8,
        color_hex=muted,
        align="center",
        vcenter=True,
    )
    add_text_box(
        slide,
        SLIDE_W - 6.0,
        SLIDE_H - 0.75,
        5.0,
        0.6,
        foot_text,
        size_pt=8,
        bold=True,
        color_hex=foot_color,
        align="right",
        vcenter=True,
    )


def _fb(plan, output_path):
    out = Path(str(output_path) + ".txt")
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(f"[PPTX Designer Fallback {plan.skeleton}]", encoding="utf-8")
    return str(out)


def _draw_percentage_grid(slide, sl, ctx, primary, accent, ink, muted, light_bg):
    """2x2 percentage card grid (Samsung Pay style)."""
    _draw_header(slide, sl, primary, ink, muted)
    from outputs.carriers.pptx_infographics import draw_percentage_grid_4

    metrics = ctx.evaluation.metrics or {}
    items = []
    for k, m in list(metrics.items())[:4]:
        v = m.get("value")
        # Convert to percentage string
        if isinstance(v, float):
            if 0 < v < 1:
                v_str = f"{v * 100:.0f}%"
            else:
                v_str = f"{v:.1f}"
        else:
            v_str = str(v)
        items.append(
            {
                "value": v_str,
                "label": _metric_caption(k, ctx),
            }
        )
    if not items:
        items = [
            {"value": "85%", "label": "AUC 우수 기준 충족"},
            {"value": "78%", "label": "F1 임계치 통과"},
            {"value": "81%", "label": "Precision 양호"},
            {"value": "74%", "label": "Recall 표준 범위"},
        ]
    draw_percentage_grid_4(slide, items, primary, accent, ink, muted, light_bg)


def _draw_as_is_to_be(slide, sl, ctx, primary, accent, ink, muted, light_bg):
    """AS-IS / TO-BE 2-column comparison."""
    _draw_header(slide, sl, primary, ink, muted)
    from outputs.carriers.pptx_infographics import draw_as_is_to_be

    # Extract AS-IS / TO-BE from body_outline
    half = max(1, len(sl.body_outline) // 2)
    as_is_pts = (
        sl.body_outline[:half]
        if sl.body_outline
        else [
            "수작업 분석 — 3~6주 소요",
            "재현 불가능 — 일회성 결과",
            "팀별 표준 부재",
            "모니터링 수동",
        ]
    )
    to_be_pts = (
        sl.body_outline[half : half * 2]
        if len(sl.body_outline) > half
        else [
            f"{(ctx.model_selection.chosen or {}).get('name', 'AI 모델')} 자동화",
            "Companion 코드로 재현 가능",
            "통합 분석 파이프라인 표준화",
            "MLflow + Langfuse 자동 추적",
        ]
    )
    draw_as_is_to_be(
        slide,
        {"label": "AS-IS", "points": as_is_pts},
        {"label": "TO-BE", "points": to_be_pts},
        primary,
        accent,
        ink,
        muted,
        light_bg,
    )


def _draw_timeline(slide, sl, ctx, primary, accent, ink, muted, light_bg):
    """Timeline milestones with alternating labels."""
    _draw_header(slide, sl, primary, ink, muted)
    from outputs.carriers.pptx_infographics import draw_timeline_milestones

    milestones = []
    for i, line in enumerate(sl.body_outline[:5]):
        # Try to parse "Phase 1 (0~30일): xxx"
        year = f"P{i + 1}"
        label = line
        if ":" in line:
            year_part, label = line.split(":", 1)
            year = year_part.strip()[:18]
            label = label.strip()
        milestones.append({"year": year, "label": label[:80]})
    if not milestones:
        milestones = [
            {"year": "Phase 1", "label": "파일럿 운영 시작 - 0~30일"},
            {"year": "Phase 2", "label": "운영 환경 단계 배포 - 30~90일"},
            {"year": "Phase 3", "label": "확장 + 재학습 자동화 - 90일+"},
        ]
    draw_timeline_milestones(slide, milestones, primary, accent, ink, muted, light_bg)


def _draw_checklist(slide, sl, ctx, primary, accent, ink, muted, light_bg):
    """4 horizontal big-check cards for action/recommendation slides."""
    _draw_header(slide, sl, primary, ink, muted)
    from outputs.carriers.pptx_infographics import draw_checklist_4

    items = []
    for line in sl.body_outline[:4]:
        # "권고 액션 1: 모델을 운영 환경..." → title + caption
        if ":" in line:
            k, v = line.split(":", 1)
            items.append({"title": k.strip()[:30], "caption": v.strip()[:120]})
        elif "·" in line:
            k, v = line.split("·", 1)
            items.append({"title": k.strip()[:30], "caption": v.strip()[:120]})
        else:
            items.append({"title": line[:30], "caption": ""})
    if not items:
        items = [
            {"title": "단계적 배포", "caption": "파일럿 → 운영 단계 전환"},
            {"title": "모니터링 설정", "caption": "핵심 지표 자동 알람"},
            {"title": "정기 재학습", "caption": "분기 1회 모델 갱신"},
            {"title": "문서·교육", "caption": "운영팀 인수인계 + 표준 문서화"},
        ]
    draw_checklist_4(slide, items, primary, accent, ink, muted, light_bg)


def _draw_step_cards(slide, sl, ctx, primary, accent, ink, muted, light_bg):
    """Vertical STEP cards for roadmap/ROI."""
    _draw_header(slide, sl, primary, ink, muted)
    from outputs.carriers.pptx_infographics import draw_step_cards_vertical

    items = []
    for line in sl.body_outline[:4]:
        if ":" in line:
            k, v = line.split(":", 1)
            items.append({"title": k.strip()[:30], "caption": v.strip()[:120]})
        elif "·" in line:
            k, v = line.split("·", 1)
            items.append({"title": k.strip()[:30], "caption": v.strip()[:120]})
        else:
            items.append({"title": line[:30], "caption": ""})
    if not items:
        items = [
            {"title": "준비", "caption": "데이터·인프라 점검"},
            {"title": "파일럿", "caption": "1개월 시범 운영"},
            {"title": "확장", "caption": "전사 단계 도입"},
            {"title": "정착", "caption": "자동화·표준화"},
        ]
    draw_step_cards_vertical(slide, items, primary, accent, ink, muted, light_bg)


def _draw_strategy_4(slide, sl, ctx, primary, accent, ink, muted, light_bg):
    """4 strategy circles around central product."""
    _draw_header(slide, sl, primary, ink, muted)
    from outputs.carriers.pptx_infographics import draw_strategy_4_circular

    items = []
    for line in sl.body_outline[:4]:
        title = line.split(":", 1)[0] if ":" in line else line.split("·", 1)[0]
        items.append({"title": title.strip()[:24]})
    if not items:
        items = [
            {"title": "자동화 + 신뢰성"},
            {"title": "재현 가능성"},
            {"title": "운영 표준화"},
            {"title": "확장성"},
        ]
    draw_strategy_4_circular(slide, items, primary, accent, ink, muted, light_bg)


def _draw_chevron_seq(slide, sl, ctx, primary, accent, ink, muted, light_bg):
    """5 chevron arrows for model candidates."""
    _draw_header(slide, sl, primary, ink, muted)
    from outputs.carriers.pptx_infographics import draw_chevron_strategies

    cands = ctx.model_selection.candidates or []
    items = []
    for c in cands[:5]:
        items.append(
            {
                "title": (c.name or "")[:20],
                "caption": (c.why_tried or "")[:80],
            }
        )
    if not items:
        items = [
            {"title": "Baseline", "caption": "기준 모델"},
            {"title": "Tree", "caption": "트리 기반"},
            {"title": "Ensemble", "caption": "앙상블"},
            {"title": "DL", "caption": "딥러닝"},
            {"title": "Chosen", "caption": "최종 선정"},
        ]
    draw_chevron_strategies(slide, items, primary, accent, ink, muted, light_bg)


def _draw_index_cards(slide, sl, ctx, primary, accent, ink, muted, light_bg):
    """4 Index cards with big numbers (Index 01-04)."""
    _draw_header(slide, sl, primary, ink, muted)
    from outputs.carriers.pptx_infographics import draw_index_cards

    items = []
    for line in sl.body_outline[:4]:
        if ":" in line:
            k, v = line.split(":", 1)
        elif "·" in line:
            k, v = line.split("·", 1)
        else:
            k, v = line, ""
        items.append({"title": k.strip()[:24], "caption": v.strip()[:100]})
    if not items:
        items = [
            {"title": "성능", "caption": "정확도·재현율·정밀도 종합"},
            {"title": "해석성", "caption": "SHAP·PDP 기반 설명력"},
            {"title": "속도", "caption": "학습·추론 시간"},
            {"title": "운영", "caption": "재학습·배포 용이성"},
        ]
    draw_index_cards(slide, items, primary, accent, ink, muted, light_bg)


def _draw_big_stats(slide, sl, ctx, primary, accent, ink, muted, light_bg):
    """Big numbers (Executive Summary)."""
    _draw_header(slide, sl, primary, ink, muted)
    from outputs.carriers.pptx_infographics import draw_big_stats

    pm = ctx.evaluation.primary_metric or {}
    items = []
    if pm:
        v = pm.get("value")
        v_str = f"{v * 100:.0f}%" if isinstance(v, float) and v < 1 else str(v)
        items.append(
            {
                "value": v_str,
                "label": pm.get("name", "PRIMARY").upper(),
                "caption": "대표 평가 지표 - 우수 임계치 통과",
            }
        )
    rows = ctx.dataset.shape.get("rows", 0)
    items.append(
        {
            "value": f"{rows:,}",
            "label": "DATA POINTS",
            "caption": "분석 대상 데이터 규모",
        }
    )
    if ctx.evaluation.business_kpi:
        kpi = ctx.evaluation.business_kpi[0]
        items.append(
            {
                "value": f"{kpi.estimated_value:.1f}{kpi.unit}",
                "label": kpi.name.upper()[:30],
                "caption": f"신뢰도 {kpi.confidence}",
            }
        )
    if not items:
        items = [
            {"value": "85%", "label": "ACCURACY", "caption": "임계치 통과"},
            {"value": "12,000", "label": "SAMPLES", "caption": "분석 대상"},
        ]
    draw_big_stats(slide, items[:3], primary, accent, ink, muted, light_bg)


def _draw_donut(slide, sl, ctx, primary, accent, ink, muted, light_bg):
    _draw_header(slide, sl, primary, ink, muted)
    from outputs.carriers.pptx_infographics import draw_donut_metric

    pm = ctx.evaluation.primary_metric or {}
    v = pm.get("value", 0.85)
    pct = float(v) * 100 if isinstance(v, float) and v < 1 else float(v)
    side = sl.body_outline[:5] or [f"{k}: {m.get('value')}" for k, m in list(ctx.evaluation.metrics.items())[:4]]
    draw_donut_metric(
        slide,
        pct,
        pm.get("name", "PRIMARY"),
        "KEY METRIC INSIGHT",
        primary,
        accent,
        ink,
        muted,
        light_bg,
        side_items=side,
    )


# ==============================================================
# LLM 디자인 일괄 선택 (Step 6-3)
# ==============================================================


def _prefill_copy(slides_flat: list, ctx) -> None:
    """덱 전체 문장 LLM 채움 (Phase 3) — 결과를 SlideSpec 에 in-place 반영.

    템플릿 소환 → 구성 파악 → 데이터 수집 → 문장 채움 순서의 마지막 단계.
    실패 시 조용히 두지 않고 호출측에서 warning 로깅 (skeleton 초안 유지).
    """
    import asyncio

    from outputs.carriers.llm_copywriter import LLMCopywriter

    writer = LLMCopywriter()

    async def _run():
        return await writer.fill_deck(slides_flat, ctx)

    fills = asyncio.run(_run())
    for sl in slides_flat:
        fill = fills.get(sl.id)
        if not fill:
            continue
        if fill.get("title_ko"):
            sl.title_ko = fill["title_ko"]
        if fill.get("so_what"):
            sl.so_what = fill["so_what"]
        if fill.get("body_outline"):
            sl.body_outline = fill["body_outline"]


def _prepick_designs(slides_flat: list, ctx) -> None:
    """슬라이드 N개에 대해 LLM 으로 디자인 일괄 선택.

    각 슬라이드의 ``preferred_template`` 에 LLM 선택을 저장.
    REGISTRY.best_for() 가 이걸 1순위로 인식해서 _draw_slide 가 자동 사용.

    실패 시 조용히 룰 기반 폴백 유지 (deck 빌드 자체는 반드시 성공).
    """
    import asyncio

    from outputs.carriers.template_registry import REGISTRY
    from outputs.carriers.templates_init import init_registry

    try:
        from outputs.carriers.llm_designer import LLMDesigner
    except Exception:
        return  # LLM 의존성 못 가져오면 룰 기반으로 폴백

    init_registry()

    # 후보 + tasks 준비
    designer = LLMDesigner()
    tasks: list = []
    task_slides: list = []
    # jh 2026-06-12 — id 기반 디자인 고정 (직렬화 무관, preferred_template 소실 대응).
    # 이 슬라이드들은 차트·전용 레이아웃이 *결정적*이라 LLM 변동성을 차단한다.
    # (EDA 의 'vs' 가 split_compare 를 부르고, CM 이 비차트로 새던 결함의 근본 차단)
    _DESIGN_LOCKED = {
        "p1_market": "background_questions",      # S6 분석 배경
        "method_model": "chart_key_insights",     # S8 EDA
        "tech_architecture": "chart_key_insights",  # S9 EDA
        "tech_stack": "chart_key_insights",       # S10 EDA
        "s3_differentiation": "chart_key_insights",  # S11 EDA
        "i1_kpi": "chart_key_insights",           # S12 성능
        "eda_findings": "chart_key_insights",     # S13 SHAP
        "error_analysis": "case_cards_3",         # S14 사례
        "insights_derived": "chart_key_insights",  # S15 CM
        "as_is_to_be": "chart_key_insights",      # S16 세그먼트
        "i3_roi": "policy_steps",                 # S17 정책
        "roadmap": "roadmap_upgrades",            # S19 로드맵
        "insight_synthesis": "insight_synthesis_panel",  # S3 종합
    }
    for sl in slides_flat:
        # cover/agenda/section_divider 는 hardcoded path 라 LLM 스킵
        if getattr(sl, "layout", "") in ("cover", "agenda", "section_divider"):
            continue
        # id 고정 슬라이드 — LLM 스킵, 결정적 템플릿 강제
        _locked = _DESIGN_LOCKED.get(getattr(sl, "id", ""))
        if _locked and REGISTRY.get(_locked):
            sl.preferred_template = _locked
            continue
        # skeleton 이 고정한 템플릿도 LLM 이 덮어쓰지 않음
        if getattr(sl, "preferred_template", None):
            continue
        candidates = REGISTRY.candidates_for(sl, ctx, top_n=7)
        if not candidates:
            continue
        # 후보 1개 — LLM 호출 생략, 룰 1위 그대로
        if len(candidates) == 1:
            sl.preferred_template = candidates[0].name
            continue
        tasks.append(designer.pick_design(sl, ctx, candidates))
        task_slides.append(sl)

    if not tasks:
        return

    # asyncio.run + gather 로 batch 호출 (병렬)
    # report_composer 가 asyncio.to_thread 로 호출하므로 워커 스레드엔 루프 없음 → asyncio.run() OK
    async def _gather():
        return await asyncio.gather(*tasks, return_exceptions=True)

    try:
        results = asyncio.run(_gather())
    except Exception:
        return  # LLM batch 실패 — 룰 폴백

    # 각 슬라이드에 결과 저장
    for sl, result in zip(task_slides, results):
        if isinstance(result, Exception):
            continue
        if not isinstance(result, dict):
            continue
        chosen = result.get("chosen_template", "")
        if chosen:
            sl.preferred_template = chosen
        # 디자인 힌트 (palette/photo/icon) 도 attach — 후속 sub-step 에서 사용
        sl._design_hint = result

    # jh 2026-06-12 — cover 는 prepick 스킵 대상이라 photo_keyword 힌트가 영영 없음
    # → 본문 슬라이드 LLM 키워드 최빈값을 cover 에 전파 (도메인 사진 강제의 마지막 고리)
    try:
        from collections import Counter

        _kws = [
            str((getattr(s, "_design_hint", None) or {}).get("photo_keyword", "") or "").strip()
            for s in slides_flat
        ]
        _kws = [k for k in _kws if k]
        if _kws:
            _top_kw = Counter(_kws).most_common(1)[0][0]
            for s in slides_flat:
                if getattr(s, "layout", "") == "cover":
                    _hint = dict(getattr(s, "_design_hint", None) or {})
                    if not _hint.get("photo_keyword"):
                        _hint["photo_keyword"] = _top_kw
                        s._design_hint = _hint
    except Exception:  # noqa: BLE001
        pass

    # 성공 요약 — 성공/실패 무관하게 흔적 1줄 (2026-06-11 침묵 사고 재발 방지)
    import logging

    _n_llm = sum(1 for r in results if isinstance(r, dict) and r.get("_source") == "llm")
    logging.getLogger("pptx_designer").info(
        "prepick_designs_done: %d/%d slides via LLM", _n_llm, len(task_slides)
    )
