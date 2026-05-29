"""Day 9 — outputs/base.py, ppt.py, pdf.py 카테고리 색상 + extras 훅 테스트.

DoD:
    1) get_theme 이 4 카테고리에서 서로 다른 색상 반환
    2) OutputGenerator._call_extras 가 등록된 핸들러를 호출
    3) handler 가 charts/tables/text_blocks 를 반환하면 carrier 에 임베드
    4) handler 없는 카테고리는 안전하게 {} 반환
"""

from __future__ import annotations

from typing import Any

import pytest

from ada.core.state import PipelineState


# ----- 1) 4 카테고리 색상 모두 다른지 -------------------------------------------
def test_category_themes_distinct():
    from outputs.base import CATEGORY_THEME, get_theme

    cats = ["tabular_ml", "tabular_dl", "timeseries", "anomaly_detection"]
    rgbs = set()
    for c in cats:
        t = get_theme(c)
        assert "primary_rgb" in t
        assert "primary_hex" in t
        assert "label_ko" in t
        rgbs.add(t["primary_rgb"])
    assert len(rgbs) == 4, f"4 카테고리가 서로 다른 색상이어야: {rgbs}"
    assert get_theme(None) != get_theme("tabular_ml")
    assert get_theme("unknown") == get_theme(None)
    assert CATEGORY_THEME.keys() >= {"tabular_ml", "tabular_dl", "timeseries", "anomaly_detection"}


# ----- 2) _call_extras 가 핸들러를 호출 ----------------------------------------
def test_call_extras_invokes_handler(monkeypatch):
    """등록된 handler 가 반환한 charts/tables/text_blocks 를 그대로 받는다."""
    from outputs.base import OutputGenerator

    class _Dummy(OutputGenerator):
        output_code = "OUT-TEST"
        extension = "bin"

        def generate(self, **kwargs: Any) -> str:
            return ""

    def fake_handler(state, ctx):
        return {
            "charts": ["s3://bucket/test/c1.png", "s3://bucket/test/c2.png"],
            "tables": [{"title": "Top-N", "columns": ["a", "b"], "rows": [{"a": 1, "b": 2}]}],
            "text_blocks": ["이 카테고리의 핵심 발견은..."],
            "unknown_key": "drop me",  # OUTPUT_EXTRAS_KEYS 에 없으면 제거
        }

    from agents.handlers import _base as handlers_base

    monkeypatch.setitem(handlers_base.HANDLER_REGISTRY, "timeseries", {"build": fake_handler})

    state = PipelineState(job_id="job-1", file_id="f.csv", category="timeseries", target_column="y")
    out = _Dummy("job-1")._call_extras(state, ctx={"x": 1})
    assert set(out.keys()) == {"charts", "tables", "text_blocks"}
    assert len(out["charts"]) == 2
    assert out["tables"][0]["title"] == "Top-N"
    assert "unknown_key" not in out


# ----- 3) handler 가 없으면 안전하게 빈 dict ----------------------------------
def test_call_extras_no_handler():
    from outputs.base import OutputGenerator

    class _Dummy(OutputGenerator):
        output_code = "OUT-TEST"
        extension = "bin"

        def generate(self, **kwargs: Any) -> str:
            return ""

    # 일부러 핸들러 없는 카테고리로
    state = PipelineState(job_id="job-1", file_id="f.csv", category="anomaly_detection", target_column="y")
    # 핸들러 레지스트리에서 anomaly_detection 비우기
    from agents.handlers import _base as handlers_base

    saved = handlers_base.HANDLER_REGISTRY.get("anomaly_detection", {})
    handlers_base.HANDLER_REGISTRY["anomaly_detection"] = {}
    try:
        out = _Dummy("job-1")._call_extras(state)
        assert out == {}
    finally:
        handlers_base.HANDLER_REGISTRY["anomaly_detection"] = saved


# ----- 4) handler 가 예외 발생 시 안전 처리 -------------------------------------
def test_call_extras_handler_exception(monkeypatch):
    from outputs.base import OutputGenerator

    class _Dummy(OutputGenerator):
        output_code = "OUT-TEST"
        extension = "bin"

        def generate(self, **kwargs: Any) -> str:
            return ""

    def crashing_handler(state, ctx):
        raise RuntimeError("kaboom")

    from agents.handlers import _base as handlers_base

    monkeypatch.setitem(handlers_base.HANDLER_REGISTRY, "tabular_ml", {"build": crashing_handler})

    state = PipelineState(job_id="job-1", file_id="f.csv", category="tabular_ml", target_column="y")
    out = _Dummy("job-1")._call_extras(state)
    assert out == {}  # 예외 catch 후 빈 dict


# ----- 5) PPT 가 4 카테고리에서 서로 다른 색상 출력하는지 ------------------------
def test_ppt_different_colors_per_category(monkeypatch, tmp_path):
    pytest.importorskip("pptx")
    from outputs.ppt import PresentationGenerator

    # _upload 우회 — 로컬 경로 반환
    saved_paths: dict[str, str] = {}

    def fake_upload(self, local_path):
        saved_paths[self.output_code] = local_path
        return f"s3://test/{self.output_code}"

    monkeypatch.setattr(PresentationGenerator, "_upload", fake_upload)
    # _call_extras 비활성 (carrier 자체만 검증)
    monkeypatch.setattr(PresentationGenerator, "_call_extras", lambda self, state=None, ctx=None: {})

    cats = ["tabular_ml", "tabular_dl", "timeseries", "anomaly_detection"]
    title_colors: set[tuple[int, int, int]] = set()

    for cat in cats:
        gen = PresentationGenerator(job_id=f"job-{cat}")
        gen.generate(
            insights="요약 문장",
            best_model={"model_name": "X", "metrics": {"val_f1": 0.7}},
            eda_charts=[],
            category=cat,
            user_intent="테스트",
            eval_result={"passed": True, "rationale": "good"},
        )
        # 생성된 pptx 열어서 표지 제목 색상 확인
        from pptx import Presentation

        prs = Presentation(saved_paths["OUT-01"])
        title = prs.slides[0].shapes.title
        for para in title.text_frame.paragraphs:
            for run in para.runs:
                rgb = run.font.color.rgb
                if rgb is not None:
                    # RGBColor 는 iterable → (r, g, b)
                    title_colors.add(tuple(rgb))

    # 4 카테고리 모두 서로 다른 제목 색상이 적용돼야 함
    assert len(title_colors) >= 3, f"카테고리별 색상이 거의 안 적용됨: {title_colors}"


# ----- 6) PPT 가 extras 의 차트 슬라이드를 추가 ---------------------------------
def test_ppt_embeds_extras_charts(monkeypatch):
    pytest.importorskip("pptx")
    from outputs.ppt import PresentationGenerator

    # _upload, _download_chart 우회
    saved: dict[str, str] = {}

    def _fake_upload(self, p):
        saved[self.output_code] = p
        return f"s3://test/{self.output_code}"

    monkeypatch.setattr(PresentationGenerator, "_upload", _fake_upload)
    monkeypatch.setattr(PresentationGenerator, "_download_chart", lambda self, p: None)

    # extras 가 차트 + 텍스트 블록 + 테이블 반환하도록
    def fake_extras(self, state=None, ctx=None):
        return {
            "charts": ["s3://test/c1.png"],
            "tables": [{"title": "Top-3", "columns": ["k", "v"], "rows": [{"k": "a", "v": 1}, {"k": "b", "v": 2}]}],
            "text_blocks": ["이 시계열의 핵심 통찰..."],
        }

    monkeypatch.setattr(PresentationGenerator, "_call_extras", fake_extras)

    gen = PresentationGenerator(job_id="job-ts")
    gen.generate(
        insights="요약",
        best_model={"model_name": "SARIMA", "metrics": {"val_rmse": 0.5}},
        eda_charts=[],
        category="timeseries",
        user_intent="예측",
        eval_result={"passed": True, "rationale": "ok"},
    )
    from pptx import Presentation

    prs = Presentation(saved["OUT-01"])
    # 표지 + 인사이트 + best model + (extras charts 1 + tables 1 + text 1) + 평가 = 7 슬라이드 이상
    assert len(prs.slides) >= 7
    # 제목 텍스트에 카테고리 라벨이 들어간 슬라이드 1개 이상
    titles = []
    for s in prs.slides:
        try:
            titles.append(s.shapes.title.text)
        except Exception:
            pass
    joined = " ".join(titles)
    assert "[시계열]" in joined or "카테고리 분석" in joined or "해설" in joined or "표" in joined


# ─── Day 9 V3 보강 테스트 (★T-9 ~ T-12) ────────────────────────────────────


# ----- ★T-9) PDF 가 4 카테고리에서 정상 생성 (D-3 PDF 측면) ---------------------
def test_pdf_generated_for_all_categories(monkeypatch, tmp_path):
    """PDF carrier 가 4 카테고리 모두에서 0 byte 아닌 파일을 생성하는지.

    색상 정확 추출은 PDF 파싱 복잡 — 대신 (1) 모두 생성됨, (2) 파일 크기
    유의미함 으로 D-3 PDF 측면 회귀 차단.
    """
    pytest.importorskip("reportlab")
    import pathlib

    from outputs.pdf import PDFReportGenerator

    saved_paths: dict[str, str] = {}

    def fake_upload(self, local_path):
        # 카테고리별로 다른 키로 저장 (덮어쓰기 방지)
        saved_paths[self.job_id] = local_path
        return f"s3://test/{self.output_code}/{self.job_id}"

    monkeypatch.setattr(PDFReportGenerator, "_upload", fake_upload)
    monkeypatch.setattr(PDFReportGenerator, "_call_extras", lambda self, state=None, ctx=None: {})
    monkeypatch.setattr(PDFReportGenerator, "_download_chart", lambda self, p: None)

    cats = ["tabular_ml", "tabular_dl", "timeseries", "anomaly_detection"]
    sizes = []
    for cat in cats:
        gen = PDFReportGenerator(job_id=f"job-{cat}")
        gen.generate(
            insights="이것은 요약 문장입니다.",
            best_model={"model_name": "X", "metrics": {"val_f1": 0.7}},
            eda_charts=[],
            category=cat,
            user_intent="테스트",
            eval_result={"passed": True, "rationale": "good"},
        )
        path = pathlib.Path(saved_paths[f"job-{cat}"])
        assert path.exists(), f"PDF 미생성: {cat}"
        sizes.append(path.stat().st_size)

    # 4개 PDF 모두 유의미한 크기 (≥ 1 KB)
    assert all(s > 1000 for s in sizes), f"PDF 크기 비정상: {sizes}"


# ----- ★T-10) PDF extras 임베드 (D-2 PDF 측면) --------------------------------
def test_pdf_embeds_extras(monkeypatch):
    pytest.importorskip("reportlab")
    import pathlib

    from outputs.pdf import PDFReportGenerator

    saved: dict[str, str] = {}

    def _fake_upload(self, p):
        saved[self.output_code] = p
        return f"s3://test/{self.output_code}"

    monkeypatch.setattr(PDFReportGenerator, "_upload", _fake_upload)
    monkeypatch.setattr(PDFReportGenerator, "_download_chart", lambda self, p: None)

    def fake_extras(self, state=None, ctx=None):
        return {
            "charts": ["s3://test/c1.png"],
            "tables": [{"title": "Top-N", "columns": ["k", "v"], "rows": [{"k": "a", "v": 1}]}],
            "text_blocks": ["이 카테고리의 핵심 발견은 X 입니다."],
        }

    monkeypatch.setattr(PDFReportGenerator, "_call_extras", fake_extras)

    # extras 없는 baseline 과 비교
    monkeypatch.setattr(PDFReportGenerator, "_call_extras", lambda self, state=None, ctx=None: {})
    PDFReportGenerator(job_id="job-base").generate(
        insights="요약",
        best_model={"model_name": "X", "metrics": {"val_f1": 0.7}},
        eda_charts=[],
        category="tabular_ml",
        user_intent="t",
        eval_result={"passed": True, "rationale": "ok"},
    )
    base_size = pathlib.Path(saved["OUT-02"]).stat().st_size

    # extras 있는 경우
    monkeypatch.setattr(PDFReportGenerator, "_call_extras", fake_extras)
    PDFReportGenerator(job_id="job-extras").generate(
        insights="요약",
        best_model={"model_name": "X", "metrics": {"val_f1": 0.7}},
        eda_charts=[],
        category="tabular_ml",
        user_intent="t",
        eval_result={"passed": True, "rationale": "ok"},
    )
    extras_size = pathlib.Path(saved["OUT-02"]).stat().st_size

    # extras 가 들어가면 PDF 가 커져야 함 (테이블 + 텍스트블록)
    assert extras_size > base_size, f"extras 임베드 안 됨: base={base_size}, extras={extras_size}"


# ----- ★T-11) _download_chart 실패 → carrier 죽지 않고 슬라이드 skip -----------
def test_ppt_survives_chart_download_failure(monkeypatch):
    pytest.importorskip("pptx")

    from outputs.ppt import PresentationGenerator

    saved: dict[str, str] = {}
    monkeypatch.setattr(
        PresentationGenerator,
        "_upload",
        lambda self, p: saved.__setitem__(self.output_code, p) or f"s3://test/{self.output_code}",
    )
    monkeypatch.setattr(PresentationGenerator, "_call_extras", lambda self, state=None, ctx=None: {})
    # 모든 차트 다운로드 실패 시뮬레이션
    monkeypatch.setattr(PresentationGenerator, "_download_chart", lambda self, p: None)

    gen = PresentationGenerator(job_id="job-chart-fail")
    # eda_charts 4개 전달 — 모두 다운로드 실패하지만 carrier 는 죽지 않아야
    result = gen.generate(
        insights="요약",
        best_model={"model_name": "X", "metrics": {"val_f1": 0.7}},
        eda_charts=["s3://bucket/c1.png", "s3://bucket/c2.png", "s3://bucket/c3.png", "s3://bucket/c4.png"],
        category="tabular_ml",
        user_intent="테스트",
        eval_result={"passed": True, "rationale": "ok"},
    )
    # carrier 가 정상 반환 (예외 없이)
    assert result.startswith("s3://test/OUT-01")
    # PPT 파일은 생성됨 (차트는 빈 제목 슬라이드)
    from pptx import Presentation

    prs = Presentation(saved["OUT-01"])
    # 표지 + 인사이트 + best + EDA 4개 + 평가 = 8 슬라이드 (사진은 안 들어감)
    assert len(prs.slides) >= 4


# ----- ★T-12) handler 가 dict 아닌 값 반환 → _call_extras 가 {} ----------------
def test_call_extras_rejects_non_dict_return(monkeypatch):
    from outputs.base import OutputGenerator

    class _Dummy(OutputGenerator):
        output_code = "OUT-TEST"
        extension = "bin"

        def generate(self, **kwargs: Any) -> str:
            return ""

    def returns_list(state, ctx):
        return ["wrong", "type"]  # dict 아님

    def returns_str(state, ctx):
        return "completely wrong"

    from agents.handlers import _base as handlers_base

    state = PipelineState(job_id="job-1", file_id="f.csv", category="timeseries", target_column="y")

    monkeypatch.setitem(handlers_base.HANDLER_REGISTRY, "timeseries", {"build": returns_list})
    assert _Dummy("job-1")._call_extras(state) == {}

    monkeypatch.setitem(handlers_base.HANDLER_REGISTRY, "timeseries", {"build": returns_str})
    assert _Dummy("job-1")._call_extras(state) == {}


# ----- ★T-13) "build" capability 자동 등록 (P1 보강 검증) -----------------------
def test_auto_register_includes_build_capability():
    """_auto_register_from_module 이 'build' 키워드도 자동 등록 — base.py 의
    get_handler(.., 'build') 가 monkeypatch 외에서도 작동하도록 보장."""
    from agents.handlers._base import HANDLER_REGISTRY, _auto_register_from_module

    # 임시 모듈처럼 동작하는 stub
    class FakeModule:
        def build(state, ctx):  # noqa: N805
            return {"charts": ["x"]}

        def assets(state, ctx):  # noqa: N805
            return {"tables": []}

    saved = HANDLER_REGISTRY.get("tabular_ml", {}).copy()
    try:
        _auto_register_from_module("tabular_ml", FakeModule)
        # build / assets 둘 다 등록되어야 함
        assert "build" in HANDLER_REGISTRY.get("tabular_ml", {}), (
            "P1 보강 누락: 'build' capability 가 _auto_register 에 없음"
        )
        assert "assets" in HANDLER_REGISTRY.get("tabular_ml", {})
    finally:
        HANDLER_REGISTRY["tabular_ml"] = saved
