"""tests.integration.test_phase12_templates — 카탈로그 보강 (Phase 12) 회귀.

견본 덱 레이아웃 이식 4종: chart_key_insights (legacy 등록) +
solution_overview / lineage_2col / icon_columns (신규).
"""

from __future__ import annotations

import pytest

from outputs.architect.plan import SlideSpec, VisualSpec
from outputs.carriers.template_registry import REGISTRY
from outputs.carriers.templates_init import _split_label, init_registry
from outputs.context.schema import ReportContext

# CI 등 python-pptx 미설치 환경에서는 이 모듈 전체 skip (collection 차단 방지).
# 런타임 carrier 코드는 pptx 를 함수 내부에서 lazy import 하므로 무관. jh 2026-06-12.
_pptx = pytest.importorskip("pptx")
Presentation = _pptx.Presentation
Cm = pytest.importorskip("pptx.util").Cm

PALETTE = ("#0F4C81", "#2E86AB", "#0B1F33", "#5B7C99", "#F4F7FA")  # primary/accent/ink/muted/light_bg


def _blank_slide():
    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)
    return prs.slides.add_slide(prs.slide_layouts[6]), prs


def _slide_spec(**kw):
    base = dict(id="s", layout="", role="", title_ko="제목", so_what="주장", body_outline=[])
    base.update(kw)
    return SlideSpec(**base)


class TestRegistration:
    def test_phase12_registered(self):
        init_registry()
        for name in ("chart_key_insights", "solution_overview", "lineage_2col", "icon_columns"):
            assert REGISTRY.get(name) is not None, f"{name} 미등록"

    def test_chart_key_insights_requires_visual(self):
        init_registry()
        spec = REGISTRY.get("chart_key_insights")
        ctx = ReportContext()
        no_visual = _slide_spec(layout="chart_callout")
        assert spec.fit(no_visual, ctx) == 0.0
        with_visual = _slide_spec(
            layout="chart_callout", visual_spec=VisualSpec(type="chart_annotated_bar", title="t", spec={})
        )
        assert spec.fit(with_visual, ctx) >= 85.0

    def test_chart_key_insights_evidence_fallback(self):
        init_registry()
        spec = REGISTRY.get("chart_key_insights")
        ctx = ReportContext()
        sl = _slide_spec(role="evidence", visual_spec=VisualSpec(type="x", title="t", spec={}))
        assert spec.fit(sl, ctx) >= 70.0


class TestSplitLabel:
    def test_middle_dot(self):
        assert _split_label("모델 · CatBoost 채택") == ("모델", "CatBoost 채택")

    def test_colon(self):
        assert _split_label("검증: 80/20 split") == ("검증", "80/20 split")

    def test_no_separator_uses_fallback(self):
        label, text = _split_label("구분자가 없는 문장", fallback="항목 1")
        assert label == "항목 1"
        assert text == "구분자가 없는 문장"

    def test_overlong_label_not_split(self):
        long_head = "라" * 30 + " · 본문"
        label, text = _split_label(long_head, fallback="F")
        assert label == "F"  # 24자 초과 라벨은 라벨로 안 봄


class TestSmokeRender:
    """blank slide 에 그려서 예외 없는지 + shape 생성 확인."""

    def _ctx(self):
        ctx = ReportContext()
        ctx.model_selection.chosen = {"name": "CatBoost"}
        return ctx

    def test_solution_overview(self):
        init_registry()
        slide, _ = _blank_slide()
        sl = _slide_spec(
            body_outline=["모델 · CatBoost GBM", "입력 피처 · 14개", "출력 · 생존 확률 + SHAP", "검증 · 80/20 split"]
        )
        REGISTRY.get("solution_overview").draw(slide, sl, self._ctx(), *PALETTE)
        assert len(slide.shapes) >= 8

    def test_lineage_2col(self):
        init_registry()
        slide, _ = _blank_slide()
        sl = _slide_spec(
            body_outline=[
                "원본 · 891행 CSV", "전처리 · 결측 대치", "피처 · 14개 생성", "모델 · CatBoost 학습",
                "범주형 · 별도 인코딩", "파생 · Title 추출",
            ]
        )
        REGISTRY.get("lineage_2col").draw(slide, sl, self._ctx(), *PALETTE)
        assert len(slide.shapes) >= 10

    def test_icon_columns(self):
        init_registry()
        slide, _ = _blank_slide()
        sl = _slide_spec(body_outline=["데이터 · 891명 학습", "패턴 · 성별 55%p", "모델 · 0.788 달성", "액션 · 정책 제언"])
        REGISTRY.get("icon_columns").draw(slide, sl, self._ctx(), *PALETTE)
        assert len(slide.shapes) >= 12

    def test_icon_columns_two_items_ok(self):
        init_registry()
        slide, _ = _blank_slide()
        sl = _slide_spec(body_outline=["하나 · 첫째", "둘 · 둘째"])
        REGISTRY.get("icon_columns").draw(slide, sl, self._ctx(), *PALETTE)
        assert len(slide.shapes) >= 6

    def test_chart_key_insights_no_visual_no_crash(self):
        """visual_spec 없어도 (fit 상 안 오지만) draw 가 예외 없이 패널만 그림."""
        init_registry()
        slide, _ = _blank_slide()
        sl = _slide_spec(body_outline=["인사이트 1", "인사이트 2"])
        REGISTRY.get("chart_key_insights").draw(slide, sl, self._ctx(), *PALETTE)
        assert len(slide.shapes) >= 3

    def test_chart_key_insights_panel_fullwidth_when_no_chart(self):
        """2026-06-12 — 차트 미배치 시 패널 전폭 확장 (좌측 절반 공백 결함)."""
        from pptx.util import Emu

        init_registry()
        slide, _ = _blank_slide()
        sl = _slide_spec(body_outline=["인사이트 1", "인사이트 2"])
        REGISTRY.get("chart_key_insights").draw(slide, sl, self._ctx(), *PALETTE)
        widths = [Emu(sh.width).cm for sh in slide.shapes if Emu(sh.width).cm > 20]
        assert widths, "전폭(>20cm) 패널이 없음 — 빈 절반 결함 재발"


class TestVisualFixes20260612:
    """운영 첫 덱 결함 회귀 — cover 제목 폭발 / 문장 절단 / chevron 반복."""

    def test_display_intent_strips_tags(self):
        from outputs.carriers.pptx_designer import _display_intent

        raw = "분석 방향: 고객 가치 예측: 생존자 분류 (주제: 이커머스 인사이트) (방법론: 앙상블) (모델 전략: 부스팅)"
        assert _display_intent(raw) == "고객 가치 예측: 생존자 분류"

    def test_display_intent_plain_passthrough(self):
        from outputs.carriers.pptx_designer import _display_intent

        assert _display_intent("타이타닉 생존 분석") == "타이타닉 생존 분석"
        assert _display_intent("") == "분석 보고서"

    def test_display_intent_truncates(self):
        from outputs.carriers.pptx_designer import _display_intent

        out = _display_intent("가" * 100, max_len=44)
        assert len(out) <= 44

    def test_cover_title_size_shrinks(self):
        from outputs.carriers.pptx_designer import _cover_title_size

        # jh 2026-06-12 — 사이즈 사다리 재보정(48→40) 동기화: 짧은 제목 40, 34자 초과 24
        assert _cover_title_size("짧은 제목") == 40
        assert _cover_title_size("가" * 44) == 24

    def test_clip_sentence_no_midword_cut(self):
        from outputs.carriers.templates_init import _clip_sentence

        s = "데이터 특성 기반 점수 매트릭스로 상위 3종 자동 선정, 고카디널리티 변수 처리 강점이 결정적이었으며 추가 검증 예정"
        out = _clip_sentence(s, 40)
        assert len(out) <= 41
        assert out.endswith("…")
        assert not out.rstrip("…").endswith(("데이", "매트"))  # 단어 중간 절단 아님 (공백 경계)

    def test_qa_shrink_overflow(self):
        """편집자 QA v1 — 박스보다 큰 텍스트의 폰트 자동 축소."""
        from pptx.util import Cm, Pt

        from outputs.carriers.pptx_designer import _qa_shrink_overflow

        slide, prs = _blank_slide()
        box = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(8), Cm(2))
        p = box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "아주 긴 한국어 문장이 박스를 가득 넘쳐 흐르는 경우를 시뮬레이션 합니다 " * 3
        run.font.size = Pt(24)
        fixed = _qa_shrink_overflow(prs)
        assert fixed >= 1
        assert run.font.size.pt < 24

    def test_qa_leaves_fitting_text_alone(self):
        from pptx.util import Cm, Pt

        from outputs.carriers.pptx_designer import _qa_shrink_overflow

        slide, prs = _blank_slide()
        box = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(20), Cm(5))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "짧은 제목"
        run.font.size = Pt(24)
        _qa_shrink_overflow(prs)
        assert run.font.size.pt == 24

    def test_chevron_dedupes_identical_justifications(self):
        from types import SimpleNamespace

        from outputs.carriers.templates_init import t_chevron_5

        init_registry()
        slide, _ = _blank_slide()
        ctx = ReportContext()
        same = "데이터 특성 기반 점수 매트릭스로 상위 3종 자동 선정"
        ctx.model_selection.candidates = [
            SimpleNamespace(name="CatBoost", why_tried=same, family="boosting", score=0.788),
            SimpleNamespace(name="LightGBM", why_tried=same, family="boosting", score=0.771),
            SimpleNamespace(name="RandomForest", why_tried=same, family="bagging", score=0.745),
        ]
        sl = _slide_spec(id="p3_alt_limits")
        t_chevron_5(slide, sl, ctx, *PALETTE)
        texts = " ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
        assert texts.count(same[:20]) == 0, "동일 사유 반복이 여전히 렌더됨"
        assert "score 0.788" in texts
